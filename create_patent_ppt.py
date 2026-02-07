import collections 
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_presentation():
    prs = Presentation()

    # defined colors
    MG_RED = RGBColor(186, 12, 47)  # Approximate MG Red
    JSW_BLUE = RGBColor(0, 51, 153) # Approximate JSW Blue
    BLACK = RGBColor(0, 0, 0)
    GRAY = RGBColor(128, 128, 128)

    def add_slide(title_text, content_text_list, is_title=False):
        if is_title:
            slide_layout = prs.slide_layouts[0] # Title Slide
        else:
            slide_layout = prs.slide_layouts[1] # Title and Content

        slide = prs.slides.add_slide(slide_layout)

        # Set Title
        title = slide.shapes.title
        title.text = title_text
        
        # Style Title
        title.text_frame.paragraphs[0].font.name = 'Arial'
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.color.rgb = MG_RED
        
        if is_title:
            subtitle = slide.placeholders[1]
            subtitle.text = "\n".join(content_text_list)
            subtitle.text_frame.paragraphs[0].font.size = Pt(20)
        else:
            # Set Content
            body_shape = slide.placeholders[1]
            tf = body_shape.text_frame
            tf.clear() # Clear default empty paragraph

            for item in content_text_list:
                p = tf.add_paragraph()
                p.text = item
                p.font.size = Pt(20)
                p.space_after = Pt(10)
                if item.startswith("    "):
                    p.level = 1
                    p.text = item.strip()
                elif item.startswith("        "):
                    p.level = 2
                    p.text = item.strip()

    # Slide 1: Title
    add_slide(
        "Patent Strategy: VAVE AI System", 
        [
            "Technical Novelty & Patent Filing Roadmap",
            "Prepared for: JSW MG Motor India",
            "Context: Internship Project IP"
        ], 
        is_title=True
    )

    # Slide 2: Executive Summary
    add_slide(
        "Executive Summary",
        [
            "Objective: Transform the 'VAVE AI' intern project into secured Intellectual Property (IP) for JSW MG.",
            "Core Innovation: 'Physics-Informed Generative AI'—solving the hallucination problem in engineering.",
            "Recommendation: File a Provisional Patent Application immediately to secure priority date.",
            "Ownership: Inventor = Intern; Applicant/Owner = JSW MG Motor India (standard corporate practice)."
        ]
    )

    # Slide 3: The Problem vs. The Solution
    add_slide(
        "The Technical Problem",
        [
            "Common LLMs (GPT-4, Gemini) lack 'Engineering Physics'.",
            "Risk: They hallucinate impossible components (e.g., 'Carbon Fiber Exhaust' on an EV).",
            "Result: Generic AI is dangerous for Engineering Decision Making.",
            "",
            "Our Solution: The VAVE AI System",
            "    1. Physics-Informed Matrix (Pre-Filter)",
            "    2. Autonomous Validation Engine (Post-Filter)",
            "    3. Visual Language Model (VLM) for 'Before vs. After' visualization"
        ]
    )

    # Slide 4: Patentable Claim 1 (The Matrix)
    add_slide(
        "Claim A: Context Resolution System",
        [
            "Invention: 'Physics-Informed Context Resolution System'",
            "Novelty: A look-up mechanism that maps fuzzy user queries to rigid Vehicle x Component constraints.",
            "Technical Effect: Prevents the AI from processing physically impossible requests.",
            "Code Evidence: `_resolve_engineering_context` algorithm.",
            "    'It strictly defines the Engineering Reality before the AI is allowed to speak.'"
        ]
    )

    # Slide 5: Patentable Claim 2 (The Gatekeeper)
    add_slide(
        "Claim B: Validation Engine",
        [
            "Invention: 'Deterministic Autonomous Validation Engine'",
            "Novelty: A post-generation filter using hard-coded physics rules (not AI probabilities).",
            "Example Rule: 'IF Weight > 1600kg AND Component == Brake Disc -> REJECT Thinning'.",
            "Technical Effect: Ensures only safe, viable ideas are presented to engineers.",
            "Code Evidence: `_validate_and_filter_ideas` function."
        ]
    )

    # Slide 6: Patenting the "Whole System" (Product)
    add_slide(
        "Patenting the 'Whole System'",
        [
            "Q: Can we patent the entire VAVE Product?",
            "A: Yes, as a 'System Patent'.",
            "",
            "Strategy: Avoid Section 3(k) (Software per se).",
            "Claim It As: 'A System for Automated Engineering Optimization' comprising:",
            "    1. Data Ingestion Module (DB/Web)",
            "    2. Physics Enforcement Module (The Matrix)",
            "    3. Generative Validation Module (The Gatekeeper)",
            "    4. Hardware Elements (Processors/Servers)",
            "Outcome: This protects the *entire workflow* as a product."
        ]
    )

    # Slide 7: Indian Patent Process (Interns)
    add_slide(
        "Filing Process: India",
        [
            "Step 1: File Provisional Specification (Month 0)",
            "    - Secures the 'Priority Date'.",
            "    - Low cost, less formal.",
            "    - Allows 'Patent Pending' status.",
            "Step 2: R&D & Refinement (Months 1-11)",
            "Step 3: File Complete Specification (Month 12)",
            "    - Full claims, flowcharts, and final code logic."
        ]
    )

    # Slide 7: Cost Breakdown & Ownership
    add_slide(
        "Cost & Ownership Structure",
        [
            "Ownership Model:",
            "    - True & First Inventor: Student Intern (You)",
            "    - Applicant/Assignee: JSW MG Motor India (The Company)",
            "    - Note: Company usually pays all fees.",
            "",
            "Estimated Official Fees (2025):",
            "    - Provisional Filing: ~₹8,000 (Large Entity) / ~₹1,600 (Student)",
            "    - Examination Request: ~₹20,000 (Large Entity)",
            "    - Professional Legal Fees: ~₹50k - ₹1.5L (varies by firm)"
        ]
    )

    # Slide 8: Next Actions
    add_slide(
        "Next Steps",
        [
            "1. Review `patent_strategy.md` for full claim details.",
            "2. Confirm IP ownership policy with HR/Legal.",
            "3. Draft the Provisional Specification using the System Design doc.",
            "4. File Form 1 & Form 2 via Patent Attorney."
        ]
    )

    output_file = "VAVE_Patent_Strategy_Presentation_System.pptx"
    try:
        prs.save(output_file)
        print(f"Successfully generated {output_file}")
    except PermissionError:
        print(f"Error: Could not save to {output_file}. Is the file open?")
        # Fallback to a timestamped name if needed, but for now just fail gracefully or try a v3
        output_file_v3 = "VAVE_Patent_Strategy_Presentation_v3.pptx"
        prs.save(output_file_v3)
        print(f"Saved to {output_file_v3} instead.")

if __name__ == "__main__":
    create_presentation()

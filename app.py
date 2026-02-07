import os
import re
import pandas as pd
import numpy as np
from flask import Flask, request, jsonify, send_file, url_for, render_template, redirect, flash, session, abort
from flask_login import LoginManager, UserMixin, login_user, logout_user, login_required, current_user
from werkzeug.security import generate_password_hash, check_password_hash
import sqlite3
from transformers import GPT2Tokenizer, BlipProcessor, BlipForConditionalGeneration, AutoTokenizer
import torch
from gtts import gTTS
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import markdown
from bs4 import BeautifulSoup
from flask_cors import CORS
import logging
from concurrent.futures import ThreadPoolExecutor
from sentence_transformers import SentenceTransformer, util
import faiss
import psycopg2                # <-- NEW
from psycopg2.extras import DictCursor # <-- NEW
from dotenv import load_dotenv # <-- NEW
from pathlib import Path
from PIL import Image
from datetime import datetime
import shutil
import threading
import nltk
from nltk.tokenize import word_tokenize
from nltk.corpus import stopwords
from werkzeug.utils import secure_filename
from werkzeug.security import generate_password_hash
import tempfile
import uuid
from functools import wraps
import string
import random
import data_processor
import data_lake
import json
import zipfile
import google.generativeai as genai
from agent import VAVEAgent
from vave_presentation_engine import generate_deep_dive_ppt
from excel_generator_engine import generate_excel_from_table_in_memory

'''# Download NLTK resources
nltk.download('punkt', quiet=True)
nltk.download('stopwords', quiet=True)
'''
# Load environment variables BEFORE using them
load_dotenv()

# Set up logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(levelname)s - %(message)s',
    handlers=[
        logging.FileHandler("logger.log"),
        logging.StreamHandler()
    ]
)
logger = logging.getLogger(__name__)

app = Flask(__name__)
CORS(app)
# Configuration - SECRET_KEY now reads from .env file correctly
app.config['SECRET_KEY'] = os.getenv('SECRET_KEY', 'a_very_secret_key_fallback')
BASE_DIR = Path(__file__).parent.resolve()
USER_DB_PATH = BASE_DIR / "users.db" # <-- NEW
DB_PATH = BASE_DIR / "cost_reduction.db" # <-- This is now just a reference
MODEL_DIR = BASE_DIR / "model"
MODEL_DIR = BASE_DIR / "model"
MODEL_ZIP_PATH = BASE_DIR / "model.zip"
TORCHSCRIPT_MODEL_PATH = MODEL_DIR / "gpt2.pt"
FAISS_INDEX_PATH = MODEL_DIR / "faiss_index.bin"
STATIC_DIR = BASE_DIR / "static"
VLM_MODEL_PATH = MODEL_DIR / "fine_tuned_blip_model"
TEMP_DIR = BASE_DIR / "temp"
IMAGE_DIRS = {
    "proposal": BASE_DIR / "images" / "proposal",
    "mg": BASE_DIR / "images" / "mg",
    "competitor": BASE_DIR / "images" / "competitor"
}
# DATA_PATH = BASE_DIR / "AIML Dummy Ideas Data.xlsx"
UPLOAD_FOLDER = TEMP_DIR / "uploads"
# Increase max upload size to handle larger Excel knowledge base files (e.g. 50 MB)
MAX_UPLOAD_SIZE = 50 * 1024 * 1024  # 50MB
ALLOWED_EXTENSIONS = {'png', 'jpg', 'jpeg', 'gif'}

app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER
app.config['STATIC_FOLDER'] = STATIC_DIR
app.config['MAX_CONTENT_LENGTH'] = MAX_UPLOAD_SIZE

# Create directories
TEMP_DIR.mkdir(exist_ok=True)
UPLOAD_FOLDER.mkdir(exist_ok=True)
STATIC_DIR.mkdir(exist_ok=True)
for dir_path in IMAGE_DIRS.values():
    dir_path.mkdir(parents=True, exist_ok=True)

# Initialize data lake (big-data style storage) on startup
data_lake.init_data_lake()

STATIC_IMAGE_DIR = STATIC_DIR / "images"
STATIC_IMAGE_DIR.mkdir(exist_ok=True)
# Create directory for VLM-generated images
STATIC_GENERATED_DIR = STATIC_DIR / "generated"
STATIC_GENERATED_DIR.mkdir(exist_ok=True)
for key, dir_path in IMAGE_DIRS.items():
    static_image_subdir = STATIC_IMAGE_DIR / key
    static_image_subdir.mkdir(exist_ok=True)
    if not any(static_image_subdir.glob("*")) and any(dir_path.glob("*")):
        try:
            shutil.copytree(dir_path, static_image_subdir, dirs_exist_ok=True)
            logger.info(f"Copied images from {dir_path} to {static_image_subdir}")
        except Exception as e:
            logger.error(f"Failed to copy images from {dir_path}: {str(e)}")

# Global variables
embedding_model = None
faiss_index = None
idea_texts = []
idea_rows = []
vave_agent = None  # VAVE Agent instance (initialized after vector DB is built)

# Ensure local model directory is available by extracting model.zip if present
def ensure_model_dir():
    try:
        if MODEL_DIR.exists():
            return
        if MODEL_ZIP_PATH.exists():
            with zipfile.ZipFile(MODEL_ZIP_PATH, 'r') as zf:
                zf.extractall(MODEL_DIR)
            logger.info(f"Extracted {MODEL_ZIP_PATH} to {MODEL_DIR}")
        else:
            logger.warning(f"Missing model directory {MODEL_DIR} and {MODEL_ZIP_PATH}; will attempt online fallback if allowed")
    except Exception as e:
        logger.error(f"Failed extracting model zip: {e}")

def find_gpt2_tokenizer_dir(base_dir: Path) -> Path | None:
    try:
        candidates = []
        # Check for tokenizer in gpt2_tokenizer subdirectory first
        gpt2_tok_dir = base_dir / "gpt2_tokenizer"
        if gpt2_tok_dir.exists():
            if (gpt2_tok_dir / "vocab.json").exists() and (gpt2_tok_dir / "merges.txt").exists():
                return gpt2_tok_dir
        
        # Also check all subdirectories
        for root, dirs, files in os.walk(base_dir):
            if 'vocab.json' in files and 'merges.txt' in files:
                candidates.append(Path(root))
        if candidates:
            # Prefer the shallowest directory
            candidates.sort(key=lambda p: len(p.parts))
            return candidates[0]
        return None
    except Exception as e:
        logger.error(f"Tokenizer discovery failed: {e}")
        return None
tokenizer = None
model = None
vlm_processor = None
vlm_model = None
proposal_embeddings = None
valid_df = None
vlm_initialized = False
device = 'cuda' if torch.cuda.is_available() else 'cpu'
executor = ThreadPoolExecutor(max_workers=4)
vlm_init_lock = threading.Lock()
# db_lock = threading.Lock()

# FIXED: Safe float conversion function
def safe_float_convert(value, default=0.0):
    """Safely convert value to float, handling various edge cases including backslashes"""
    try:
        if value is None or pd.isna(value):
            return default
        
        # Convert to string and clean
        str_value = str(value).strip()
        
        # Handle empty strings
        if not str_value or str_value.lower() in ['', 'n/a', 'na', 'null', 'none']:
            return default
            
        # Handle backslashes and other invalid characters
        if str_value in ['\\', '\\\\', '/', '#N/A', '#VALUE!', '#ERROR!', '#DIV/0!']:
            logger.warning(f"Invalid numeric value encountered: '{str_value}', using default: {default}")
            return default
            
        # Remove commas and currency symbols but preserve decimal points and negative signs
        str_value = re.sub(r'[^\d\.\-]', '', str_value)
        
        # Handle multiple decimal points or other edge cases
        if str_value.count('.') > 1 or str_value.count('-') > 1:
            return default
            
        # Try to convert to float
        return float(str_value) if str_value and str_value not in ['.', '-', '-.'] else default
        
    except (ValueError, TypeError) as e:
        logger.warning(f"Float conversion error for value '{value}': {e}, using default: {default}")
        return default

def get_db_connection():
    """Creates a new connection to the PostgreSQL database."""
    conn = None
    try:
        # Check for both DB_PASS and DB_PASSWORD for compatibility
        password = os.getenv("DB_PASS") or os.getenv("DB_PASSWORD")
        conn = psycopg2.connect(
            dbname=os.getenv("DB_NAME"),
            user=os.getenv("DB_USER"),
            password=password,
            host=os.getenv("DB_HOST", "localhost"), # Default to localhost for Windows run
            port=os.getenv("DB_PORT", "5432")
        )
        return conn
    except psycopg2.OperationalError as e:
        logger.error(f"FATAL: Could not connect to database: {e}")
        raise


def generate_table_csv(idea_ids=None) -> str:
    """
    Export the current contents of the PostgreSQL 'ideas' table to a CSV file
    in the UPLOAD_FOLDER and return the filename.

    The optional idea_ids parameter is currently ignored and all ideas are exported.
    """
    conn = None
    try:
        conn = get_db_connection()
        # Simple snapshot of all ideas; adjust columns if needed
        query = "SELECT * FROM ideas"
        df = pd.read_sql(query, conn)

        if df.empty:
            logger.info("generate_table_csv: 'ideas' table is empty, nothing to export.")
            return ""

        ts = datetime.now().strftime("%Y%m%d_%H%M%S")
        filename = f"ideas_snapshot_{ts}.csv"
        csv_path = UPLOAD_FOLDER / filename

        df.to_csv(csv_path, index=False)
        logger.info(f"Exported ideas snapshot to {csv_path}")
        return filename
    except Exception as e:
        logger.error(f"generate_table_csv error: {e}")
        return ""
    finally:
        if conn is not None:
            conn.close()

def get_user_db_conn():
    """Creates a new connection to the SQLite users.db."""
    conn = sqlite3.connect(USER_DB_PATH)
    conn.row_factory = sqlite3.Row
    return conn


def init_user_db_schema():
    """
    Ensure users.db has the expected schema and default roles.
    This is idempotent and safe to call on every startup.
    """
    conn = get_user_db_conn()
    cursor = conn.cursor()

    # Base table (includes role and created_at)
    cursor.execute(
        """
        CREATE TABLE IF NOT EXISTS users (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            username TEXT UNIQUE NOT NULL,
            password_hash TEXT NOT NULL,
            role TEXT DEFAULT 'User',
            created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
        );
        """
    )

    # Add missing columns for older databases
    try:
        cursor.execute("ALTER TABLE users ADD COLUMN role TEXT DEFAULT 'User'")
    except sqlite3.OperationalError:
        pass  # Column already exists

    try:
        cursor.execute("ALTER TABLE users ADD COLUMN created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP")
    except sqlite3.OperationalError:
        pass  # Column already exists

    # Upgrade default users to proper roles if they exist
    try:
        cursor.execute(
            "UPDATE users SET role = 'SuperAdmin' WHERE username IN ('superadmin', 'superuser') AND (role IS NULL OR role = 'User')"
        )
        cursor.execute(
            "UPDATE users SET role = 'Admin' WHERE username = 'admin' AND (role IS NULL OR role = 'User')"
        )
    except sqlite3.OperationalError:
        pass

    # Ensure at least one SuperAdmin exists
    cursor.execute("SELECT COUNT(*) FROM users WHERE role = 'SuperAdmin'")
    count_superadmin = cursor.fetchone()[0] or 0
    if count_superadmin == 0:
        try:
            default_username = "superadmin"
            default_password = "admin123"
            hashed_pw = generate_password_hash(default_password)
            cursor.execute(
                "INSERT INTO users (username, password_hash, role) VALUES (?, ?, ?)",
                (default_username, hashed_pw, "SuperAdmin"),
            )
            logger.info("Created default SuperAdmin user: superadmin / admin123")
        except sqlite3.IntegrityError:
            logger.info("SuperAdmin user already exists.")

    conn.commit()
    conn.close()


def log_event(event_type: str, username: str | None = None, payload: dict | None = None):
    """
    Append an event into PostgreSQL 'events' table.
    This simulates an event-stream / Kafka topic for later analytics.
    """
    try:
        conn = get_db_connection()
        cur = conn.cursor()
        cur.execute(
            """
            INSERT INTO events (event_type, username, payload)
            VALUES (%s, %s, %s)
            """,
            (event_type, username, json.dumps(payload or {})),
        )
        conn.commit()
        cur.close()
        conn.close()
    except Exception as e:
        logger.error(f"Failed to log event {event_type}: {e}")

# Database Setup
'''def init_db():
    """Initialize SQLite database and create or migrate ideas table."""
    try:
        with sqlite3.connect(DB_PATH) as conn:
            cursor = conn.cursor()
            # Define expected table schema based on Excel columns
            expected_columns = [
                'id', 'idea_id', 'cost_reduction_idea', 'reason', 'mgi_gross_saving',
                'estimated_cost_savings', 'saving_value_inr', 'weight_saving', 'group_id', 'status',
                'way_forward', 'dept', 'target_date', 'kd_lc', 'wtd_avg', 'est_impl_date',
                'investment_cr', 'capex', 'mgi_carline', 'benchmarking_carline', 'mg_product_scenario',
                'competitor_product_scenario', 'purpose_mg_product', 'purpose_competitor_product',
                'impact_other_systems', 'client_statement', 'cae_required', 'homologation_required',
                'styling_change', 'part_level_testing', 'assembly_trials', 'cad_drawing_update',
                'ecn_required', 'part_production_trials', 'vehicle_level_testing', 'idea_generated_by',
                'new_tool_required', 'new_tool_cost', 'tool_modification_required', 'tool_modification_cost',
                'variants', 'current_status', 'resp', 'mix', 'volume', 'purchase_proposal',
                'interest', 'payback_months', 'mgi_pe_feasibility', 'homeroom_approval', 'pp_approval',
                'supplier_feasibility', 'financial_feasibility', 'proposal_image_filename',
                'mg_vehicle_image', 'competitor_vehicle_image', 'user_input', 'response_text',
                'image_path', 'created_at'
            ]

            column_definitions = [
                'id INTEGER PRIMARY KEY AUTOINCREMENT',
                'idea_id TEXT', 'cost_reduction_idea TEXT', 'reason TEXT', 'mgi_gross_saving REAL',
                'estimated_cost_savings TEXT', 'saving_value_inr REAL', 'weight_saving REAL',
                'group_id TEXT', 'status TEXT', 'way_forward TEXT', 'dept TEXT', 'target_date TEXT',
                'kd_lc TEXT', 'wtd_avg REAL', 'est_impl_date TEXT', 'investment_cr REAL', 'capex REAL',
                'mgi_carline TEXT', 'benchmarking_carline TEXT', 'mg_product_scenario TEXT',
                'competitor_product_scenario TEXT', 'purpose_mg_product TEXT', 'purpose_competitor_product TEXT',
                'impact_other_systems TEXT', 'client_statement TEXT', 'cae_required TEXT',
                'homologation_required TEXT', 'styling_change TEXT', 'part_level_testing TEXT',
                'assembly_trials TEXT', 'cad_drawing_update TEXT', 'ecn_required TEXT',
                'part_production_trials TEXT', 'vehicle_level_testing TEXT', 'idea_generated_by TEXT',
                'new_tool_required TEXT', 'new_tool_cost REAL', 'tool_modification_required TEXT',
                'tool_modification_cost REAL', 'variants TEXT', 'current_status TEXT', 'resp TEXT',
                'mix TEXT', 'volume TEXT', 'purchase_proposal TEXT', 'interest TEXT',
                'payback_months TEXT', 'mgi_pe_feasibility TEXT', 'homeroom_approval TEXT',
                'pp_approval TEXT', 'supplier_feasibility TEXT', 'financial_feasibility TEXT',
                'proposal_image_filename TEXT', 'mg_vehicle_image TEXT', 'competitor_vehicle_image TEXT',
                'user_input TEXT', 'response_text TEXT', 'image_path TEXT',
                'created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP'
            ]

            # Check if ideas table exists
            cursor.execute("SELECT name FROM sqlite_master WHERE type='table' AND name='ideas'")
            table_exists = cursor.fetchone()

            if not table_exists:
                # Create new table
                cursor.execute(f"""
                    CREATE TABLE ideas (
                        {', '.join(column_definitions)}
                    )
                """)
                conn.commit()
                logger.info("Created new ideas table with correct schema")
                return

            # Verify existing table schema
            cursor.execute("PRAGMA table_info(ideas)")
            columns = [col[1] for col in cursor.fetchall()]
            if 'capex' not in columns:
                try:
                    cursor.execute("ALTER TABLE ideas ADD COLUMN capex REAL")
                    conn.commit()
                    columns.append('capex')
                except Exception:
                    pass
            if set(columns) == set(expected_columns):
                logger.info("Existing ideas table schema matches expected schema")
                return

            logger.warning(f"Schema mismatch: expected {expected_columns}, found {columns}")

            # Create new table with correct schema
            cursor.execute("DROP TABLE IF EXISTS ideas_new")
            cursor.execute(f"""
                CREATE TABLE ideas_new (
                    {', '.join(column_definitions)}
                )
            """)

            # Get common columns for data migration
            common_columns = [col for col in columns if col in expected_columns and col != 'id']

            if common_columns:
                # Copy data from old table to new table
                cursor.execute(f"""
                    INSERT INTO ideas_new ({', '.join(common_columns)})
                    SELECT {', '.join(common_columns)} FROM ideas
                """)
                conn.commit()
                logger.info(f"Migrated data for columns: {common_columns}")

            # Drop old table and rename new table
            cursor.execute("DROP TABLE ideas")
            cursor.execute("ALTER TABLE ideas_new RENAME TO ideas")
            conn.commit()
            logger.info("Successfully migrated ideas table to new schema")

            # Re-verify schema
            cursor.execute("PRAGMA table_info(ideas)")
            columns = [col[1] for col in cursor.fetchall()]
            if set(columns) != set(expected_columns):
                raise ValueError(f"Database schema migration failed: expected {expected_columns}, found {columns}")

            logger.info(f"Database initialized at {DB_PATH} with correct schema")

    except Exception as e:
        logger.error(f"Error initializing database: {str(e)}")
        raise'''

login_manager = LoginManager()
login_manager.init_app(app)
login_manager.login_view = 'login' # Redirect to /login if @login_required fails
login_manager.login_message_category = 'error'

def role_required(allowed_roles):
    def decorator(f):
        @wraps(f)
        def decorated_function(*args, **kwargs):
            if not current_user.is_authenticated:
                return redirect(url_for('login'))
            if current_user.role not in allowed_roles:
                flash("You do not have permission to access this page.", "error")
                return redirect(url_for('chat_app'))
            return f(*args, **kwargs)
        return decorated_function
    return decorator

class User(UserMixin):
    def __init__(self, id, username, password_hash, role='User'):
        self.id = id
        self.username = username
        self.password_hash = password_hash
        self.role = role

    def is_superadmin(self):
        return self.role == 'SuperAdmin'
    
    def is_admin(self):
        return self.role in ['SuperAdmin', 'Admin']

@login_manager.user_loader
def load_user(user_id):
    conn = get_user_db_conn()
    cursor = conn.cursor()
    cursor.execute("SELECT * FROM users WHERE id = ?", (user_id,))
    user_row = cursor.fetchone()
    conn.close()
    
    if user_row:
        # Default to 'User' if role column is missing during migration
        role = user_row['role'] if 'role' in user_row.keys() else 'User'
        return User(user_row['id'], user_row['username'], user_row['password_hash'], role)
    return None

def init_db():
    """Check if production database tables exist."""
    conn = None
    cursor = None
    try:
        conn = get_db_connection()
        cursor = conn.cursor()
        
        # Check if 'ideas' table exists
        cursor.execute("SELECT to_regclass('public.ideas');")
        table_exists = cursor.fetchone()[0]

        if not table_exists:
            logger.error("Database table 'ideas' not found.")
            logger.error("Please run the 'create_database.py' script first.")
            raise ValueError("Database not populated. Run 'create_database.py'.")
            
        logger.info("Database connection and 'ideas' table verified successfully.")

        # --- Big-data related tables in PostgreSQL ---

        # Track knowledge base uploads (for admin dashboard)
        cursor.execute(
            """
            CREATE TABLE IF NOT EXISTS data_uploads (
                id SERIAL PRIMARY KEY,
                excel_filename TEXT,
                zip_filename TEXT,
                uploaded_by TEXT,
                status TEXT DEFAULT 'active',
                message TEXT,
                records_loaded INTEGER,
                created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
            );
            """
        )

        # Event stream (simulated Kafka topic)
        cursor.execute(
            """
            CREATE TABLE IF NOT EXISTS events (
                id SERIAL PRIMARY KEY,
                event_type TEXT NOT NULL,
                username TEXT,
                payload JSONB,
                created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
            );
            """
        )

        # Simple feature store for personalization
        cursor.execute(
            """
            CREATE TABLE IF NOT EXISTS user_features (
                username TEXT PRIMARY KEY,
                preferred_depts TEXT[],
                avg_query_length REAL,
                avg_response_length REAL,
                active_hours TEXT[],
                last_updated TIMESTAMP DEFAULT CURRENT_TIMESTAMP
            );
            """
        )

        # Time-series metrics (simulated TimescaleDB)
        cursor.execute(
            """
            CREATE TABLE IF NOT EXISTS metrics_timeseries (
                id SERIAL PRIMARY KEY,
                metric_name TEXT NOT NULL,
                metric_value REAL,
                labels JSONB,
                ts TIMESTAMP DEFAULT CURRENT_TIMESTAMP
            );
            """
        )

        # Idea index (search engine style metadata)
        cursor.execute(
            """
            CREATE TABLE IF NOT EXISTS idea_index (
                id SERIAL PRIMARY KEY,
                idea_id TEXT,
                title TEXT,
                dept TEXT,
                carline TEXT,
                tags TEXT[],
                last_updated TIMESTAMP DEFAULT CURRENT_TIMESTAMP
            );
            """
        )

        # Knowledge graph style relations: problem -> idea
        cursor.execute(
            """
            CREATE TABLE IF NOT EXISTS problem_idea_relations (
                id SERIAL PRIMARY KEY,
                problem_id TEXT,
                idea_id TEXT,
                relation_type TEXT,
                created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
            );
            """
        )

        # Ensure chat_history table exists for logging and dashboard/history views
        cursor.execute(
            """
            CREATE TABLE IF NOT EXISTS chat_history (
                id SERIAL PRIMARY KEY,
                username TEXT NOT NULL,
                user_query TEXT,
                response_text TEXT,
                image_path TEXT,
                created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
            );
            """
        )
            
    except Exception as e:
        logger.error(f"Error initializing database: {str(e)}")
        raise
    finally:
        if cursor:
            cursor.close()
        if conn:
            conn.close()

# Model and Vector DB Setup
def setup_model():
    """Load GPT-2 model and tokenizer."""
    try:
        global tokenizer, model

        ensure_model_dir()

        tok_dir = find_gpt2_tokenizer_dir(MODEL_DIR)
        if tok_dir:
            try:
                tokenizer = AutoTokenizer.from_pretrained(tok_dir, local_files_only=True)
                logger.info(f"Tokenizer loaded from {tok_dir}")
            except Exception as e:
                logger.warning(f"AutoTokenizer local load failed: {e}")
                tokenizer = None
        else:
            logger.warning("No local GPT-2 tokenizer files found (vocab.json, merges.txt)")
            tokenizer = None

        if tokenizer is None:
            try:
                tokenizer = AutoTokenizer.from_pretrained('gpt2')
                logger.info("Fetched GPT-2 tokenizer from Hugging Face")
                # Save tokenizer files locally to avoid future warnings
                try:
                    tokenizer.save_pretrained(MODEL_DIR / "gpt2_tokenizer")
                    logger.info(f"Saved GPT-2 tokenizer to {MODEL_DIR / 'gpt2_tokenizer'}")
                except Exception as save_err:
                    logger.warning(f"Could not save tokenizer locally: {save_err}")
            except Exception as e:
                logger.error(f"Failed to load any tokenizer: {e}")
                raise

        if not TORCHSCRIPT_MODEL_PATH.exists():
            # Try to locate a .pt in MODEL_DIR
            pt_candidates = list(MODEL_DIR.glob('**/*.pt'))
            if pt_candidates:
                # Prefer file named gpt2.pt else take first
                preferred = next((p for p in pt_candidates if p.name.lower() == 'gpt2.pt'), pt_candidates[0])
                logger.warning(f"TorchScript model not at expected path; using {preferred}")
                model_path = preferred
            else:
                raise FileNotFoundError(f"TorchScript model not found in {MODEL_DIR}")
        else:
            model_path = TORCHSCRIPT_MODEL_PATH

        model = torch.jit.load(model_path, map_location=device)
        model.eval()
        logger.info(f"TorchScript GPT-2 model loaded from {model_path} on {device}")
    except Exception as e:
        logger.error(f"Error loading GPT-2 model: {str(e)}")
        raise

def normalize_text(text):
    """Normalize text for better similarity matching."""
    try:
        text = str(text).lower()
        text = re.sub(r'[^\w\s]', '', text)
        tokens = word_tokenize(text)
        stop_words = set(stopwords.words('english'))
        tokens = [t for t in tokens if t not in stop_words]
        return ' '.join(tokens)
    except Exception as e:
        logger.error(f"Error normalizing text: {str(e)}")
        return str(text).lower()

'''def setup_vlm():
    """Load BLIP VLM and process Excel data."""
    global vlm_processor, vlm_model, embedding_model, proposal_embeddings, valid_df, vlm_initialized

    with vlm_init_lock:
        if vlm_initialized:
            logger.info("VLM already initialized, skipping setup.")
            return

        try:
            vlm_processor = BlipProcessor.from_pretrained(VLM_MODEL_PATH, local_files_only=True)
            vlm_model = BlipForConditionalGeneration.from_pretrained(VLM_MODEL_PATH, local_files_only=True).to(device)
            vlm_model.eval()
            logger.info(f"VLM model (BLIP) loaded from {VLM_MODEL_PATH} on {device}")

            embedding_model = SentenceTransformer('all-MiniLM-L6-v2', device=device)

            df = pd.read_excel(DATA_PATH)
            logger.info(f"Loaded Excel file with {len(df)} rows")

            # Map Excel column names to database column names
            column_mapping = {
                'Idea Id': 'idea_id',
                'Cost Reduction Idea Proposal': 'cost_reduction_idea',
                'Reason of cost reduction proposal': 'reason',
                'MGI Estimated Gross saving': 'mgi_gross_saving',
                'Estimated Cost Savings - Breakup': 'estimated_cost_savings',
                'Saving Value(INR)': 'saving_value_inr',
                'Weight Saving(Kg)': 'weight_saving',
                'Group ID': 'group_id',
                'Status (OK/TBD/NG)': 'status',
                'Way Forward': 'way_forward',
                'Dept': 'dept',
                'Target Date': 'target_date',
                'KD/LC': 'kd_lc',
                'Wtd. Avg.': 'wtd_avg',
                'Est. Impl. Date': 'est_impl_date',
                'Investment (Cr)': 'investment_cr',
                'CAPEX': 'capex',
                'MGI carline': 'mgi_carline',
                'Banchmarking carline': 'benchmarking_carline',
                'MG Product Scenario': 'mg_product_scenario',
                'Competitor Product Scenario': 'competitor_product_scenario',
                'Purpose on MG product': 'purpose_mg_product',
                'Purpose on Competitor Product': 'purpose_competitor_product',
                'Impact on Other Systems': 'impact_other_systems',
                'Client Proposal/Rejection Statement': 'client_statement',
                'CAE Required (No/Yes)': 'cae_required',
                'Homologation Required (No/Yes)': 'homologation_required',
                'Styling Change (No/Yes)': 'styling_change',
                'Part Level Testing (No/Yes)': 'part_level_testing',
                'Assembly Trials (No/Yes)': 'assembly_trials',
                'CAD Drawing Update (No/Yes)': 'cad_drawing_update',
                'ECN Required (No/Yes)': 'ecn_required',
                'Part Production Trials (No/Yes)': 'part_production_trials',
                'Vehicle Level Testing (No/Yes)': 'vehicle_level_testing',
                'Idea Generated By': 'idea_generated_by',
                'New Tool Required (No/Yes)': 'new_tool_required',
                'New Tool Cost': 'new_tool_cost',
                'Tool Modification Required (No/Yes)': 'tool_modification_required',
                'Tool Modification Cost': 'tool_modification_cost',
                'Variants': 'variants',
                'Current Status': 'current_status',
                'Resp': 'resp',
                'Mix': 'mix',
                'Volume': 'volume',
                'Purchase Praposal': 'purchase_proposal',
                'Interest': 'interest',
                'Payback months': 'payback_months',
                'MGI PE Feasibilty': 'mgi_pe_feasibility',
                'Homeroom Approval': 'homeroom_approval',
                'PP Approval': 'pp_approval',
                'Supplier feasibility': 'supplier_feasibility',
                'Financial Feasibility': 'financial_feasibility',
                'Proposal Image': 'proposal_image_filename',
                'MG Vehicle Image': 'mg_vehicle_image',
                'Competitor Vehicle Image': 'competitor_vehicle_image'
            }

            df.rename(columns=column_mapping, inplace=True)
            df['normalized_cost_reduction_idea'] = df['cost_reduction_idea'].apply(normalize_text)
            df = df[df['cost_reduction_idea'].notnull()].reset_index(drop=True)

            image_files = sorted([f.name for f in IMAGE_DIRS["proposal"].glob("*.jpg")])
            df['proposal_image_filename'] = image_files[:len(df)]

            def file_exists(row):
                fname = row['proposal_image_filename']
                exists = fname in image_files
                if not exists:
                    logger.warning(f"Missing proposal image file: {fname}")
                return exists

            global valid_df
            valid_df = df[df.apply(file_exists, axis=1)].reset_index(drop=True)

            if valid_df.empty:
                logger.error("No valid data found after filtering")
                vlm_initialized = False
                return

            proposal_texts = valid_df['cost_reduction_idea'].apply(str).tolist()
            global proposal_embeddings
            proposal_embeddings = embedding_model.encode(proposal_texts, convert_to_tensor=True, device=device)
            logger.info(f"Generated embeddings for {len(proposal_texts)} proposals")

            # Populate database with Excel data
            with db_lock:
                with sqlite3.connect(DB_PATH) as conn:
                    cursor = conn.cursor()
                    for _, row in valid_df.iterrows():
                        cursor.execute("""
                            INSERT OR IGNORE INTO ideas (
                                idea_id, cost_reduction_idea, reason, mgi_gross_saving,
                                estimated_cost_savings, saving_value_inr, weight_saving,
                                group_id, status, way_forward, dept, target_date, kd_lc,
                                wtd_avg, est_impl_date, investment_cr, capex, mgi_carline,
                                benchmarking_carline, mg_product_scenario, competitor_product_scenario,
                                purpose_mg_product, purpose_competitor_product, impact_other_systems,
                                client_statement, cae_required, homologation_required, styling_change,
                                part_level_testing, assembly_trials, cad_drawing_update, ecn_required,
                                part_production_trials, vehicle_level_testing, idea_generated_by,
                                new_tool_required, new_tool_cost, tool_modification_required,
                                tool_modification_cost, variants, current_status, resp, mix,
                                volume, purchase_proposal, interest, payback_months, mgi_pe_feasibility,
                                homeroom_approval, pp_approval, supplier_feasibility, financial_feasibility,
                                proposal_image_filename, mg_vehicle_image, competitor_vehicle_image
                            ) VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?)
                        """, (
                            str(row.get('idea_id', 'N/A')),
                            str(row.get('cost_reduction_idea', 'N/A')),
                            str(row.get('reason', 'N/A')),
                            safe_float_convert(row.get('mgi_gross_saving')),  # FIXED: Using safe conversion
                            str(row.get('estimated_cost_savings', 'N/A')),
                            safe_float_convert(row.get('saving_value_inr')),  # FIXED: Using safe conversion
                            safe_float_convert(row.get('weight_saving')),     # FIXED: Using safe conversion
                            str(row.get('group_id', 'N/A')),
                            str(row.get('status', 'N/A')),
                            str(row.get('way_forward', 'N/A')),
                            str(row.get('dept', 'N/A')),
                            str(row.get('target_date', 'N/A')),
                            str(row.get('kd_lc', 'N/A')),
                            safe_float_convert(row.get('wtd_avg')),           # FIXED: Using safe conversion
                            str(row.get('est_impl_date', 'N/A')),
                            safe_float_convert(row.get('investment_cr')),     # FIXED: Using safe conversion
                            safe_float_convert(row.get('capex') or row.get('investment_cr')),
                            str(row.get('mgi_carline', 'N/A')),
                            str(row.get('benchmarking_carline', 'N/A')),
                            str(row.get('mg_product_scenario', 'N/A')),
                            str(row.get('competitor_product_scenario', 'N/A')),
                            str(row.get('purpose_mg_product', 'N/A')),
                            str(row.get('purpose_competitor_product', 'N/A')),
                            str(row.get('impact_other_systems', 'N/A')),
                            str(row.get('client_statement', 'N/A')),
                            str(row.get('cae_required', 'N/A')),
                            str(row.get('homologation_required', 'N/A')),
                            str(row.get('styling_change', 'N/A')),
                            str(row.get('part_level_testing', 'N/A')),
                            str(row.get('assembly_trials', 'N/A')),
                            str(row.get('cad_drawing_update', 'N/A')),
                            str(row.get('ecn_required', 'N/A')),
                            str(row.get('part_production_trials', 'N/A')),
                            str(row.get('vehicle_level_testing', 'N/A')),
                            str(row.get('idea_generated_by', 'N/A')),
                            str(row.get('new_tool_required', 'N/A')),
                            safe_float_convert(row.get('new_tool_cost')),     # FIXED: Using safe conversion
                            str(row.get('tool_modification_required', 'N/A')),
                            safe_float_convert(row.get('tool_modification_cost')),  # FIXED: Using safe conversion
                            str(row.get('variants', 'N/A')),
                            str(row.get('current_status', 'N/A')),
                            str(row.get('resp', 'N/A')),
                            str(row.get('mix', 'N/A')),
                            str(row.get('volume', 'N/A')),
                            str(row.get('purchase_proposal', 'N/A')),
                            str(row.get('interest', 'N/A')),
                            str(row.get('payback_months', 'N/A')),
                            str(row.get('mgi_pe_feasibility', 'N/A')),
                            str(row.get('homeroom_approval', 'N/A')),
                            str(row.get('pp_approval', 'N/A')),
                            str(row.get('supplier_feasibility', 'N/A')),
                            str(row.get('financial_feasibility', 'N/A')),
                            row.get('proposal_image_filename'),
                            row.get('mg_vehicle_image'),
                            row.get('competitor_vehicle_image')
                        ))
                    conn.commit()
                    logger.info(f"Populated database with {len(valid_df)} ideas from Excel")

            vlm_initialized = True
            logger.info(f"VLM setup complete with {len(valid_df)} valid entries")

        except Exception as e:
            logger.error(f"Error setting up VLM: {str(e)}")
            vlm_processor = None
            vlm_model = None
            valid_df = None
            proposal_embeddings = None
            vlm_initialized = False'''

# SQL - without endpoint

'''def setup_vlm():
    """Load BLIP VLM and build embeddings from the database."""
    global vlm_processor, vlm_model, embedding_model, proposal_embeddings, valid_df, vlm_initialized

    with vlm_init_lock:
        if vlm_initialized:
            logger.info("VLM already initialized, skipping setup.")
            return

        try:
            # 1. Load models (this is unchanged)
            vlm_processor = BlipProcessor.from_pretrained(VLM_MODEL_PATH, local_files_only=True)
            vlm_model = BlipForConditionalGeneration.from_pretrained(VLM_MODEL_PATH, local_files_only=True).to(device)
            vlm_model.eval()
            logger.info(f"VLM model (BLIP) loaded from {VLM_MODEL_PATH} on {device}")

            embedding_model = SentenceTransformer('all-MiniLM-L6-v2', device=device)

            # 2. Read data FROM the database (This replaces the Excel/pandas logic)
            logger.info("Loading VLM data from database...")
            with db_lock:
                with sqlite3.connect(DB_PATH) as conn:
                    # We only need data that has an associated proposal image for VLM
                    query = "SELECT * FROM ideas WHERE proposal_image_filename IS NOT NULL AND cost_reduction_idea IS NOT NULL"
                    valid_df = pd.read_sql_query(query, conn)
            
            if valid_df.empty:
                logger.error("No valid data with proposal images found in database for VLM setup. Database might be empty or data is missing image filenames.")
                vlm_initialized = False
                return

            # 3. Generate embeddings from the database data
            proposal_texts = valid_df['cost_reduction_idea'].apply(str).tolist()
            global proposal_embeddings
            proposal_embeddings = embedding_model.encode(proposal_texts, convert_to_tensor=True, device=device)
            logger.info(f"Generated VLM embeddings for {len(proposal_texts)} proposals from database")

            vlm_initialized = True
            logger.info(f"VLM setup complete with {len(valid_df)} valid entries from database")

        except Exception as e:
            logger.error(f"Error setting up VLM: {str(e)}")
            vlm_processor = None
            vlm_model = None
            valid_df = None
            proposal_embeddings = None
            vlm_initialized = False'''

def setup_vlm():
    """Load BLIP VLM and build embeddings from the database."""
    global vlm_processor, vlm_model, embedding_model, proposal_embeddings, valid_df, vlm_initialized

    with vlm_init_lock:
        if vlm_initialized:
            logger.info("VLM already initialized, skipping setup.")
            return

        conn = None
        try:
            # 0. Validate local VLM directory presence and required files
            if not VLM_MODEL_PATH.exists():
                logger.warning(f"VLM path missing: {VLM_MODEL_PATH}. Skipping VLM initialization.")
                vlm_initialized = False
                return

            preproc = VLM_MODEL_PATH / 'preprocessor_config.json'
            if not preproc.exists():
                logger.warning(f"VLM preprocessor_config.json not found in {VLM_MODEL_PATH}. Skipping VLM initialization.")
                vlm_initialized = False
                return

            # 1. Load models
            vlm_processor = BlipProcessor.from_pretrained(VLM_MODEL_PATH, local_files_only=True)
            vlm_model = BlipForConditionalGeneration.from_pretrained(VLM_MODEL_PATH, local_files_only=True).to(device)
            vlm_model.eval()
            logger.info(f"VLM model (BLIP) loaded from {VLM_MODEL_PATH} on {device}")

            embedding_model = SentenceTransformer('all-MiniLM-L6-v2', device=device)

            # 2. Read data FROM the database
            logger.info("Loading VLM data from database...")
            conn = get_db_connection() # <-- NEW
            query = "SELECT * FROM ideas WHERE proposal_image_filename IS NOT NULL AND cost_reduction_idea IS NOT NULL"
            valid_df = pd.read_sql_query(query, conn) # <-- Pandas reads directly
            
            if valid_df.empty:
                logger.error("No valid data with proposal images found in database for VLM setup.")
                vlm_initialized = False
                return

            # 3. Generate embeddings
            proposal_texts = valid_df['cost_reduction_idea'].apply(str).tolist()
            global proposal_embeddings
            proposal_embeddings = embedding_model.encode(proposal_texts, convert_to_tensor=True, device=device)
            logger.info(f"Generated VLM embeddings for {len(proposal_texts)} proposals from database")

            vlm_initialized = True
            logger.info(f"VLM setup complete with {len(valid_df)} valid entries from database")

        except Exception as e:
            logger.error(f"Error setting up VLM: {str(e)}")
            vlm_processor, vlm_model, valid_df, proposal_embeddings = None, None, None, None
            vlm_initialized = False
        finally:
            if conn:
                conn.close()

'''def build_vector_db():
    """Build FAISS index from SQLite database."""
    global embedding_model, faiss_index, idea_texts, idea_rows

    try:
        logger.info("Building vector database...")

        if embedding_model is None:
            embedding_model = SentenceTransformer('all-MiniLM-L6-v2', device=device)
            logger.info("SentenceTransformer model initialized successfully")

        faiss_index = faiss.read_index(str(FAISS_INDEX_PATH))
        logger.info(f"FAISS index loaded from {FAISS_INDEX_PATH}")

        with db_lock:
            with sqlite3.connect(DB_PATH) as conn:
                # Verify table exists and has required columns
                cursor = conn.cursor()
                cursor.execute("PRAGMA table_info(ideas)")
                columns = [col[1] for col in cursor.fetchall()]

                required_columns = [
                    'idea_id', 'cost_reduction_idea', 'reason', 'mgi_gross_saving',
                    'estimated_cost_savings', 'saving_value_inr', 'weight_saving', 'group_id', 'status',
                    'way_forward', 'dept', 'target_date', 'kd_lc', 'wtd_avg', 'est_impl_date',
                    'investment_cr', 'mgi_carline', 'benchmarking_carline', 'mg_product_scenario',
                    'competitor_product_scenario', 'purpose_mg_product', 'purpose_competitor_product',
                    'impact_other_systems', 'client_statement', 'cae_required', 'homologation_required',
                    'styling_change', 'part_level_testing', 'assembly_trials', 'cad_drawing_update',
                    'ecn_required', 'part_production_trials', 'vehicle_level_testing', 'idea_generated_by',
                    'new_tool_required', 'new_tool_cost', 'tool_modification_required', 'tool_modification_cost',
                    'variants', 'current_status', 'resp', 'mix', 'volume', 'purchase_proposal',
                    'interest', 'payback_months', 'mgi_pe_feasibility', 'homeroom_approval', 'pp_approval',
                    'supplier_feasibility', 'financial_feasibility', 'proposal_image_filename',
                    'mg_vehicle_image', 'competitor_vehicle_image'
                ]

                if not all(col in columns for col in required_columns):
                    missing = [col for col in required_columns if col not in columns]
                    logger.error(f"Missing columns in ideas table: {missing}")
                    raise ValueError(f"Missing columns in ideas table: {missing}")

                query = f"""
                    SELECT {', '.join(required_columns)}
                    FROM ideas
                    WHERE idea_id IS NOT NULL AND cost_reduction_idea IS NOT NULL
                """
                df = pd.read_sql_query(query, conn)
                logger.info(f"Retrieved {len(df)} ideas from SQLite database")

                if df.empty:
                    logger.warning("No valid ideas found in the database")
                    # Populate with Excel data as fallback
                    setup_vlm()
                    with sqlite3.connect(DB_PATH) as conn:
                        df = pd.read_sql_query(query, conn)
                        logger.info(f"Retrieved {len(df)} ideas after fallback population")

                    if df.empty:
                        raise ValueError("No valid ideas found in the database after fallback")

                idea_texts.clear()
                idea_rows.clear()

                for _, row in df.iterrows():
                    text = ' '.join(str(row.get(col, '')) for col in df.columns if pd.notna(row.get(col)))
                    if text.strip():
                        idea_texts.append(text)
                        idea_rows.append(row)

                logger.info(f"Vector database built with {len(idea_texts)} ideas")

    except Exception as e:
        logger.error(f"Error building vector DB: {str(e)}")
        raise'''

# SQL without Endpoints
'''def build_vector_db():
    """Build FAISS index from SQLite database."""
    global embedding_model, faiss_index, idea_texts, idea_rows

    try:
        logger.info("Building vector database...")

        if embedding_model is None:
            embedding_model = SentenceTransformer('all-MiniLM-L6-v2', device=device)
            logger.info("SentenceTransformer model initialized successfully")

        faiss_index = faiss.read_index(str(FAISS_INDEX_PATH))
        logger.info(f"FAISS index loaded from {FAISS_INDEX_PATH}")

        with db_lock:
            with sqlite3.connect(DB_PATH) as conn:
                # Verify table exists and has required columns
                cursor = conn.cursor()
                cursor.execute("PRAGMA table_info(ideas)")
                columns = [col[1] for col in cursor.fetchall()]

                required_columns = [
                    'idea_id', 'cost_reduction_idea', 'reason', 'mgi_gross_saving',
                    'estimated_cost_savings', 'saving_value_inr', 'weight_saving', 'group_id', 'status',
                    'way_forward', 'dept', 'target_date', 'kd_lc', 'wtd_avg', 'est_impl_date',
                    'investment_cr', 'mgi_carline', 'benchmarking_carline', 'mg_product_scenario',
                    'competitor_product_scenario', 'purpose_mg_product', 'purpose_competitor_product',
                    'impact_other_systems', 'client_statement', 'cae_required', 'homologation_required',
                    'styling_change', 'part_level_testing', 'assembly_trials', 'cad_drawing_update',
                    'ecn_required', 'part_production_trials', 'vehicle_level_testing', 'idea_generated_by',
                    'new_tool_required', 'new_tool_cost', 'tool_modification_required', 'tool_modification_cost',
                    'variants', 'current_status', 'resp', 'mix', 'volume', 'purchase_proposal',
                    'interest', 'payback_months', 'mgi_pe_feasibility', 'homeroom_approval', 'pp_approval',
                    'supplier_feasibility', 'financial_feasibility', 'proposal_image_filename',
                    'mg_vehicle_image', 'competitor_vehicle_image'
                ]

                if not all(col in columns for col in required_columns):
                    missing = [col for col in required_columns if col not in columns]
                    logger.error(f"Missing columns in ideas table: {missing}")
                    raise ValueError(f"Missing columns in ideas table: {missing}")

                query = f"""
                    SELECT {', '.join(required_columns)}
                    FROM ideas
                    WHERE idea_id IS NOT NULL AND cost_reduction_idea IS NOT NULL
                """
                df = pd.read_sql_query(query, conn)
                logger.info(f"Retrieved {len(df)} ideas from SQLite database")

                if df.empty:
                    logger.error("No valid ideas found in the database. The database must be populated before running the application.")
                    raise ValueError("No valid ideas found in the database. Run the create_database.py script first.")

                idea_texts.clear()
                idea_rows.clear()

                for _, row in df.iterrows():
                    text = ' '.join(str(row.get(col, '')) for col in df.columns if pd.notna(row.get(col)))
                    if text.strip():
                        idea_texts.append(text)
                        idea_rows.append(row)

                logger.info(f"Vector database built with {len(idea_texts)} ideas")

    except Exception as e:
        logger.error(f"Error building vector DB: {str(e)}")
        raise'''

def build_vector_db():
    """Build FAISS index from PostgreSQL database."""
    global embedding_model, faiss_index, idea_texts, idea_rows

    conn = None
    cursor = None
    try:
        logger.info("Building vector database...")

        if embedding_model is None:
            embedding_model = SentenceTransformer('all-MiniLM-L6-v2', device=device)
            logger.info("SentenceTransformer model initialized successfully")

        try:
            faiss_index = faiss.read_index(str(FAISS_INDEX_PATH))
            logger.info(f"FAISS index loaded from {FAISS_INDEX_PATH}")
        except Exception as e:
            logger.warning(f"Could not load FAISS index: {e}. Building a new one.")
            faiss_index = None

        conn = get_db_connection()
        cursor = conn.cursor(cursor_factory=DictCursor) # <-- Use DictCursor

        required_columns = [
            'idea_id', 'cost_reduction_idea', 'reason', 'mgi_gross_saving',
            'estimated_cost_savings', 'saving_value_inr', 'weight_saving', 'group_id', 'status',
            'way_forward', 'dept', 'target_date', 'kd_lc', 'wtd_avg', 'est_impl_date',
            'investment_cr', 'mgi_carline', 'benchmarking_carline', 'mg_product_scenario',
            'competitor_product_scenario', 'purpose_mg_product', 'purpose_competitor_product',
            'impact_other_systems', 'client_statement', 'cae_required', 'homologation_required',
            'styling_change', 'part_level_testing', 'assembly_trials', 'cad_drawing_update',
            'ecn_required', 'part_production_trials', 'vehicle_level_testing', 'idea_generated_by',
            'new_tool_required', 'new_tool_cost', 'tool_modification_required', 'tool_modification_cost',
            'variants', 'current_status', 'resp', 'mix', 'volume', 'purchase_proposal',
            'interest', 'payback_months', 'mgi_pe_feasibility', 'homeroom_approval', 'pp_approval',
            'supplier_feasibility', 'financial_feasibility', 'proposal_image_filename',
            'mg_vehicle_image', 'competitor_vehicle_image'
        ]
            
        query = f"""
            SELECT {', '.join(required_columns)}
            FROM ideas
            WHERE idea_id IS NOT NULL AND cost_reduction_idea IS NOT NULL
        """
        cursor.execute(query)
        db_rows = cursor.fetchall() # Fetches list of DictRow objects
            
        if not db_rows:
            logger.error("No valid ideas found in the database. Run 'create_database.py' first.")
            raise ValueError("Database is empty.")

        logger.info(f"Retrieved {len(db_rows)} ideas from database")

        idea_texts.clear()
        idea_rows.clear() 

        for row in db_rows:
            text_parts = [
                str(row[col] or '') for col in required_columns 
                if isinstance(row[col], str)
            ]
            text = ' '.join(text_parts)
                
            if text.strip():
                idea_texts.append(normalize_text(text))
                idea_rows.append(dict(row)) # Convert DictRow to a plain dict

        # If FAISS index wasn't loaded, build it
        if faiss_index is None and idea_texts:
            logger.info("Building new FAISS index...")
            text_embeddings = embedding_model.encode(idea_texts, convert_to_numpy=True, device=device, show_progress_bar=True)
            dimension = text_embeddings.shape[1]
            faiss_index = faiss.IndexFlatL2(dimension)
            faiss_index.add(text_embeddings.astype('float32'))
            logger.info("Saving new FAISS index...")
            faiss.write_index(faiss_index, str(FAISS_INDEX_PATH))
            logger.info(f"New FAISS index saved to {FAISS_INDEX_PATH}")
            
        logger.info(f"Vector database built with {len(idea_texts)} ideas")
        
        # Initialize VAVE Agent after vector DB is ready
        global vave_agent
        if vave_agent is None:
            try:
                # Note: DB_PATH points to SQLite, but app uses PostgreSQL
                # You may need to export PostgreSQL data to SQLite or modify tools.py to use PostgreSQL
                vave_agent = VAVEAgent(
                    db_path=str(DB_PATH),
                    vector_db_func=retrieve_context,
                    llm_client=None,  # Set to your LLM client if available
                    pg_conn_func=get_db_connection,  # Use PostgreSQL connection function
                    db_conn=get_db_connection,  # Pass DB connection function for VLM Engine
                    faiss_index=faiss_index,  # Pass FAISS index for VLM Engine
                    sentence_model=embedding_model  # Pass sentence model for VLM Engine
                )
                logger.info("VAVE Agent initialized successfully")
            except Exception as e:
                logger.warning(f"Failed to initialize VAVE Agent: {e}. Agentic RAG will use fallback mode.")

    except Exception as e:
        logger.error(f"Error building vector DB: {str(e)}")
        raise
    finally:
        if cursor:
            cursor.close()
        if conn:
            conn.close()

# Query Parsing and Filtering
def parse_query_filters(query):
    """Parse query to extract filters with improved pattern matching."""
    filters = {}
    query_lower = query.lower()

    # Enhanced cost/saving filters - detect both "more than" and "less than" patterns
    cost_patterns = [
        r'(?:save|saving|cost|price|value).*?(?:more than|greater than|above|at least|over|minimum)\s*(?:inr|rs\.?|₹)?\s*([\d,]+\.?\d*)',
        r'(?:more than|greater than|above|at least|over|minimum)\s*(?:inr|rs\.?|₹)?\s*([\d,]+\.?\d*).*?(?:save|saving|cost|price|value)',
        r'(?:inr|rs\.?|₹)\s*([\d,]+\.?\d*).*?(?:or more|and above|and higher)',
    ]
    
    for pattern in cost_patterns:
        match = re.search(pattern, query_lower)
        if match:
            try:
                cost_val = float(match.group(1).replace(',', ''))
                filters['saving value'] = ('greater', cost_val)
                logger.info(f"Extracted cost filter: greater than INR {cost_val}")
                break
            except ValueError:
                continue

    # Capture explicit cost saving constraints including "less than"
    cost_dir_patterns = [
        r'(?:cost saving|cost savings|saving value|savings|cost reduction).*?(greater|more|above|over|at least|minimum|less|below|under)\s*(?:than)?\s*(?:inr|rs\.?|₹)?\s*([\d,]+\.?\d*)',
        r'(greater|more|above|over|at least|minimum|less|below|under)\s*(?:than)?\s*(?:inr|rs\.?|₹)?\s*([\d,]+\.?\d*)\s*(?:cost saving|cost savings|saving value|savings|cost reduction)',
    ]
    for pattern in cost_dir_patterns:
        match = re.search(pattern, query_lower)
        if match:
            op = match.group(1)
            if op in ['greater', 'more', 'above', 'over', 'at least', 'minimum']:
                op = 'greater'
            elif op in ['less', 'below', 'under']:
                op = 'less'
            try:
                val = float(match.group(2).replace(',', ''))
                filters['saving value'] = (op, val)
                logger.info(f"Extracted cost filter: {op} than INR {val}")
                break
            except ValueError:
                continue

    # Numeric filters with improved patterns
    for field in ['saving value', 'mgi gross saving', 'estimated cost savings', 'investment', 'wtd. avg.', 'new tool cost', 'tool modification cost']:
        if field in query_lower:
            match = re.search(r'(greater|less|more|above|below|under)\s*(?:than)?\s*(?:inr|rs\.?|₹)?\s*([\d,]+\.?\d*)', query_lower)
            if match:
                op = match.group(1)
                if op in ['greater', 'more', 'above']:
                    op = 'greater'
                elif op in ['less', 'below', 'under']:
                    op = 'less'
                filters[field] = (op, float(match.group(2).replace(',', '')))

    if 'weight saving' in query_lower:
        match = re.search(r'(greater|less|more|above|below)\s*(?:than)?\s*([\d.]+)', query_lower)
        if match:
            op = match.group(1)
            if op in ['greater', 'more', 'above']:
                op = 'greater'
            elif op in ['less', 'below']:
                op = 'less'
            filters['weight_saving'] = (op, float(match.group(2)))
        else:
            match = re.search(r'weight saving\s*([\d.]+)', query_lower)
            if match:
                filters['weight_saving'] = ('equal', float(match.group(1)))

    # Status filters
    for status in ['ok', 'tbd', 'ng']:
        if f'status {status}' in query_lower:
            filters['status'] = status.upper()

    # Department filters
    for dept in ['engineering', 'manufacturing', 'procurement', 'design']:
        if dept in query_lower:
            filters['dept'] = dept

    # KD/LC filters
    for kd_lc in ['kd', 'lc']:
        if kd_lc in query_lower:
            filters['kd_lc'] = kd_lc.upper()

    # System filters
    for system in ['brake', 'suspension', 'engine', 'transmission', 'body', 'chassis', 'door', 'paint']:
        if system in query_lower:
            filters['system'] = system

    # Image display filters
    filters['show_proposal'] = 'proposal' in query_lower or 'all' in query_lower
    filters['show_mg'] = 'mg' in query_lower or 'all' in query_lower
    filters['show_competitor'] = 'competitor' in query_lower or 'all' in query_lower

    return filters


# ========================= AUTO ANALYST HELPERS =========================

IDEA_SCHEMA_CACHE = None

COLUMN_SYNONYMS = {
    "idea_id": ["idea_id", "Idea Id", "Idea ID"],
    "cost_reduction_idea": [
        "cost_reduction_idea", "Cost Reduction Idea", "Cost Reduction Idea Proposal",
        "Cost Reduction Idea Proposal Text", "Proposal Summary"
    ],
    "saving_value_inr": ["saving_value_inr", "Saving Value(INR)", "Saving Value (INR)", "Saving Value"],
    "capex": ["capex", "CAPEX", "Capex Investment", "Investment (Cr)", "investment_cr"],
    "dept": ["dept", "Dept", "Department"],
    "status": ["status", "Status", "Status (OK/TBD/NG)"],
    "responsibility": ["resp", "Responsibility", "Owner", "Idea Owner"],
    "weight_saving": ["weight_saving", "Weight Saving(Kg)", "Weight Saving (Kg)"],
    "way_forward": ["way_forward", "Way Forward"],
    "target_date": ["target_date", "Target Date"],
}

DEFAULT_ANALYST_COLUMNS = [
    "Idea Id",
    "Cost Reduction Idea",
    "CAPEX",
    "Saving Value (INR)",
    "Dept",
    "Status"
]

ANALYST_ACTION_KEYWORDS = [
    "reduce", "replace", "remove", "lightweight", "optimize", "integrate"
]

COMPONENT_FALLBACKS = [
    "stabilizer link", "damper tube", "lower control arm", "subframe gusset",
    "brake hose", "headliner overlap", "door map pocket", "wheel arch liner",
    "engine cover", "glove box bin"
]

def _normalize_column(col: str) -> str:
    return (col or "").lower().replace(" ", "").replace("_", "").replace("(", "").replace(")", "")


def _find_column(available_columns, candidate_names):
    normalized_avail = { _normalize_column(col): col for col in available_columns }
    for cand in candidate_names:
        key = _normalize_column(cand)
        if key in normalized_avail:
            return normalized_avail[key]
    return None


def _get_value(row, column_name, default="-"):
    if column_name in row:
        return row.get(column_name, default)
    
    lookup_candidates = COLUMN_SYNONYMS.get(column_name.lower(), [])
    for cand in lookup_candidates:
        for key in row.keys():
            if _normalize_column(key) == _normalize_column(cand):
                return row.get(key, default)
    # As a last resort, try direct normalized match
    for key in row.keys():
        if _normalize_column(key) == _normalize_column(column_name):
            return row.get(key, default)
    return row.get(column_name, default)


def get_idea_schema_columns():
    global IDEA_SCHEMA_CACHE
    if IDEA_SCHEMA_CACHE:
        return IDEA_SCHEMA_CACHE
    conn = None
    cursor = None
    try:
        conn = get_db_connection()
        cursor = conn.cursor()
        cursor.execute("""
            SELECT column_name 
            FROM information_schema.columns 
            WHERE table_name = 'ideas'
            ORDER BY ordinal_position
        """)
        IDEA_SCHEMA_CACHE = [row[0] for row in cursor.fetchall()] or DEFAULT_ANALYST_COLUMNS
    except Exception as exc:
        logger.warning(f"Falling back to default schema due to error: {exc}")
        IDEA_SCHEMA_CACHE = DEFAULT_ANALYST_COLUMNS.copy()
    finally:
        if cursor:
            cursor.close()
        if conn:
            conn.close()
    return IDEA_SCHEMA_CACHE


def auto_select_columns(user_query: str, available_columns: list[str]) -> list[str]:
    selected = []
    schema_candidates = get_idea_schema_columns()
    combined_avail = available_columns or schema_candidates

    def add_column(candidate_key):
        column = _find_column(combined_avail, COLUMN_SYNONYMS.get(candidate_key, [candidate_key]))
        if column and column not in selected:
            selected.append(column)

    # Always start with Idea Id + Cost Reduction Idea if available
    add_column("idea_id")
    add_column("cost_reduction_idea")

    query_lower = user_query.lower()
    keyword_map = {
        "saving": "saving_value_inr",
        "cost": "saving_value_inr",
        "value": "saving_value_inr",
        "weight": "weight_saving",
        "owner": "responsibility",
        "responsible": "responsibility",
        "dept": "dept",
        "department": "dept",
        "status": "status",
        "way forward": "way_forward",
        "timeline": "target_date",
        "target": "target_date",
    }

    for keyword, column_key in keyword_map.items():
        if keyword in query_lower:
            add_column(column_key)

    if len(selected) < 4:
        add_column("saving_value_inr")
        add_column("dept")
        add_column("status")

    # Ensure we have at most 6 columns for readability
    return selected[:6]


def extract_topic_keywords(query: str) -> list[str]:
    keywords = []
    topic_map = {
        "suspension": ["suspension", "damper", "spring", "strut", "subframe", "control arm", "stabilizer"],
        "chassis": ["chassis", "frame", "cross member"],
        "brake": ["brake", "caliper", "rotor", "booster"],
        "engine": ["engine", "powertrain"],
        "interior": ["trim", "console", "seat", "headliner", "panel"],
        "lighting": ["lamp", "lighting", "led", "headlamp"],
        "paint": ["paint", "coating"],
    }
    lower = query.lower()
    for topic, synonyms in topic_map.items():
        if topic in lower or any(word in lower for word in synonyms):
            keywords.extend(synonyms)
    if not keywords:
        tokens = re.findall(r"[a-zA-Z]+", lower)
        keywords = [token for token in tokens if len(token) > 4][:5]
    return list(dict.fromkeys(keywords))  # preserve order


def filter_rows_for_auto_analysis(table_rows: list[dict], user_query: str) -> list[dict]:
    if not table_rows:
        return []
    filters = parse_query_filters(user_query)
    keywords = extract_topic_keywords(user_query)
    query_lower = user_query.lower()
    filtered = []
    
    # Check for status filter in query (TBD, OK, NG, Hold)
    status_filter = None
    if "tbd" in query_lower or "'tbd'" in query_lower or '"tbd"' in query_lower:
        status_filter = "TBD"
    elif "status ok" in query_lower or "approved" in query_lower:
        status_filter = "OK"
    elif "status ng" in query_lower or "rejected" in query_lower:
        status_filter = "NG"
    elif "hold" in query_lower:
        status_filter = "Hold"
    
    # Also check parse_query_filters result
    if 'status' in filters:
        status_filter = filters['status']
    
    for row in table_rows:
        # status filter
        if status_filter:
            row_status = str(_get_value(row, "status", "")).upper().strip()
            if row_status != status_filter.upper():
                continue
        
        # numeric filter: saving value
        if 'saving value' in filters:
            op, val = filters['saving value']
            cell = _get_value(row, "saving_value_inr", 0)
            saving = safe_float_convert(cell)
            if op == 'greater' and saving < val:
                continue
            if op == 'less' and saving > val:
                continue
        
        # dept filter
        if 'dept' in filters:
            dept_filter = filters['dept'].lower()
            row_dept = str(_get_value(row, "dept", "")).lower()
            if dept_filter not in row_dept and row_dept not in dept_filter:
                # Also check if query mentions specific depts
                if "finance" in query_lower and "finance" not in row_dept:
                    continue
                if "vave" in query_lower and "vave" not in row_dept:
                    continue
        
        # keyword filter (only if no specific filters applied)
        if keywords and not status_filter and 'saving value' not in filters:
            search_text = " ".join(str(v).lower() for v in row.values() if isinstance(v, str))
            if not any(keyword in search_text for keyword in keywords):
                continue
        elif keywords:
            # Even with filters, check if keywords match
            search_text = " ".join(str(v).lower() for v in row.values() if isinstance(v, str))
            # More lenient: at least one keyword should match
            if not any(keyword in search_text for keyword in keywords):
                continue
        
        filtered.append(row)
    
    # If we have specific filters but no results, return empty (don't fall back to all rows)
    if (status_filter or 'saving value' in filters or 'dept' in filters) and not filtered:
        return []
    
    if not filtered:
        return table_rows
    return filtered


def build_markdown_table(rows: list[dict], columns: list[str]) -> str:
    if not rows:
        return "No data available."
    header = "| " + " | ".join(columns) + " |"
    separator = "| " + " | ".join(["---"] * len(columns)) + " |"
    lines = [header, separator]
    for row in rows:
        values = []
        for col in columns:
            value = _get_value(row, col, "-")
            if isinstance(value, float):
                value = f"{value:.2f}"
            values.append(str(value) if value not in [None, "None"] else "-")
        lines.append("| " + " | ".join(values) + " |")
    return "\n".join(lines)


def infer_trend_from_rows(rows: list[dict]) -> str:
    if not rows:
        return "material optimization"
    combined_text = " ".join(str(_get_value(row, "cost_reduction_idea", "")) for row in rows).lower()
    if "replace" in combined_text or "material" in combined_text:
        return "material substitution"
    if "reduce" in combined_text or "thickness" in combined_text:
        return "dimensional optimization"
    if "remove" in combined_text or "delete" in combined_text:
        return "part deletion"
    if "integration" in combined_text or "integrate" in combined_text:
        return "component integration"
    return "value engineering"


def generate_ai_ideas(user_query: str, base_rows: list[dict], columns: list[str], count: int = 5) -> list[dict]:
    generated = []
    if count <= 0:
        return generated
    filters = parse_query_filters(user_query)
    trend = infer_trend_from_rows(base_rows)
    keywords = extract_topic_keywords(user_query)
    dominant_dept = None
    for row in base_rows:
        dept_val = _get_value(row, "dept")
        if dept_val and dept_val != "-":
            dominant_dept = dept_val
            break
    dominant_dept = dominant_dept or "PE"

    savings = [
        safe_float_convert(_get_value(row, "saving_value_inr", 0))
        for row in base_rows if safe_float_convert(_get_value(row, "saving_value_inr", 0)) > 0
    ]
    avg_saving = sum(savings) / len(savings) if savings else 40.0

    # Extract actual idea patterns from base_rows for more realistic generation
    idea_patterns = []
    for row in base_rows[:3]:  # Use top 3 as templates
        idea_text = str(_get_value(row, "cost_reduction_idea", "")).lower()
        if idea_text and len(idea_text) > 20:
            idea_patterns.append(idea_text)

    # Generate more realistic ideas based on patterns
    suspension_components = ["stabilizer link", "damper dust cover", "lower control arm", "subframe gusset", "shock absorber mount"]
    brake_components = ["brake hose", "brake pad backing plate", "brake caliper bracket", "parking brake cable", "brake fluid reservoir"]
    interior_components = ["door map pocket", "headliner overlap", "glove box bin", "console trim", "seat rail cover"]
    
    # Select component list based on query keywords
    if any(kw in user_query.lower() for kw in ["suspension", "chassis", "damper", "spring", "strut"]):
        component_list = suspension_components
    elif any(kw in user_query.lower() for kw in ["brake", "caliper", "rotor"]):
        component_list = brake_components
    elif any(kw in user_query.lower() for kw in ["interior", "trim", "console", "door", "seat"]):
        component_list = interior_components
    else:
        component_list = keywords if keywords else COMPONENT_FALLBACKS

    # Extract actual material/component patterns from existing ideas for more realistic generation
    material_patterns = []
    component_patterns = []
    for row in base_rows[:5]:  # Analyze top 5 ideas
        idea_text = str(_get_value(row, "cost_reduction_idea", "")).lower()
        if "replace" in idea_text:
            # Extract material patterns: "Replace X with Y"
            parts = idea_text.split("replace")[1].split("with") if "with" in idea_text else []
            if len(parts) >= 2:
                old_mat = parts[0].strip().split()[0] if parts[0].strip() else ""
                new_mat = parts[1].strip().split()[0] if parts[1].strip() else ""
                if old_mat and new_mat and len(old_mat) > 2:
                    material_patterns.append((old_mat, new_mat))
        # Extract component names
        for comp in component_list:
            if comp in idea_text:
                component_patterns.append(comp)
    
    # Use actual patterns if available
    if not material_patterns:
        material_patterns = [("PC-ABS", "PP"), ("Steel", "Plastic"), ("Brass", "Steel")]
    if not component_patterns:
        component_patterns = component_list if component_list else COMPONENT_FALLBACKS
    
    for idx in range(count):
        idea = {}
        component = component_patterns[idx % len(component_patterns)] if component_patterns else COMPONENT_FALLBACKS[idx % len(COMPONENT_FALLBACKS)]
        
        # Tune savings to match queried constraints when present
        est_saving = max(20.0, avg_saving * (0.8 + 0.15 * idx))
        if 'saving value' in filters:
            op, val = filters['saving value']
            if op == 'less':
                est_saving = max(5.0, min(val * 0.9, val - 1))
            elif op == 'greater':
                est_saving = max(val + 5, val * 1.05)

        # Tune weight saving to align with query constraints
        base_weight = 0.15 + 0.05 * idx
        if 'weight_saving' in filters:
            w_op, w_val = filters['weight_saving']
            if w_op == 'greater':
                base_weight = max(w_val + 0.05, w_val * 1.1)
            elif w_op == 'less':
                base_weight = max(0.05, min(w_val - 0.05, w_val * 0.9))
            elif w_op == 'equal':
                base_weight = w_val
        
        # Generate realistic idea description based on trend and actual patterns
        if trend == "material substitution" and material_patterns:
            old_mat, new_mat = material_patterns[idx % len(material_patterns)]
            idea_desc = f"Replace {old_mat} with {new_mat} in the manufacture of {component} within MGI"
        elif trend == "material substitution":
            idea_desc = f"Replace higher-cost material with lower-cost alternative in {component} for cost reduction"
        elif trend == "dimensional optimization":
            idea_desc = f"Reduce the thickness/coating thickness of {component} to optimize material usage and cost"
        elif trend == "part deletion":
            idea_desc = f"Remove or de-content {component} where functionality can be maintained without it"
        elif trend == "component integration":
            idea_desc = f"Integrate {component} into adjacent assembly to reduce parts count and assembly cost"
        else:
            # Use pattern from existing ideas if available
            if idea_patterns and idx < len(idea_patterns):
                base_pattern = idea_patterns[idx]
                # Adapt the pattern to new component
                idea_desc = base_pattern.replace(component_patterns[0] if component_patterns else "component", component)
            else:
                idea_desc = f"Apply {trend} approach to {component} for cost reduction in MGI"
        
        # Find the correct column names
        idea_id_col = _find_column(columns, ["idea_id", "idea id", "ideaid"])
        cost_idea_col = _find_column(columns, ["cost_reduction_idea", "cost reduction idea", "proposal", "summary"])
        saving_col = _find_column(columns, ["saving_value_inr", "saving value (inr)", "saving value", "mgi_gross_saving"])
        dept_col = _find_column(columns, ["dept", "department"])
        status_col = _find_column(columns, ["status"])
        weight_col = _find_column(columns, ["weight_saving", "weight saving (kg)", "weight saving"])
        way_forward_col = _find_column(columns, ["way_forward", "way forward", "imp"])
        
        # Set values using found column names
        for col in columns:
            norm = _normalize_column(col)
            if idea_id_col and col == idea_id_col:
                idea[col] = f"AI-GEN-{idx+1:02d}"
            elif cost_idea_col and col == cost_idea_col:
                idea[col] = idea_desc
            elif saving_col and col == saving_col:
                idea[col] = f"₹ {est_saving:,.0f}" if "inr" in norm else est_saving
            elif dept_col and col == dept_col:
                idea[col] = dominant_dept
            elif status_col and col == status_col:
                idea[col] = "AI-PROPOSAL"
            elif weight_col and col == weight_col:
                idea[col] = round(base_weight, 2)
            elif way_forward_col and col == way_forward_col:
                idea[col] = "Review feasibility with cross-functional team"
            elif "respons" in norm or "owner" in norm:
                idea[col] = "Digital Analyst"
            else:
                # Default: try to match by normalization
                if "idea" in norm and "id" in norm:
                    idea[col] = f"AI-GEN-{idx+1:02d}"
                elif ("costreductionidea" in norm or "proposal" in norm or "summary" in norm) and not cost_idea_col:
                    idea[col] = idea_desc
                elif ("saving" in norm and "value" in norm) and not saving_col:
                    idea[col] = f"₹ {est_saving:,.0f}" if "inr" in norm else est_saving
                elif "dept" in norm and not dept_col:
                    idea[col] = dominant_dept
                elif "status" in norm and not status_col:
                    idea[col] = "AI-PROPOSAL"
                elif "weight" in norm and not weight_col:
                    idea[col] = round(0.05 + 0.02 * idx, 2)
                elif "wayforward" in norm or "imp" in norm:
                    idea[col] = "Review feasibility with cross-functional team"
                else:
                    idea[col] = "-"
        generated.append(idea)
    return generated


def build_strategy_summary(user_query: str, selected_columns: list[str], existing_count: int, filtered_count: int) -> str:
    return (
        f"Strategy Summary: Interpreted the query \"{user_query}\" and selected "
        f"{len(selected_columns)} priority column(s): {', '.join(selected_columns)}. "
        f"Filtered {filtered_count} retrieved ideas down to {existing_count} highly relevant entries."
    )


def build_auto_analyst_response(user_query: str, table_rows: list[dict], max_existing: int = 10) -> tuple[str, list[dict], list[dict]]:
    # Get available columns from table_rows or use defaults
    if table_rows:
        available_columns = sorted({col for row in table_rows for col in row.keys()})
    else:
        # Use default columns if no data available
        available_columns = DEFAULT_ANALYST_COLUMNS.copy()
        # Try to get schema from database
        try:
            schema_cols = get_idea_schema_columns()
            if schema_cols:
                available_columns = schema_cols[:10]  # Limit to first 10 for readability
        except:
            pass
    
    selected_columns = auto_select_columns(user_query, available_columns)
    
    # Filter rows if we have data
    if table_rows:
        filtered_rows = filter_rows_for_auto_analysis(table_rows, user_query)
        existing_rows = filtered_rows[:max_existing]
    else:
        # If no data, try to fetch from database directly using SQL
        logger.info("No table_rows provided, attempting direct SQL query for status/dept filters")
        filtered_rows = []
        existing_rows = []
        
        query_lower = user_query.lower()
        if "tbd" in query_lower or "status" in query_lower:
            try:
                conn = get_db_connection()
                cursor = conn.cursor(cursor_factory=DictCursor)
                
                status_val = "TBD"
                if "'tbd'" in query_lower or '"tbd"' in query_lower:
                    status_val = "TBD"
                elif "ok" in query_lower and "status" in query_lower:
                    status_val = "OK"
                elif "ng" in query_lower and "status" in query_lower:
                    status_val = "NG"
                
                sql_params = [status_val]
                dept_conditions = []
                
                if "finance" in query_lower:
                    dept_conditions.append("(dept ILIKE %s OR dept ILIKE %s)")
                    sql_params.extend(['%Finance%', '%VAVE%'])
                elif "vave" in query_lower:
                    dept_conditions.append("dept ILIKE %s")
                    sql_params.append('%VAVE%')
                
                where_clause = "status = %s"
                if dept_conditions:
                    where_clause += " AND (" + " OR ".join(dept_conditions) + ")"
                
                sql = f"""
                    SELECT idea_id, cost_reduction_idea, saving_value_inr, dept, status,
                           weight_saving, way_forward, resp, proposal_image_filename
                    FROM ideas
                    WHERE {where_clause}
                    ORDER BY saving_value_inr DESC NULLS LAST
                    LIMIT 20
                """
                cursor.execute(sql, sql_params)
                db_rows = cursor.fetchall()
                filtered_rows = [dict(row) for row in db_rows]
                existing_rows = filtered_rows[:max_existing]
                
                # Convert to table_data format
                for row in existing_rows:
                    row["Idea Id"] = str(row.get('idea_id', 'N/A'))
                    row["Cost Reduction Idea"] = str(row.get('cost_reduction_idea', 'N/A'))
                    row["Saving Value (INR)"] = f"INR {safe_float_convert(row.get('saving_value_inr', 0)):,}"
                    row["CAPEX"] = f"{safe_float_convert(row.get('capex') or row.get('investment_cr')):,.2f}" if (row.get('capex') is not None or row.get('investment_cr') is not None) else "N/A"
                    row["Dept"] = str(row.get('dept', 'N/A'))
                    row["Status"] = str(row.get('status', 'N/A'))
                
                cursor.close()
                conn.close()
                logger.info(f"Direct SQL query returned {len(existing_rows)} rows")
            except Exception as sql_err:
                logger.error(f"Direct SQL query failed: {sql_err}")
                filtered_rows = []
                existing_rows = []
    
    # Always generate AI ideas, even if no existing rows
    ai_generated_rows = generate_ai_ideas(user_query, filtered_rows if filtered_rows else [], selected_columns, count=5)

    table1 = build_markdown_table(existing_rows, selected_columns)
    table2 = build_markdown_table(ai_generated_rows, selected_columns)
    
    if existing_rows:
        summary = build_strategy_summary(user_query, selected_columns, len(existing_rows), len(filtered_rows))
    else:
        summary = (
            f"Strategy Summary: Query \"{user_query}\" requested analysis of specific status/dept ideas. "
            f"Selected {len(selected_columns)} priority column(s): {', '.join(selected_columns)}. "
            f"No matching existing ideas found in database, but generated 5 new AI proposals based on query context."
        )

    response_text = "\n\n".join([
        summary,
        "",
        "Table 1: Existing Ideas",
        table1,
        "",
        "Table 2: AI-Generated Ideas",
        table2
    ])
    
    # Format existing_rows for frontend (same format as retrieve_context)
    formatted_existing_rows = []
    for row in existing_rows:
        # Check if already formatted (has "Idea Id" key)
        if "Idea Id" in row:
            formatted_existing_rows.append(row)
        else:
            # Format from raw database row
            image_filename = row.get('proposal_image_filename')
            image_path = None
            if image_filename:
                src_path = IMAGE_DIRS["proposal"] / image_filename
                static_path = STATIC_IMAGE_DIR / "proposal" / image_filename
                if src_path.exists() and not static_path.exists():
                    try:
                        shutil.copy(src_path, static_path)
                    except Exception as e:
                        logger.error(f"Failed to copy image {image_filename}: {str(e)}")
                if static_path.exists() or src_path.exists():
                    image_path = url_for('static', filename=f'images/proposal/{image_filename}')
            
            # Get current scenario image if available
            current_scenario_image = row.get('current_scenario_image') or row.get('Current Scenario Image')
            current_image_url = "N/A"
            if current_scenario_image and current_scenario_image != "N/A":
                if os.path.exists(current_scenario_image):
                    filename = os.path.basename(current_scenario_image)
                    static_path = STATIC_DIR / "generated" / filename
                    if not static_path.exists():
                        try:
                            shutil.copy(current_scenario_image, static_path)
                        except Exception:
                            pass
                    current_image_url = url_for('static', filename=f'generated/{filename}')
                elif current_scenario_image.startswith("static/") or current_scenario_image.startswith("/static/"):
                    current_image_url = url_for('static', filename=current_scenario_image.replace("static/", ""))
                else:
                    current_image_url = current_scenario_image if current_scenario_image.startswith("http") or current_scenario_image.startswith("/") else "N/A"
            
            formatted_row = {
                "Idea Id": str(row.get('idea_id', 'N/A')),
                "Cost Reduction Idea": str(row.get('cost_reduction_idea', 'N/A')),
                "Proposal Image Path": image_path or "N/A",
                "Current Scenario Image": current_image_url,
                "Proposal Scenario Image": image_path or "N/A",
                "MGI Gross Saving": f"INR {safe_float_convert(row.get('mgi_gross_saving', 0)):,}",
                "Way Forward": str(row.get('way_forward', 'N/A')),
                "Estimated Cost Savings": str(row.get('estimated_cost_savings', 'N/A')),
                "MGI Carline": str(row.get('mgi_carline', 'N/A')),
                "Saving Value (INR)": f"INR {safe_float_convert(row.get('saving_value_inr', 0)):,}",
                "Weight Saving (Kg)": f"{safe_float_convert(row.get('weight_saving', 0)):.2f}",
                "CAPEX": f"{safe_float_convert(row.get('capex') or row.get('investment_cr')):,.2f}" if (row.get('capex') is not None or row.get('investment_cr') is not None) else "N/A",
                "Status": str(row.get('status', 'N/A')),
                "Dept": str(row.get('dept', 'N/A')),
                "KD/LC": str(row.get('kd_lc', 'N/A')),
                # Additional fields for PPT
                "Idea ID": str(row.get('idea_id', 'N/A')),
                "Proposal": str(row.get('cost_reduction_idea', 'N/A')),
                "Imp": str(row.get('way_forward', 'N/A')),
                "Saving Value": f"INR {safe_float_convert(row.get('saving_value_inr', 0)):,}",
                "Responsibility": str(row.get('resp', 'N/A')),
                "Date": str(row.get('target_date', 'N/A'))
            }
            formatted_existing_rows.append(formatted_row)
    
    # Format AI-generated rows similarly
    formatted_ai_rows = []
    for row in ai_generated_rows:
        # Get image paths from row (may come from VLM engine enrichment)
        current_image = row.get("current_scenario_image") or row.get("Current Scenario Image")
        proposal_image = row.get("proposal_scenario_image") or row.get("Proposal Scenario Image")
        
        # Format image paths as URLs if they exist
        current_image_url = "N/A"
        proposal_image_url = "N/A"
        
        if current_image and current_image != "N/A":
            # If it's already a URL path, use it; otherwise convert to URL
            if current_image.startswith("static/") or current_image.startswith("/static/"):
                current_image_url = url_for('static', filename=current_image.replace("static/", ""))
            elif os.path.exists(current_image):
                # Copy to static if needed
                filename = os.path.basename(current_image)
                static_path = STATIC_DIR / "generated" / filename
                if not static_path.exists():
                    try:
                        shutil.copy(current_image, static_path)
                    except Exception:
                        pass
                current_image_url = url_for('static', filename=f'generated/{filename}')
            else:
                current_image_url = current_image if current_image.startswith("http") or current_image.startswith("/") else "N/A"
        
        if proposal_image and proposal_image != "N/A":
            # If it's already a URL path, use it; otherwise convert to URL
            if proposal_image.startswith("static/") or proposal_image.startswith("/static/"):
                proposal_image_url = url_for('static', filename=proposal_image.replace("static/", ""))
            elif os.path.exists(proposal_image):
                # Copy to static if needed
                filename = os.path.basename(proposal_image)
                static_path = STATIC_DIR / "generated" / filename
                if not static_path.exists():
                    try:
                        shutil.copy(proposal_image, static_path)
                    except Exception:
                        pass
                proposal_image_url = url_for('static', filename=f'generated/{filename}')
            else:
                proposal_image_url = proposal_image if proposal_image.startswith("http") or proposal_image.startswith("/") else "N/A"
        
        formatted_ai_row = {
            "Idea Id": row.get("idea_id", "AI-GEN-01"),
            "Cost Reduction Idea": row.get("cost_reduction_idea", row.get("Cost Reduction Idea", "N/A")),
            "Source": "AI Generated",
            "Proposal Image Path": proposal_image_url if proposal_image_url != "N/A" else "N/A",
            "Current Scenario Image": current_image_url,
            "Proposal Scenario Image": proposal_image_url,
            "MGI Gross Saving": "N/A",
            "Way Forward": row.get("way_forward", row.get("Way Forward", "Review feasibility")),
            "Estimated Cost Savings": "N/A",
            "MGI Carline": "N/A",
            "Saving Value (INR)": row.get("saving_value_inr", row.get("Saving Value (INR)", "INR 0")),
            "Weight Saving (Kg)": str(row.get("weight_saving", row.get("Weight Saving (Kg)", "0.00"))),
            "Status": row.get("status", "AI-PROPOSAL"),
            "Dept": row.get("dept", row.get("Dept", "PE")),
            "KD/LC": "N/A",
            "Idea ID": row.get("idea_id", "AI-GEN-01"),
            "Proposal": row.get("cost_reduction_idea", row.get("Cost Reduction Idea", "N/A")),
            "Imp": row.get("way_forward", row.get("Way Forward", "Review feasibility")),
            "Saving Value": row.get("saving_value_inr", row.get("Saving Value (INR)", "INR 0")),
            "Responsibility": "Digital Analyst",
            "Date": "N/A"
        }
        formatted_ai_rows.append(formatted_ai_row)
    
    # Combine existing + AI-generated for table display
    combined_table_data = formatted_existing_rows + formatted_ai_rows
    
    return response_text, formatted_existing_rows, formatted_ai_rows

# ========================= END AUTO ANALYST HELPERS ======================
def filter_ideas(rows, filters):
    """Filter ideas based on parsed filters."""
    filtered = []

    for row in rows:
        # Numeric filters with safe conversion
        if 'saving value' in filters:
            op, val = filters['saving value']
            saving = safe_float_convert(row.get('saving_value_inr'))  # FIXED: Using safe conversion
            if op == 'greater' and saving <= val:
                continue
            if op == 'less' and saving >= val:
                continue

        if 'mgi gross saving' in filters:
            op, val = filters['mgi gross saving']
            saving = safe_float_convert(row.get('mgi_gross_saving'))  # FIXED: Using safe conversion
            if op == 'greater' and saving <= val:
                continue
            if op == 'less' and saving >= val:
                continue

        if 'estimated cost savings' in filters:
            op, val = filters['estimated cost savings']
            saving = safe_float_convert(row.get('estimated_cost_savings'))  # FIXED: Using safe conversion
            if op == 'greater' and saving <= val:
                continue
            if op == 'less' and saving >= val:
                continue

        if 'investment' in filters:
            op, val = filters['investment']
            investment = safe_float_convert(row.get('investment_cr'))  # FIXED: Using safe conversion
            if op == 'greater' and investment <= val:
                continue
            if op == 'less' and investment >= val:
                continue

        if 'wtd. avg.' in filters:
            op, val = filters['wtd. avg.']
            wtd_avg = safe_float_convert(row.get('wtd_avg'))  # FIXED: Using safe conversion
            if op == 'greater' and wtd_avg <= val:
                continue
            if op == 'less' and wtd_avg >= val:
                continue

        if 'new tool cost' in filters:
            op, val = filters['new tool cost']
            cost = safe_float_convert(row.get('new_tool_cost'))  # FIXED: Using safe conversion
            if op == 'greater' and cost <= val:
                continue
            if op == 'less' and cost >= val:
                continue

        if 'tool modification cost' in filters:
            op, val = filters['tool modification cost']
            cost = safe_float_convert(row.get('tool_modification_cost'))  # FIXED: Using safe conversion
            if op == 'greater' and cost <= val:
                continue
            if op == 'less' and cost >= val:
                continue

        if 'weight_saving' in filters:
            op, val = filters['weight_saving']
            weight = safe_float_convert(row.get('weight_saving'))  # FIXED: Using safe conversion
            if op == 'greater' and weight <= val:
                continue
            if op == 'less' and weight >= val:
                continue
            if op == 'equal' and weight != val:
                continue

        # Text filters
        if 'status' in filters:
            status = filters['status']
            if str(row.get('status', '')).upper() != status:
                continue

        if 'dept' in filters:
            dept = filters['dept']
            if str(row.get('dept', '')).lower() != dept:
                continue

        if 'kd_lc' in filters:
            kd_lc = filters['kd_lc']
            if str(row.get('kd_lc', '')).upper() != kd_lc:
                continue

        if 'system' in filters:
            system = filters['system']
            proposal = str(row.get('cost_reduction_idea', '')).lower()
            if system not in proposal:
                continue

        filtered.append(row)

    return filtered

# Max chars for embedding (avoids model limits and improves relevance for long prompts)
MAX_QUERY_LENGTH_FOR_EMBEDDING = 2000

def _query_for_embedding(query: str) -> str:
    """Normalize and optionally truncate query for embedding so complex/long prompts don't fail."""
    if not query or not query.strip():
        return ""
    q = query.strip()
    if len(q) <= MAX_QUERY_LENGTH_FOR_EMBEDDING:
        return q
    # Preserve start (intent) and end (target) for long prompts; take middle summary
    head = q[:800].strip()
    tail = q[-600:].strip() if len(q) > 1400 else ""
    if tail and head != tail:
        return head + " ... " + tail
    return head


def retrieve_context(query, top_k=10):
    """Retrieve relevant ideas from FAISS index or SQL (for status/dept filters). Uses full query for filters; safe-length query for embedding."""
    if embedding_model is None or faiss_index is None:
        logger.warning("Vector DB not initialized")
        return [], []

    try:
        filters = parse_query_filters(query)
        query_lower = (query or "").lower()
        search_query = _query_for_embedding(query or "")
        
        # Check if query has status or dept filters - use SQL for better accuracy
        has_status_filter = 'status' in filters or any(s in query_lower for s in ["'tbd'", '"tbd"', 'tbd status', 'status tbd', 'ok status', 'status ok', 'ng status', 'status ng'])
        has_dept_filter = 'dept' in filters or any(d in query_lower for d in ["finance", "vave", "pe", "purchase"])
        
        # If status or dept filter detected, use SQL query directly
        if has_status_filter or has_dept_filter:
            logger.info(f"Status/Dept filter detected, using SQL query instead of vector search")
            conn = None
            cursor = None
            try:
                conn = get_db_connection()
                cursor = conn.cursor(cursor_factory=DictCursor)
                
                # Build SQL query
                sql_conditions = []
                sql_params = []
                
                if has_status_filter:
                    status_val = filters.get('status', 'TBD')
                    if 'tbd' in query_lower:
                        status_val = 'TBD'
                    elif 'ok' in query_lower and 'status' in query_lower:
                        status_val = 'OK'
                    elif 'ng' in query_lower and 'status' in query_lower:
                        status_val = 'NG'
                    sql_conditions.append("status = %s")
                    sql_params.append(status_val)
                
                if has_dept_filter:
                    dept_val = filters.get('dept', '')
                    if 'finance' in query_lower:
                        sql_conditions.append("(dept ILIKE %s OR dept ILIKE %s)")
                        sql_params.extend(['%Finance%', '%VAVE%'])
                    elif 'vave' in query_lower:
                        sql_conditions.append("dept ILIKE %s")
                        sql_params.append('%VAVE%')
                    elif dept_val:
                        sql_conditions.append("dept ILIKE %s")
                        sql_params.append(f'%{dept_val}%')
                
                # Add saving value filter if present
                if 'saving value' in filters:
                    op, val = filters['saving value']
                    if op == 'greater':
                        sql_conditions.append("saving_value_inr > %s")
                        sql_params.append(val)
                    elif op == 'less':
                        sql_conditions.append("saving_value_inr < %s")
                        sql_params.append(val)
                
                where_clause = " AND ".join(sql_conditions) if sql_conditions else "1=1"
                sql_query = f"""
                    SELECT idea_id, cost_reduction_idea, reason, mgi_gross_saving,
                           estimated_cost_savings, saving_value_inr, weight_saving, group_id, status,
                           way_forward, dept, target_date, kd_lc, wtd_avg, est_impl_date,
                           investment_cr, mgi_carline, benchmarking_carline, mg_product_scenario,
                           competitor_product_scenario, purpose_mg_product, purpose_competitor_product,
                           impact_other_systems, client_statement, cae_required, homologation_required,
                           styling_change, part_level_testing, assembly_trials, cad_drawing_update,
                           ecn_required, part_production_trials, vehicle_level_testing, idea_generated_by,
                           new_tool_required, new_tool_cost, tool_modification_required, tool_modification_cost,
                           variants, current_status, resp, mix, volume, purchase_proposal, interest,
                           payback_months, mgi_pe_feasibility, homeroom_approval, pp_approval,
                           supplier_feasibility, financial_feasibility, proposal_image_filename,
                           mg_vehicle_image, competitor_vehicle_image
                    FROM ideas
                    WHERE {where_clause}
                    ORDER BY saving_value_inr DESC NULLS LAST
                    LIMIT 50
                """
                
                cursor.execute(sql_query, sql_params)
                db_rows = cursor.fetchall()
                selected_rows = [dict(row) for row in db_rows]
                logger.info(f"SQL query returned {len(selected_rows)} rows")
                
            except Exception as sql_err:
                logger.error(f"SQL fallback error: {sql_err}, falling back to vector search")
                selected_rows = []
            finally:
                if cursor:
                    cursor.close()
                if conn:
                    conn.close()
            
            # If SQL returned results, use them; otherwise fall back to vector search
            if not selected_rows:
                logger.info("SQL returned no results, falling back to vector search with increased top_k")
                q_emb = search_query or query
                try:
                    query_emb = embedding_model.encode([q_emb], convert_to_numpy=True, device=device)
                    D, I = faiss_index.search(query_emb, min(50, len(idea_rows)))
                    selected_rows = [idea_rows[i] for i in I[0] if i < len(idea_rows)]
                except Exception as enc_err:
                    logger.warning(f"Embedding failed for long query, retrying with truncated: {enc_err}")
                    q_emb = (q_emb[:1500] + " ...") if len(q_emb) > 1500 else q_emb
                    query_emb = embedding_model.encode([q_emb], convert_to_numpy=True, device=device)
                    D, I = faiss_index.search(query_emb, min(50, len(idea_rows)))
                    selected_rows = [idea_rows[i] for i in I[0] if i < len(idea_rows)]
        else:
            # Normal vector search path (use safe-length query for embedding)
            search_k = top_k
            if any(k in filters for k in ['weight_saving', 'saving value', 'mgi gross saving', 'estimated cost savings']):
                search_k = max(top_k, 50)
            q_emb = search_query or query
            try:
                query_emb = embedding_model.encode([q_emb], convert_to_numpy=True, device=device)
                D, I = faiss_index.search(query_emb, min(search_k, len(idea_rows)))
                selected_rows = [idea_rows[i] for i in I[0] if i < len(idea_rows)]
            except Exception as enc_err:
                logger.warning(f"Embedding failed for complex query, retrying with truncated: {enc_err}")
                q_emb = (q_emb[:1500] + " ...") if len(q_emb) > 1500 else q_emb
                query_emb = embedding_model.encode([q_emb], convert_to_numpy=True, device=device)
                D, I = faiss_index.search(query_emb, min(search_k, len(idea_rows)))
                selected_rows = [idea_rows[i] for i in I[0] if i < len(idea_rows)]
        
        filtered_rows = filter_ideas(selected_rows, filters)

        table_data = []
        contexts = []

        for row in filtered_rows:
            image_filename = row.get('proposal_image_filename')
            image_path = image_filename if image_filename else None

            if image_path:
                src_path = IMAGE_DIRS["proposal"] / image_path
                static_path = STATIC_IMAGE_DIR / "proposal" / image_path
                if src_path.exists() and not static_path.exists():
                    try:
                        shutil.copy(src_path, static_path)
                        logger.info(f"Copied image {image_path} to {static_path}")
                    except Exception as e:
                        logger.error(f"Failed to copy image {image_path}: {str(e)}")
                        image_path = None

                image_path = url_for('static', filename=f'images/proposal/{image_path}') if image_path else "N/A"

            # Get current scenario image if available from database
            current_scenario_image = row.get('current_scenario_image') or row.get('Current Scenario Image')
            current_image_url = "N/A"
            if current_scenario_image and current_scenario_image != "N/A":
                if os.path.exists(current_scenario_image):
                    filename = os.path.basename(current_scenario_image)
                    static_path = STATIC_DIR / "generated" / filename
                    if not static_path.exists():
                        try:
                            shutil.copy(current_scenario_image, static_path)
                        except Exception:
                            pass
                    current_image_url = url_for('static', filename=f'generated/{filename}')
                elif current_scenario_image.startswith("static/") or current_scenario_image.startswith("/static/"):
                    current_image_url = url_for('static', filename=current_scenario_image.replace("static/", ""))
                else:
                    current_image_url = current_scenario_image if current_scenario_image.startswith("http") or current_scenario_image.startswith("/") else "N/A"
            
            idea_data = {
                "Idea Id": str(row.get('idea_id', 'N/A')),
                "Cost Reduction Idea": str(row.get('cost_reduction_idea', 'N/A')),
                "Proposal Image Path": image_path,
                "Current Scenario Image": current_image_url,
                "Proposal Scenario Image": image_path if image_path != "N/A" else "N/A",
                "MGI Gross Saving": f"INR {safe_float_convert(row.get('mgi_gross_saving')):,}",  # FIXED: Using safe conversion
                "Way Forward": str(row.get('way_forward', 'N/A')),
                "Estimated Cost Savings": str(row.get('estimated_cost_savings', 'N/A')),
                "MGI Carline": str(row.get('mgi_carline', 'N/A')),
                "Saving Value (INR)": f"INR {safe_float_convert(row.get('saving_value_inr')):,}",  # FIXED: Using safe conversion
                "Weight Saving (Kg)": f"{safe_float_convert(row.get('weight_saving')):.2f}",  # FIXED: Using safe conversion
                "Status": str(row.get('status', 'N/A')),
                "Dept": str(row.get('dept', 'N/A')),
                "KD/LC": str(row.get('kd_lc', 'N/A')),
                # Additional fields for PPT
                "Idea ID": str(row.get('idea_id', 'N/A')),
                "Proposal": str(row.get('cost_reduction_idea', 'N/A')),
                "Imp": str(row.get('way_forward', 'N/A')),
                "Saving Value": f"INR {safe_float_convert(row.get('saving_value_inr')):,}",
                "Responsibility": str(row.get('resp', 'N/A')),
                "Date": str(row.get('target_date', 'N/A')),
                # Preserve raw database fields for agent/VLM processing
                "mg_vehicle_image": row.get('mg_vehicle_image'),
                "competitor_vehicle_image": row.get('competitor_vehicle_image'),
                "proposal_image_filename": row.get('proposal_image_filename'),
                "idea_id": row.get('idea_id'),
                "cost_reduction_idea": row.get('cost_reduction_idea'),
                "way_forward": row.get('way_forward'),
                "status": row.get('status'),
                "dept": row.get('dept'),
                "saving_value_inr": row.get('saving_value_inr'),
                "weight_saving": row.get('weight_saving'),
                "resp": row.get('resp'),
                "target_date": row.get('target_date')
            }

            table_data.append(idea_data)
            
            # Build concise, relevant context string (exclude redundant/nan fields)
            context_parts = []
            if idea_data.get("Cost Reduction Idea") and idea_data["Cost Reduction Idea"] != "N/A":
                context_parts.append(f"Idea: {idea_data['Cost Reduction Idea']}")
            if idea_data.get("Saving Value (INR)") and idea_data["Saving Value (INR)"] not in ["INR 0.0", "INR N/A", "N/A"]:
                context_parts.append(f"Savings: {idea_data['Saving Value (INR)']}")
            if idea_data.get("Weight Saving (Kg)") and idea_data["Weight Saving (Kg)"] not in ["0.00", "N/A"]:
                context_parts.append(f"Weight Saved: {idea_data['Weight Saving (Kg)']} kg")
            if idea_data.get("Status") and idea_data["Status"] not in ["N/A", "nan"]:
                context_parts.append(f"Status: {idea_data['Status']}")
            if idea_data.get("Dept") and idea_data["Dept"] not in ["N/A", "nan"]:
                context_parts.append(f"Department: {idea_data['Dept']}")
            if idea_data.get("Estimated Cost Savings") and idea_data["Estimated Cost Savings"] not in ["N/A", "nan"] and len(str(idea_data["Estimated Cost Savings"])) < 200:
                # Only include if not too verbose
                savings_text = str(idea_data["Estimated Cost Savings"])
                if len(savings_text) < 150:
                    context_parts.append(f"Details: {savings_text[:150]}")
            
            contexts.append(" | ".join(context_parts) if context_parts else f"Idea ID: {idea_data.get('Idea Id', 'N/A')}")

        return table_data, "\n\n".join([f"{i+1}. {ctx}" for i, ctx in enumerate(contexts)])

    except Exception as e:
        logger.error(f"Error in retrieve_context: {str(e)}")
        return [], []

def infer_vlm(image, query, top_k=10):
    """Run VLM inference on an image and query."""
    if not vlm_initialized:
        logger.warning("VLM not initialized")
        return "VLM not initialized", None, [], [], {}

    try:
        image = image.convert("RGB")
        inputs = vlm_processor(images=image, text=query, return_tensors="pt").to(device)

        with torch.no_grad():
            output = vlm_model.generate(**inputs, max_new_tokens=50)

        generated_text = vlm_processor.decode(output[0], skip_special_tokens=True)
        logger.info(f"BLIP generated text: {generated_text}")

        query_embedding = embedding_model.encode(generated_text, convert_to_tensor=True, device=device)
        scores = util.pytorch_cos_sim(query_embedding, proposal_embeddings)[0]

        top_k = min(top_k, len(valid_df))
        top_scores, top_indices = torch.topk(scores, k=top_k)

        result_rows = valid_df.iloc[top_indices.tolist()].copy()
        logger.info(f"Retrieved {len(result_rows)} rows from valid_df")

        def get_proposal_image_path(x):
            if pd.isna(x):
                logger.warning("Proposal image filename is NaN")
                return "N/A"

            path = IMAGE_DIRS["proposal"] / x
            static_path = STATIC_IMAGE_DIR / "proposal" / x

            if path.exists():
                if not static_path.exists():
                    try:
                        shutil.copy(path, static_path)
                        logger.info(f"Copied image {x} to {static_path}")
                    except Exception as e:
                        logger.error(f"Failed to copy image {x}: {str(e)}")
                return url_for('static', filename=f'images/proposal/{x}')

            logger.warning(f"Proposal image missing: {path}")
            return "N/A"

        result_rows['Proposal Image Path'] = result_rows['proposal_image_filename'].apply(get_proposal_image_path)

        table_data = []
        for _, row in result_rows.iterrows():
            idea_data = {
                "Idea Id": str(row.get('idea_id', 'N/A')),
                "Cost Reduction Idea": str(row.get('cost_reduction_idea', 'N/A')),
                "Proposal Image Path": row.get('Proposal Image Path', 'N/A'),
                "MGI Gross Saving": f"INR {safe_float_convert(row.get('mgi_gross_saving')):,}",  # FIXED: Using safe conversion
                "Way Forward": str(row.get('way_forward', 'N/A')),
                "Estimated Cost Savings": str(row.get('estimated_cost_savings', 'N/A')),
                "MGI Carline": str(row.get('mgi_carline', 'N/A')),
                "Saving Value (INR)": f"INR {safe_float_convert(row.get('saving_value_inr')):,}",  # FIXED: Using safe conversion
                "Weight Saving (Kg)": f"{safe_float_convert(row.get('weight_saving')):.2f}",  # FIXED: Using safe conversion
                "Status": str(row.get('status', 'N/A')),
                "Dept": str(row.get('dept', 'N/A')),
                "KD/LC": str(row.get('kd_lc', 'N/A')),
                # Additional fields for PPT
                "Idea ID": str(row.get('idea_id', 'N/A')),
                "Proposal": str(row.get('cost_reduction_idea', 'N/A')),
                "Imp": str(row.get('way_forward', 'N/A')),
                "Saving Value": f"INR {safe_float_convert(row.get('saving_value_inr')):,}",
                "Responsibility": str(row.get('resp', 'N/A')),
                "Date": str(row.get('target_date', 'N/A'))
            }
            table_data.append(idea_data)

        image_urls = {
            'proposal': [url for url in result_rows['Proposal Image Path'].dropna().tolist() if url != "N/A"],
            'mg': [url_for('static', filename=f'images/mg/{x}') for x in result_rows['mg_vehicle_image'].dropna() if (IMAGE_DIRS['mg'] / x).exists()],
            'competitor': [url_for('static', filename=f'images/competitor/{x}') for x in result_rows['competitor_vehicle_image'].dropna() if (IMAGE_DIRS['competitor'] / x).exists()]
        }

        return None, generated_text, table_data, top_scores.tolist(), image_urls

    except Exception as e:
        logger.error(f"Error in infer_vlm: {str(e)}")
        return str(e), None, [], [], {}

def log_query_to_db(username, user_input, response_text, image_path):
    """Logs the chat interaction to the chat_history table in PostgreSQL."""
    conn = None
    cursor = None
    try:
        conn = get_db_connection()
        cursor = conn.cursor()
        
        cursor.execute("""
            INSERT INTO chat_history (username, user_query, response_text, image_path)
            VALUES (%s, %s, %s, %s)
        """, (username, user_input, response_text, image_path))
        
        conn.commit()
        logger.info(f"Logged chat for user: {username}")
    except Exception as e:
        logger.error(f"Failed to log chat to DB: {str(e)}")
        if conn:
            conn.rollback()
    finally:
        if cursor:
            cursor.close()
        if conn:
            conn.close()

def generate_response(username, user_query, image_path=None, pre_retrieved_data=None):
    """Generate response using VLM or LLM based on input."""
    try:
        if image_path:
            # VLM processing
            if not vlm_initialized:
                logger.warning("VLM not initialized, attempting to set up now...")
                setup_vlm() # Attempt to initialize
                if not vlm_initialized:
                    logger.error("VLM initialization failed, falling back to LLM")
                    table_data, response_text = retrieve_context(user_query)
                    return {
                        "response_text": response_text or "No relevant cost reduction ideas found",
                        "table_data": table_data, "image_urls": [],
                        "excel_url": "",
                        "ppt_url": "",
                        "vlm_warning": "VLM unavailable due to initialization failure"
                    }

            image = Image.open(image_path)
            error, generated_text, table_data, similarity_scores, image_urls = infer_vlm(
                image, user_query, top_k=10
            )

            image_path_to_log = os.path.basename(image_path) if image_path else None
            executor.submit(log_query_to_db, 
                                        username,         # <-- ADD THIS
                                        user_query, 
                                        generated_text, 
                                        image_path_to_log)
            if error:
                logger.error(f"VLM processing failed: {error}")
                table_data, response_text = retrieve_context(user_query)
                return {
                    "response_text": response_text or "No relevant cost reduction ideas found",
                    "table_data": table_data, "image_urls": [],
                    "excel_url": "",
                    "ppt_url": "",
                    "vlm_warning": f"VLM processing failed: {error}"
                }

            # --- Store VLM result in database (asynchronously) ---
            idea_id_to_log = table_data[0].get('Idea Id') if table_data else 'N/A'
            image_path_to_log = os.path.basename(image_path) if image_path else None
            executor.submit(log_query_to_db, 
                            username,         # <-- ADD THIS
                            user_query, 
                            response, 
                            None)
            
            return {
                "response_text": generated_text,
                "table_data": table_data,
                "image_urls": image_urls['proposal'],
                "excel_url": "",
                "ppt_url": "",
                "vlm_warning": None
            }

        else:
            # LLM processing
            # Use pre-retrieved data if provided (from agentic RAG), otherwise retrieve fresh
            if pre_retrieved_data is not None:
                table_data, context_str = pre_retrieved_data
                logger.info(f"Using pre-retrieved data: {len(table_data)} ideas")
            else:
                table_data, context_str = retrieve_context(user_query)

            # Always try Auto-Analyst, even if table_data is empty (it will do SQL fallback)
            ai_generated_rows = []
            existing_rows_formatted = []
            auto_response = None
            try:
                auto_response, existing_rows_formatted, ai_generated_rows = build_auto_analyst_response(user_query, table_data)
                logger.info("Auto-Analyst response generated successfully.")
                # Use the formatted rows from Auto-Analyst instead of original table_data
                if existing_rows_formatted:
                    table_data = existing_rows_formatted
                    logger.info(f"Using {len(table_data)} existing rows from Auto-Analyst")
            except Exception as auto_exc:
                logger.error(f"Auto-Analyst formatting failed: {auto_exc}")
                auto_response = None
            
            # If Auto-Analyst returned empty and we have no table_data, still try to get some data
            if not table_data and not auto_response:
                logger.warning(f"No relevant ideas found for query: {user_query}")
                return {
                    "response_text": f"No relevant cost reduction ideas found for query: {user_query}",
                    "table_data": [], "image_urls": [], 
                    "excel_url": "",
                    "ppt_url": ""
                }

            if auto_response:
                response = auto_response
                # Combine existing + AI-generated for table display
                if ai_generated_rows:
                    table_data = table_data + ai_generated_rows if table_data else ai_generated_rows
                    logger.info(f"Combined table_data: {len(table_data)} rows (existing: {len(existing_rows_formatted)}, AI: {len(ai_generated_rows)})")
            else:
                # Fallback to legacy GPT-2 generation
                prompt = f"""You are a professional cost-optimization consultant for MG Motor India. A colleague has asked you:

"{user_query}"

Based on the following relevant cost reduction ideas from our database:

{context_str}

Please provide a clear, concise response that:
1. Directly answers their question
2. Highlights the most relevant ideas (focus on those matching their query)
3. Summarizes key savings and benefits
4. Mentions any important implementation considerations

Write in a professional, conversational tone. Do not repeat the query or context verbatim. Provide actionable insights.

Response:"""
                MAX_MODEL_LENGTH = 1024
                MAX_NEW_TOKENS = 200

                inputs = tokenizer(prompt, return_tensors="pt", truncation=True, max_length=MAX_MODEL_LENGTH - MAX_NEW_TOKENS)
                inputs = {k: v.to(device) for k, v in inputs.items()}

                with torch.no_grad():
                    generated = inputs['input_ids']
                    attn = inputs['attention_mask']
                    
                    # Define timeout parameters before the loop
                    max_timeout = 10  # Maximum consecutive errors before stopping
                    timeout_count = 0
                    
                    # Repetition detection
                    recent_tokens = []  # Track last N tokens to detect repetition
                    repetition_threshold = 5  # If same token appears N times in a row, stop
                    last_token_count = {}
                    max_repetition_count = 0

                    for i in range(MAX_NEW_TOKENS):
                        try:
                            # Single model call per iteration - optimized for performance
                            logits = model(generated, attn)
                            probs = torch.softmax(logits[:, -1, :], dim=-1)
                            next_token = torch.multinomial(probs, 1)
                            token_id = next_token.item()
                            
                            # Reset timeout counter on successful token generation
                            timeout_count = 0
                            
                            # Repetition detection
                            if token_id in last_token_count:
                                last_token_count[token_id] += 1
                            else:
                                last_token_count = {token_id: 1}
                            
                            if last_token_count[token_id] >= repetition_threshold:
                                logger.warning(f"Repetition detected at token {i}, stopping generation")
                                break
                            
                            # Track recent tokens (keep last 10)
                            recent_tokens.append(token_id)
                            if len(recent_tokens) > 10:
                                recent_tokens.pop(0)
                            
                            if token_id == tokenizer.eos_token_id:
                                logger.info(f"Generation completed at token {i} with EOS token")
                                break
                            
                            generated = torch.cat([generated, next_token], dim=-1)
                            attn = torch.cat([attn, torch.ones_like(next_token)], dim=-1)
                            
                            if generated.size(1) >= MAX_MODEL_LENGTH:
                                logger.warning("Generation stopped: Reached model context limit (1024)")
                                break

                        except Exception as e:
                            logger.error(f"Token generation error at step {i}: {e}")
                            timeout_count += 1
                            if timeout_count > max_timeout:
                                logger.error("Generation stopped: repeated errors")
                                break
                            continue

                response = tokenizer.decode(generated[0], skip_special_tokens=True)
                # Extract response after "Response:" marker, or use full response if marker not found
                if "Response:" in response:
                    response = response.split("Response:")[-1].strip()
                elif "### Response:" in response:
                    response = response.split("### Response:")[-1].strip()
                
                # Clean up common artifacts
                response = response.replace("### Instruction:", "").replace("### Context:", "").replace("### Guidelines:", "").strip()
                
                # If response is too short or seems incomplete, try to extract meaningful part
                if len(response) < 20:
                    # Try to get text after the last meaningful separator
                    parts = response.split("\n")
                    if len(parts) > 1:
                        response = "\n".join(parts[-2:]).strip()

                # Add AI-generated ideas to complement DB ideas when Auto-Analyst is unavailable
                if not ai_generated_rows:
                    # Use existing table columns if available, otherwise fall back to defaults
                    ai_columns = list(table_data[0].keys()) if table_data else DEFAULT_ANALYST_COLUMNS
                    ai_generated_rows = generate_ai_ideas(user_query, table_data, ai_columns, count=3)

                if ai_generated_rows:
                    table_data = table_data + ai_generated_rows
                    # Summarize AI ideas inline so user sees both DB and new proposals
                    ai_summary_lines = []
                    for idea in ai_generated_rows:
                        idea_title = idea.get("Cost Reduction Idea", idea.get("cost_reduction_idea", "New proposal"))
                        saving_val = idea.get("Saving Value (INR)", idea.get("saving_value_inr", "N/A"))
                        ai_summary_lines.append(f"- {idea_title} (est. saving: {saving_val})")
                    if ai_summary_lines:
                        response += "\n\nAdditional AI-suggested ideas:\n" + "\n".join(ai_summary_lines)

            # --- Store LLM result in database (asynchronously) ---
            idea_id_to_log = table_data[0].get('Idea Id') if table_data else 'N/A'
            executor.submit(log_query_to_db, 
                            username,  # Correct parameter order: username first
                            user_query, 
                            response, 
                            None)  # image_path is None for LLM path

            # Files are now generated on-demand when download buttons are clicked
            # This saves API service costs by not generating files automatically
            
            return {
                "response_text": response,
                "table_data": table_data,
                "generated_ideas": ai_generated_rows,
                "image_urls": [row.get('Proposal Image Path', 'N/A') for row in table_data if row.get('Proposal Image Path') != 'N/A'],
                "excel_url": "",  # Will be generated on-demand
                "ppt_url": ""  # Will be generated on-demand
            }
        
    except Exception as e:
        logger.error(f"Error in generate_response: {str(e)}")
        return {
            "response_text": f"Error processing query: {str(e)}",
            "table_data": [], "image_urls": [], "excel_url": "", "ppt_url": "",
            "vlm_warning": str(e) if image_path else None
        }

# --- generate_html_table (Unchanged) ---
def generate_html_table(table_data):
    if not table_data:
        return "<p>No relevant ideas found.</p>"
    headers = [
        'Idea Id', 'Cost Reduction Idea', 'Proposal Image', 'MGI Gross Saving',
        'Way Forward', 'Estimated Cost Savings', 'MGI Carline', 'Saving Value (INR)',
        'Weight Saving (Kg)', 'Status', 'Dept', 'KD/LC'
    ]
    html = ['<table class="table table-striped table-hover">']
    html.append('<thead class="table-dark">')
    html.append('<tr>')
    html.extend([f'<th>{h}</th>' for h in headers])
    html.append('</tr>')
    html.append('</thead>')
    html.append('<tbody>')
    for row in table_data:
        html.append('<tr>')
        for key in headers:
            if key == 'Proposal Image':
                img = row.get('Proposal Image Path', 'N/A')
                html.append(f'<td><img src="{img}" alt="Proposal" style="width:50px;height:50px;"></td>' if img != 'N/A' else '<td>N/A</td>')
            else:
                html.append(f'<td>{row.get(key, "N/A")}</td>')
        html.append('</tr>')
    html.append('</tbody>')
    html.append('</table>')
    return ''.join(html)
# SQL with no endpoints
'''def generate_response(user_query, image_path=None):
    """Generate response using VLM or LLM based on input."""
    try:
        if image_path:
            # VLM processing
            if not vlm_initialized:
                setup_vlm()
                if not vlm_initialized:
                    logger.warning("VLM initialization failed, falling back to LLM")
                    table_data, response_text = retrieve_context(user_query)
                    return {
                        "response_text": response_text or "No relevant cost reduction ideas found",
                        "table_data": table_data,
                        "image_urls": [],
                        "excel_url": "",
                        "ppt_url": "",
                        "vlm_warning": "VLM unavailable due to initialization failure"
                    }

            image = Image.open(image_path)
            error, generated_text, table_data, similarity_scores, image_urls = infer_vlm(
                image, user_query, top_k=10
            )

            if error:
                logger.error(f"VLM processing failed: {error}")
                table_data, response_text = retrieve_context(user_query)
                return {
                    "response_text": response_text or "No relevant cost reduction ideas found",
                    "table_data": table_data,
                    "image_urls": [],
                    "excel_url": "",
                    "ppt_url": "",
                    "vlm_warning": f"VLM processing failed: {error}"
                }

            # Store VLM result in database
            with db_lock:
                with sqlite3.connect(DB_PATH) as conn:
                    cursor = conn.cursor()
                    cursor.execute("""
                        INSERT INTO ideas (
                            idea_id, cost_reduction_idea, mgi_gross_saving, way_forward,
                            estimated_cost_savings, mgi_carline, saving_value_inr, weight_saving,
                            status, dept, kd_lc, proposal_image_filename, user_input, response_text, image_path
                        ) VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?)
                    """, (
                        table_data[0].get('Idea Id') if table_data else 'N/A',
                        table_data[0].get('Cost Reduction Idea') if table_data else user_query,
                        safe_float_convert(str(table_data[0].get('MGI Gross Saving', '0')).replace('INR ', '').replace(',', '')) if table_data else 0,  # FIXED: Using safe conversion
                        table_data[0].get('Way Forward') if table_data else 'N/A',
                        table_data[0].get('Estimated Cost Savings') if table_data else 'N/A',
                        table_data[0].get('MGI Carline') if table_data else 'N/A',
                        safe_float_convert(str(table_data[0].get('Saving Value (INR)', '0')).replace('INR ', '').replace(',', '')) if table_data else 0,  # FIXED: Using safe conversion
                        safe_float_convert(table_data[0].get('Weight Saving (Kg)', '0')) if table_data else 0,  # FIXED: Using safe conversion
                        table_data[0].get('Status') if table_data else 'N/A',
                        table_data[0].get('Dept') if table_data else 'N/A',
                        table_data[0].get('KD/LC') if table_data else 'N/A',
                        os.path.basename(image_path) if image_path else None,
                        user_query,
                        generated_text,
                        os.path.basename(image_path) if image_path else None
                    ))
                    conn.commit()

            return {
                "response_text": generated_text,
                "table_data": table_data,
                "image_urls": image_urls['proposal'],
                "excel_url": "",
                "ppt_url": "",
                "vlm_warning": None
            }

        else:
            # LLM processing
            # Use pre-retrieved data if provided (from agentic RAG), otherwise retrieve fresh
            if pre_retrieved_data is not None:
                table_data, context_str = pre_retrieved_data
                logger.info(f"Using pre-retrieved data: {len(table_data)} ideas")
            else:
                table_data, context_str = retrieve_context(user_query)

            if not table_data:
                logger.warning(f"No relevant ideas found for query: {user_query}")
                return {
                    "response_text": f"No relevant cost reduction ideas found for query: {user_query}",
                    "table_data": [],
                    "image_urls": [],
                    "excel_url": "",
                    "ppt_url": ""
                }

            prompt = f"""
### Instruction:
You are a cost-optimization expert. Analyse the following cost-reduction query.

Query: {user_query}

### Context:
{context_str}

### Guidelines:
1. Analyse only the cost-reduction ideas relevant to the query.
2. Compare the proposed solutions, focusing on feasibility and impact.
3. Highlight potential savings and investment requirements.
4. Consider department, status, and implementation dates.
5. Provide actionable recommendations.

### Response:
"""

            inputs = tokenizer(prompt, return_tensors="pt", truncation=True, max_length=1024)
            inputs = {k: v.to(device) for k, v in inputs.items()}

            with torch.no_grad():
                generated = inputs['input_ids']
                attn = inputs['attention_mask']
                
                # FIXED: Add timeout and better loop control to prevent infinite loops
                max_tokens = 512
                timeout_count = 0
                max_timeout = 50  # Maximum iterations without progress
                
                for i in range(max_tokens):
                    try:
                        logits = model(generated, attn)
                        probs = torch.softmax(logits[:, -1, :], dim=-1)
                        next_token = torch.multinomial(probs, 1)
                        
                        # Check if we got a valid token
                        if next_token.item() == tokenizer.eos_token_id:
                            logger.info(f"Generation completed at token {i} with EOS token")
                            break
                            
                        generated = torch.cat([generated, next_token], dim=-1)
                        attn = torch.cat([attn, torch.ones_like(next_token)], dim=-1)
                        
                        # Reset timeout counter on successful token generation
                        timeout_count = 0
                        
                        # Add additional stopping conditions
                        if generated.size(1) > 2048:  # Prevent extremely long sequences
                            logger.warning("Generation stopped due to maximum sequence length")
                            break
                            
                    except Exception as e:
                        logger.error(f"Error during token generation at step {i}: {e}")
                        timeout_count += 1
                        if timeout_count > max_timeout:
                            logger.error("Generation stopped due to repeated errors")
                            break
                        continue

            response = tokenizer.decode(generated[0], skip_special_tokens=True)
            response = response.split("### Response:")[-1].strip()

            # Store LLM result in database
            with db_lock:
                with sqlite3.connect(DB_PATH) as conn:
                    cursor = conn.cursor()
                    cursor.execute("""
                        INSERT INTO ideas (
                            idea_id, cost_reduction_idea, mgi_gross_saving, way_forward,
                            estimated_cost_savings, mgi_carline, saving_value_inr, weight_saving,
                            status, dept, kd_lc, user_input, response_text
                        ) VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?)
                    """, (
                        table_data[0].get('Idea Id') if table_data else 'N/A',
                        table_data[0].get('Cost Reduction Idea') if table_data else user_query,
                        safe_float_convert(str(table_data[0].get('MGI Gross Saving', '0')).replace('INR ', '').replace(',', '')) if table_data else 0,  # FIXED: Using safe conversion
                        table_data[0].get('Way Forward') if table_data else 'N/A',
                        table_data[0].get('Estimated Cost Savings') if table_data else 'N/A',
                        table_data[0].get('MGI Carline') if table_data else 'N/A',
                        safe_float_convert(str(table_data[0].get('Saving Value (INR)', '0')).replace('INR ', '').replace(',', '')) if table_data else 0,  # FIXED: Using safe conversion
                        safe_float_convert(table_data[0].get('Weight Saving (Kg)', '0')) if table_data else 0,  # FIXED: Using safe conversion
                        table_data[0].get('Status') if table_data else 'N/A',
                        table_data[0].get('Dept') if table_data else 'N/A',
                        table_data[0].get('KD/LC') if table_data else 'N/A',
                        user_query,
                        response
                    ))
                    conn.commit()

            return {
                    "response_text": response,
                    "table_data": table_data,   # <-- This MUST be provided for the table to render
                    "image_urls": [row.get('Proposal Image Path', 'N/A') for row in table_data if row.get('Proposal Image Path') != 'N/A'],
                    "excel_url": "",
                    "ppt_url": ""
                }
        
    except Exception as e:
        logger.error(f"Error in generate_response: {str(e)}")
        return {
            "response_text": f"Error processing query: {str(e)}",
            "table_data": [],
            "image_urls": [],
            "excel_url": "",
            "ppt_url": "",
            "vlm_warning": str(e) if image_path else None
        }
'''
def generate_html_table(table_data):
    """Generate HTML table from table_data."""
    if not table_data:
        return "<p>No relevant ideas found.</p>"

    headers = [
        'Idea Id', 'Cost Reduction Idea', 'Proposal Image', 'MGI Gross Saving',
        'Way Forward', 'Estimated Cost Savings', 'MGI Carline', 'Saving Value (INR)',
        'Weight Saving (Kg)', 'Status', 'Dept', 'KD/LC'
    ]

    html = ['<table class="table table-striped table-hover">']
    html.append('<thead class="table-dark">')
    html.append('<tr>')
    html.extend([f'<th>{h}</th>' for h in headers])
    html.append('</tr>')
    html.append('</thead>')
    html.append('<tbody>')

    for row in table_data:
        html.append('<tr>')
        for key in headers:
            if key == 'Proposal Image':
                img = row.get('Proposal Image Path', 'N/A')
                html.append(f'<td><img src="{img}" alt="Proposal" style="width:50px;height:50px;"></td>' if img != 'N/A' else '<td>N/A</td>')
            else:
                html.append(f'<td>{row.get(key, "N/A")}</td>')
        html.append('</tr>')

    html.append('</tbody>')
    html.append('</table>')
    return ''.join(html)

def generate_csv_from_table(table_data):
    """
    Generates a CSV string from a list of dictionaries (table data).
    Dynamic: Automatically detects columns from the data itself.
    """
    import io
    import csv
    
    if not table_data:
        return ""
    
    output = io.StringIO()
    
    # Dynamic Headers: Take keys from the first row of data
    # This prevents the "dict contains fields not in fieldnames" error
    if len(table_data) > 0:
        headers = list(table_data[0].keys())
        writer = csv.DictWriter(output, fieldnames=headers, extrasaction='ignore')
        writer.writeheader()
        writer.writerows(table_data)
        
    return output.getvalue()

def generate_ppt_from_response(response_text, table_data):
    try:
        prs = Presentation()
        # Set theme colors
        primary_color = RGBColor(66, 133, 244)  # Google Blue
        secondary_color = RGBColor(220, 53, 69)  # Red
        font_name = "Calibri"

        # Helper: Add footer at bottom right
        def add_footer(slide, text="Generated by VAVE-AI"):
            left = Inches(6.5)
            top = Inches(6.8)
            width = Inches(3.0)
            height = Inches(0.3)
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.text = text
            p.font.size = Pt(10)
            p.font.color.rgb = RGBColor(150, 150, 150)
            p.alignment = PP_ALIGN.RIGHT

        # Slide 1: Title Slide with Logo and Custom Text
        title_slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        # Add logo
        logo_path = os.path.join(STATIC_DIR, "images", "mg.png")
        if os.path.exists(logo_path):
            try:
                title_slide.shapes.add_picture(logo_path, Inches(0.5), Inches(0.5), width=Inches(2.0))
            except Exception as e:
                logger.warning(f"Failed to add logo: {str(e)}")
        else:
            logger.warning(f"Logo file not found at {logo_path}")

        # Add text box for title content
        left, top, width, height = Inches(2.7), Inches(1.0), Inches(7.0), Inches(3.0)
        txBox = title_slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True

        # JSW MG MOTOR INDIA
        p = tf.add_paragraph()
        p.text = "JSW MG MOTOR INDIA"
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = secondary_color
        p.font.name = font_name

        # Cost Reduction Idea Analysis Presentation
        p = tf.add_paragraph()
        p.text = "Cost Reduction Idea Analysis Presentation"
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 0, 0)
        p.font.name = font_name

        # Presentation generated by VAVE AI
        p = tf.add_paragraph()
        p.text = "Presentation generated by VAVE AI"
        p.font.size = Pt(16)
        p.font.name = font_name

        add_footer(title_slide)

        # Slide 2: Ideas Summary Slide with Table 
        ideas_per_slide = 5
        for slide_idx, start_idx in enumerate(range(0, len(table_data), ideas_per_slide)):
            summary_slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank slide
            title_shape = summary_slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9.0), Inches(0.4))
            title_frame = title_shape.text_frame
            p = title_frame.paragraphs[0]
            p.text = f"Ideas Summary (Part {slide_idx + 1})"
            p.font.size = Pt(10)
            p.font.bold = True
            p.font.name = font_name

            # Prepare table layout
            num_rows = min(len(table_data) - start_idx, ideas_per_slide) + 1  # Data rows + header
            cols = 8
            left, top, width, height = Inches(0.2), Inches(0.7), Inches(9.6), Inches(6.5)  # Adjusted for fit
            table = summary_slide.shapes.add_table(num_rows, cols, left, top, width, height).table

            # Adjust column widths: optimized for content
            column_widths = [0.8, 2.5, 0.6, 1.0, 0.8, 1.8, 1.0, 1.1]  # Total ~9.6 inches
            for i, w in enumerate(column_widths):
                table.columns[i].width = Inches(w)

            # Headers
            headers = ["Idea ID", "Proposal", "Imp", "Saving Value", "Status", "Way Forward", "Responsibility", "Date"]
            for col_idx, header in enumerate(headers):
                cell = table.cell(0, col_idx)
                cell.text = header
                for para in cell.text_frame.paragraphs:
                    para.font.size = Pt(10)
                    para.font.bold = True
                    para.font.name = font_name
                cell.text_frame.word_wrap = True

            # Add data rows
            for row_idx, idea in enumerate(table_data[start_idx:start_idx + ideas_per_slide], start=1):
                row_values = [
                    idea.get("Idea ID", ""),
                    idea.get("Proposal", ""),
                    idea.get("Imp", ""),
                    idea.get("Saving Value", ""),
                    idea.get("Status", ""),
                    idea.get("Way Forward", ""),
                    idea.get("Responsibility", ""),
                    idea.get("Date", "")
                ]
                for col_idx, value in enumerate(row_values):
                    cell = table.cell(row_idx, col_idx)
                    cell.text = str(value)
                    for para in cell.text_frame.paragraphs:
                        para.font.size = Pt(10)
                        para.font.name = font_name
                    cell.text_frame.word_wrap = True

                # Set row height for compact fit
                table.rows[row_idx].height = Inches(0.35)

            add_footer(summary_slide, "Generated by VAVE-AI")
            
        # Slide 3: Chart Analysis Discussion with Bar Chart
        chart_slide = prs.slides.add_slide(prs.slide_layouts[5])
        title = chart_slide.shapes.title
        title.text = "Chart Analysis Discussion"
        for para in title.text_frame.paragraphs:
            para.font.size = Pt(32)
            para.font.bold = True
            para.font.name = font_name

        # Create bar chart data
        chart_data = CategoryChartData()
        chart_data.categories = [idea["Idea ID"] for idea in table_data]
        savings = [float(str(idea["Saving Value"]).replace("INR ", "").replace(",", "")) if 'INR' in str(idea["Saving Value"]) else 0 for idea in table_data]
        chart_data.add_series("Savings (INR)", savings)

        # Add bar chart
        x, y, cx, cy = Inches(0.5), Inches(1.5), Inches(9.0), Inches(4.5)
        chart = chart_slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
        ).chart
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.chart_title.text_frame.text = "Cost Savings by Idea"
        for para in chart.chart_title.text_frame.paragraphs:
            para.font.size = Pt(16)
            para.font.name = font_name

        add_footer(chart_slide)

        # Slide 4: Pie Chart Slide
        pie_slide = prs.slides.add_slide(prs.slide_layouts[5])
        title = pie_slide.shapes.title
        title.text = "Pie chart for ideas status"
        for para in title.text_frame.paragraphs:
            para.font.size = Pt(32)
            para.font.bold = True
            para.font.name = font_name

        # Create pie chart data
        status_counts = {"Implemented": 0, "Status OK": 0, "Status TBD": 0, "Status NG": 0}
        for idea in table_data:
            status = idea.get("Status", "TBD")
            if status == "OK":
                status_counts["Status OK"] += 1
            elif status == "TBD":
                status_counts["Status TBD"] += 1
            elif status == "NG":
                status_counts["Status NG"] += 1
            elif status == "Implemented":
                status_counts["Implemented"] += 1

        chart_data = CategoryChartData()
        chart_data.categories = ["Implemented", "Status OK", "Status TBD", "Status NG"]
        chart_data.add_series("Status", [status_counts["Implemented"], status_counts["Status OK"], status_counts["Status TBD"], status_counts["Status NG"]])

        # Add pie chart
        x, y, cx, cy = Inches(0.5), Inches(1.5), Inches(4.5), Inches(4.5)
        chart = pie_slide.shapes.add_chart(
            XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data
        ).chart
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.RIGHT
        chart.legend.include_in_layout = False
        chart.chart_title.text_frame.text = "Ideas Status Distribution"
        for para in chart.chart_title.text_frame.paragraphs:
            para.font.size = Pt(16)
            para.font.name = font_name

        # Customize pie chart colors
        chart.plots[0].series[0].points[0].format.fill.solid()  # Implemented (blue)
        chart.plots[0].series[0].points[0].format.fill.fore_color.rgb = RGBColor(0, 0, 255)
        chart.plots[0].series[0].points[1].format.fill.solid()  # Status OK (red)
        chart.plots[0].series[0].points[1].format.fill.fore_color.rgb = RGBColor(255, 0, 0)
        chart.plots[0].series[0].points[2].format.fill.solid()  # Status TBD (green)
        chart.plots[0].series[0].points[2].format.fill.fore_color.rgb = RGBColor(0, 128, 0)
        chart.plots[0].series[0].points[3].format.fill.solid()  # Status NG (purple)
        chart.plots[0].series[0].points[3].format.fill.fore_color.rgb = RGBColor(128, 0, 128)

        # Add summary text on the right side
        left, top, width, height = Inches(5.5), Inches(1.5), Inches(4.0), Inches(4.5)
        txBox = pie_slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        total_savings = sum(float(str(idea["Saving Value"]).replace("INR ", "").replace(",", "")) if 'INR' in str(idea["Saving Value"]) else 0 for idea in table_data)
        summary_text = [
            f"• Total Ideas: {len(table_data)}",
            f"• Total Potential Savings: INR {total_savings:,.2f}",
            f"• Implemented Ideas: {status_counts['Implemented']}",
            f"• Status OK: {status_counts['Status OK']}",
            f"• Status TBD: {status_counts['Status TBD']}",
            f"• Status NG: {status_counts['Status NG']}",
            "• The provided cost reduction ideas focus on process deletions and material replacements to reduce manufacturing costs.",
            "• Key proposals include removing paint coatings from brackets and mounts, replacing PC-ABS with ABS, and switching to galvanized steel for window regulators.",
            "• These changes aim to maintain functionality while reducing costs and, in some cases, weight."
        ]
        for line in summary_text:
            p = tf.add_paragraph()
            p.text = line
            for para in tf.paragraphs:
                para.font.size = Pt(10)
                para.font.name = font_name
        add_footer(pie_slide)

        # Slide 5: Summary Slide with Bullet Points
        summary_slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = summary_slide.shapes.title
        title.text = "Summary section"
        for para in title.text_frame.paragraphs:
            para.font.size = Pt(32)
            para.font.bold = True
            para.font.name = font_name
        body_shape = summary_slide.placeholders[1]
        tf = body_shape.text_frame
        tf.word_wrap = True
        summary_text = [
            f"Total Ideas: {len(table_data)} ideas analyzed.",
            f"Total Savings: INR {total_savings:,.2f} in potential cost reductions.",
            "Key Proposals: Remove paint coatings from mounts and brackets, replace PC-ABS with ABS, and use galvanized steel for window regulators.",
            "Status: All ideas currently TBD, pending further evaluation.",
            "Focus: Cost reduction through process deletion and material substitution while maintaining functionality."
        ]
        for line in summary_text:
            p = tf.add_paragraph()
            p.text = line
            p.level = 0
        for para in tf.paragraphs:
            para.font.size = Pt(18)
            para.font.name = font_name
        add_footer(summary_slide)

        ppt_filename = f"response_{id(response_text)}.pptx"
        ppt_path = os.path.join(UPLOAD_FOLDER, ppt_filename)
        prs.save(ppt_path)
        logger.info(f"PPT file generated: {ppt_filename}")
        return ppt_filename
    except Exception as e:
        logger.error(f"Error generating PPT: {str(e)}")
        import traceback
        traceback.print_exc()
        return None


def allowed_file(filename):
    """Check if file extension is allowed."""
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS
# Flask Routes
@app.route('/')
def route_index():
    """
    This is the new starting point.
    It checks if you are logged in.
    - If YES, it sends you to '/app'.
    - If NO, it sends you to '/login'.
    """
    if current_user.is_authenticated:
        return redirect(url_for('chat_app'))
    return redirect(url_for('login'))

@app.route('/app')
@login_required
def chat_app():
    # 1. Fetch Users List (Only if the user is Admin/SuperAdmin)
    users = []
    if current_user.role in ['SuperAdmin', 'Admin']:
        # Make sure schema is up to date (adds role/created_at if missing)
        try:
            init_user_db_schema()
        except Exception as e:
            logger.error(f"init_user_db_schema failed in chat_app: {e}")

        conn = get_user_db_conn()
        cursor = conn.cursor()
        try:
            # Detect columns reliably to avoid OperationalError warnings
            cursor.execute("PRAGMA table_info(users)")
            cols = [row[1] for row in cursor.fetchall()]
            has_created_at = 'created_at' in cols
            has_role = 'role' in cols

            if has_created_at and has_role:
                cursor.execute("SELECT id, username, role, created_at FROM users")
                users = cursor.fetchall()
            elif has_role:
                cursor.execute("SELECT id, username, role FROM users")
                tmp_rows = cursor.fetchall()
                users = [(r[0], r[1], r[2], None) for r in tmp_rows]
            else:
                cursor.execute("SELECT id, username FROM users")
                tmp_rows = cursor.fetchall()
                users = [(r[0], r[1], None, None) for r in tmp_rows]
        finally:
            conn.close()

    # 2. Render the MAIN Chat App (passing dashboard data to it)
    return render_template(
        'chat_app.html', 
        username=current_user.username,
        current_user=current_user,
        users=users,
        now=datetime.now().strftime("%Y-%m-%d %H:%M")
    )


@app.route('/dashboard')
@login_required
def dashboard():
    """
    Profile / dashboard page with basic stats.
    Uses PostgreSQL chat_history (if present) and SQLite users.db.
    """
    stats = {
        "total_chats": 0,
        "recent_activity": 0,
        "total_users": 0,
    }
    recent_activity = []

    # Chat stats from PostgreSQL
    try:
        conn = get_db_connection()
        cur = conn.cursor()

        # Total chats for current user
        cur.execute(
            "SELECT COUNT(*) FROM chat_history WHERE username = %s",
            (current_user.username,),
        )
        stats["total_chats"] = cur.fetchone()[0] or 0

        # Recent activity: count of chats in last 7 days
        cur.execute(
            """
            SELECT COUNT(*)
            FROM chat_history
            WHERE username = %s
              AND created_at >= NOW() - INTERVAL '7 days'
            """,
            (current_user.username,),
        )
        stats["recent_activity"] = cur.fetchone()[0] or 0

        # Recent activity list (last 5 messages)
        cur.execute(
            """
            SELECT user_query, created_at
            FROM chat_history
            WHERE username = %s
            ORDER BY created_at DESC
            LIMIT 5
            """,
            (current_user.username,),
        )
        rows = cur.fetchall()
        recent_activity = [
            {
                "message": r[0],
                "timestamp": r[1],
            }
            for r in rows
        ]

        cur.close()
        conn.close()
    except Exception as e:
        logger.error(f"Dashboard chat stats error: {e}")

    # User stats from SQLite (only relevant for Admin/SuperAdmin display)
    if current_user.role in ["SuperAdmin", "Admin"]:
        try:
            conn_u = get_user_db_conn()
            cur_u = conn_u.cursor()
            cur_u.execute("SELECT COUNT(*) FROM users")
            stats["total_users"] = cur_u.fetchone()[0] or 0
            cur_u.close()
            conn_u.close()
        except Exception as e:
            logger.error(f"Dashboard user stats error: {e}")

    return render_template(
        "dashboard.html",
        current_user=current_user,
        stats=stats,
        recent_activity=recent_activity,
    )


# --- USER MANAGEMENT & ADMIN ROUTES ---

@app.route('/change_password', methods=['POST'])
@login_required
def change_password():
    new_pass = request.form.get('new_password')
    if not new_pass:
        flash("Password cannot be empty.", "error")
        return redirect(url_for('chat_app', section='dashboard'))
        
    hashed_pw = generate_password_hash(new_pass)
    
    conn = get_user_db_conn()
    conn.execute("UPDATE users SET password_hash = ? WHERE id = ?", (hashed_pw, current_user.id))
    conn.commit()
    conn.close()
    
    flash("Your password has been updated successfully.", "success")
    return redirect(url_for('chat_app', section='dashboard'))

@app.route('/create_user', methods=['POST'])
@login_required
def create_user():
    """
    Create user.
    - If JSON request: return JSON (used by SPA UI).
    - If form request: redirect with flash messages.
    """
    # Security Check: Only SuperAdmin can create users
    if current_user.role != 'SuperAdmin':
        if request.is_json:
            return jsonify({"success": False, "error": "Permission denied"}), 403
        flash("Permission denied", "error")
        return redirect(url_for('chat_app', section='dashboard'))

    if request.is_json:
        data = request.get_json(silent=True) or {}
        username = data.get('username')
        password = data.get('password')
        role = data.get('role', 'User')
    else:
        username = request.form.get('username')
        password = request.form.get('password')
        role = request.form.get('role', 'User')

    if not username or not password:
        msg = "Username and password are required."
        if request.is_json:
            return jsonify({"success": False, "error": msg}), 400
        flash(msg, "error")
        return redirect(url_for('chat_app', section='dashboard'))

    conn = get_user_db_conn()
    try:
        hashed_pw = generate_password_hash(password)
        conn.execute(
            "INSERT INTO users (username, password_hash, role) VALUES (?, ?, ?)",
            (username, hashed_pw, role),
        )
        conn.commit()
        msg = f'User {username} created successfully!'
        if request.is_json:
            return jsonify({"success": True, "message": msg})
        flash(msg, 'success')
    except sqlite3.IntegrityError:
        msg = 'Username already exists'
        if request.is_json:
            return jsonify({"success": False, "error": msg}), 400
        flash(msg, 'error')
    finally:
        conn.close()

    # Redirects back to the main app, but opens the dashboard tab
    return redirect(url_for('chat_app', section='dashboard'))

@app.route('/delete_user/<int:user_id>')
@login_required
def delete_user(user_id):
    # Security Check: Only SuperAdmin can delete users
    if current_user.role != 'SuperAdmin':
        flash("Permission denied", "error")
        return redirect(url_for('chat_app', section='dashboard'))

    if user_id == current_user.id:
        flash("You cannot delete yourself!", "error")
        return redirect(url_for('chat_app', section='dashboard'))
        
    conn = get_user_db_conn()
    conn.execute("DELETE FROM users WHERE id = ?", (user_id,))
    conn.commit()
    conn.close()
    flash('User deleted.', 'success')
    return redirect(url_for('chat_app', section='dashboard'))


@app.route('/delete_user', methods=['POST'])
@login_required
def delete_user_api():
    """
    JSON API used by SPA to delete a user by username.
    """
    if current_user.role != 'SuperAdmin':
        return jsonify({"success": False, "error": "Permission denied"}), 403

    data = request.get_json(silent=True) or {}
    username = data.get('username')
    if not username:
        return jsonify({"success": False, "error": "Username is required"}), 400

    conn = get_user_db_conn()
    cursor = conn.cursor()
    try:
        # Prevent deleting self
        cursor.execute("SELECT id FROM users WHERE username = ?", (username,))
        row = cursor.fetchone()
        if not row:
            return jsonify({"success": False, "error": "User not found"}), 404

        user_id = row["id"] if isinstance(row, sqlite3.Row) else row[0]
        if user_id == current_user.id:
            return jsonify({"success": False, "error": "You cannot delete yourself"}), 400

        cursor.execute("DELETE FROM users WHERE id = ?", (user_id,))
        conn.commit()
        return jsonify({"success": True})
    finally:
        conn.close()


@app.route('/login', methods=['GET', 'POST'])
def login():
    """Handles user login."""
    
    # --- Logging for GET request ---
    if request.method == 'GET':
        logger.info("Serving login page (GET request).")
        if current_user.is_authenticated:
            logger.info(f"User '{current_user.username}' is already logged in, redirecting to /app.")
            return redirect(url_for('chat_app'))

    # --- Logging for POST request (login attempt) ---
    if request.method == 'POST':
        username = request.form.get('username')
        password = request.form.get('password')
        
        # +++ NEW LOG 1 +++
        logger.info(f"Login attempt started for user: '{username}'")

        conn = None
        try:
            conn = get_user_db_conn()
            cursor = conn.cursor()
            cursor.execute("SELECT * FROM users WHERE username = ?", (username,))
            user_row = cursor.fetchone()
            
            # --- Check if user was found ---
            if user_row:
                # +++ NEW LOG 2 +++
                logger.info(f"User '{username}' found in database. Checking password...")
                
                # --- Check if password hash matches ---
                password_is_correct = check_password_hash(user_row['password_hash'], password)
                
                if password_is_correct:
                    # +++ NEW LOG 3 (Success!) +++
                    logger.info(f"Password for '{username}' is CORRECT. Logging user in.")
                    
                    # Extract role from user_row, defaulting to 'User' if missing
                    role = user_row.get('role', 'User') if 'role' in user_row.keys() else 'User'
                    user = User(
                        id=user_row['id'], 
                        username=user_row['username'], 
                        password_hash=user_row['password_hash'],
                        role=role  # Include role parameter for proper RBAC
                    )
                    login_user(user, remember=True)
                    return redirect(url_for('chat_app'))
                else:
                    # +++ NEW LOG 4 (Failure) +++
                    logger.warning(f"Password for '{username}' is INCORRECT.")
                    flash('Invalid username or password.', 'error')
            
            else:
                # +++ NEW LOG 5 (Failure) +++
                logger.warning(f"User '{username}' NOT FOUND in database.")
                flash('Invalid username or password.', 'error')
                
        except Exception as e:
            # +++ MODIFIED LOG 6 (Critical Error) +++
            logger.error(f"DATABASE or APPLICATION ERROR during login for '{username}': {e}", exc_info=True)
            flash('An error occurred during login.', 'error')
        finally:
            if conn:
                conn.close()
                
    # This will run for a GET request or a failed POST request
    return render_template('login.html')

@app.route('/register', methods=['GET', 'POST'])
def register():
    """Handles new user registration."""
    if current_user.is_authenticated:
        return redirect(url_for('chat_app'))
    
    if request.method == 'POST':
        username = request.form.get('username')
        password = request.form.get('password')
        
        if not username or not password:
            flash('Username and password are required.', 'error')
            return render_template('register.html')
            
        hashed_password = generate_password_hash(password)
        
        conn = None
        try:
            conn = get_user_db_conn()
            cursor = conn.cursor()
            cursor.execute(
                "INSERT INTO users (username, password_hash, role) VALUES (?, ?, ?)",
                (username, hashed_password, 'User') # Explicitly set as User
            )
            conn.commit()
            logger.info(f"New user '{username}' registered.")
            flash('Registration successful! Please log in.', 'success')
            return redirect(url_for('login'))
            
        except sqlite3.IntegrityError:
            logger.warning(f"Registration failed: Username '{username}' already exists.")
            flash('Username already exists. Please choose another.', 'error')
        except Exception as e:
            logger.error(f"Error during registration: {e}")
            flash('An error occurred during registration.', 'error')
        finally:
            if conn:
                conn.close()
                
    return render_template('register.html')

@app.route('/logout')
@login_required
def logout():
    """Logs the current user out."""
    logger.info(f"User '{current_user.username}' logging out.")
    logout_user()
    flash('You have been logged out.', 'success')
    return redirect(url_for('login'))

@app.route('/chat', methods=['POST'])
@login_required
def chat():
    try:
        user_query = request.form.get('message', '').strip()
        if not user_query:
            logger.warning("Received empty or missing 'message' in POST request")
            return jsonify({
                "success": True,
                "response_text": "Please enter a query to proceed.",
                "table_data": [],
                "image_urls": [],
                "excel_url": "",
                "ppt_url": ""
            }), 200

        # Check if an image was uploaded
        image_path = None
        if 'image' in request.files:
            file = request.files['image']
            if file and allowed_file(file.filename):
                filename = secure_filename(f"{uuid.uuid4()}_{file.filename}")
                image_path = UPLOAD_FOLDER / filename
                file.save(str(image_path))
                logger.info(f"Image uploaded: {image_path}")

        # Log incoming chat event for big-data pipeline
        log_event(
            "chat_request",
            username=current_user.username,
            payload={"user_query": user_query, "has_image": bool(image_path)},
        )

        # Generate response via agentic_rag_chat which uses generate_response with Auto-Analyst
        # This ensures Auto-Analyst logic (dual tables with existing + AI-generated ideas) is always used
        logger.info(f"Processing query via agentic_rag_chat (Auto-Analyst enabled): {user_query[:100]}")
        result = agentic_rag_chat(
            current_user.username,
            user_query,
            image_path=image_path
        )

        # Log table_data contents for debugging & analytics
        table_rows = len(result.get('table_data', []))
        logger.info(f"Returning table_data with {table_rows} rows")

        # Log response event + metric time-series
        log_event(
            "chat_response",
            username=current_user.username,
            payload={
                "rows": table_rows,
                "response_length": len(result.get("response_text") or ""),
            },
        )
        try:
            conn = get_db_connection()
            cur = conn.cursor()
            cur.execute(
                """
                INSERT INTO metrics_timeseries (metric_name, metric_value, labels)
                VALUES (%s, %s, %s)
                """,
                (
                    "chat_response_rows",
                    float(table_rows),
                    json.dumps({"username": current_user.username}),
                ),
            )
            conn.commit()
            cur.close()
            conn.close()
        except Exception as m_err:
            logger.error(f"Failed to log metrics_timeseries: {m_err}")
        if not result.get('table_data'):
            logger.warning("No table data generated for query: %s", user_query)

        # Clean up uploaded image
        if image_path and image_path.exists():
            try:
                os.remove(image_path)
                logger.info(f"Cleaned up uploaded image: {image_path}")
            except Exception as e:
                logger.error(f"Failed to clean up image {image_path}: {str(e)}")

        # Return JSON response without the redundant 'table' field
        response_data = {
            "success": True,
            "response_text": result.get('response_text', 'No response generated'),
            "table_data": result.get('table_data', []),
            "image_urls": result.get('image_urls', []),
            "excel_url": result.get('excel_url', ''),
            "ppt_url": result.get('ppt_url', ''),
            "warning": result.get('vlm_warning')
        }
        
        # Log the response data for debugging
        logger.info(f"Response data - Excel URL: {response_data['excel_url']}, PPT URL: {response_data['ppt_url']}")
        logger.info(f"Table data length: {len(response_data['table_data'])}")
        
        return jsonify(response_data)

    except Exception as e:
        logger.error(f"Error in chat endpoint: {str(e)}")
        return jsonify({
            "success": False,
            "error": f"An error occurred: {str(e)}",
            "table_data": [],
            "image_urls": [],
            "excel_url": "",
            "ppt_url": ""
        }), 500

@app.route('/upload_data', methods=['POST'])
@login_required
def upload_data():
    # Security: Only Admin/SuperAdmin
    if current_user.role not in ['SuperAdmin', 'Admin']:
        return jsonify({'success': False, 'error': 'Permission Denied'})

    if 'excel_file' not in request.files or 'zip_file' not in request.files:
        return jsonify({'success': False, 'error': 'Both Excel and Zip files are required.'})

    excel_file = request.files['excel_file']
    zip_file = request.files['zip_file']

    if excel_file.filename == '' or zip_file.filename == '':
        return jsonify({'success': False, 'error': 'No selected file.'})

    # Save to temp
    excel_path = UPLOAD_FOLDER / secure_filename(excel_file.filename)
    zip_path = UPLOAD_FOLDER / secure_filename(zip_file.filename)
    try:
        excel_file.save(str(excel_path))
        zip_file.save(str(zip_path))
        
        # Define DB Config for the processor
        db_config = {
            "dbname": os.getenv("DB_NAME"),
            "user": os.getenv("DB_USER"),
            "password": os.getenv("DB_PASS"),
            "host": os.getenv("DB_HOST", "localhost"),
            "port": os.getenv("DB_PORT", "5432")
        }

        # Offload to processor
        # Note: This blocks the request until done. For very large files, use Celery/background tasks.
        success, message = data_processor.process_knowledge_base_upload(excel_path, zip_path, db_config)

        # Big-data concept: ingest a CSV snapshot into the data lake for future analytical queries.
        # We generate a CSV of the ideas table and push it into the lake.
        if success:
            try:
                # Use existing helper to export the current ideas into a CSV
                if 'generate_table_csv' in globals():
                    csv_filename = generate_table_csv([])
                    if csv_filename:
                        csv_path = UPLOAD_FOLDER / csv_filename
                        # Move into RAW zone and run ETL
                        raw_path = data_lake.ingest_uploaded_csv(csv_path, logical_dataset="ideas")
                        data_lake.run_etl_on_raw_csv(raw_path, logical_dataset="ideas")
            except Exception as etl_err:
                logger.error(f"ETL to data lake failed (non-fatal): {etl_err}")
        
        if success:
            # Rebuild Vector DB in memory
            build_vector_db()
            return jsonify({'success': True, 'message': message})
        else:
            return jsonify({'success': False, 'error': message})

    except Exception as e:
        logger.error(f"Upload endpoint error: {e}")
        return jsonify({'success': False, 'error': str(e)})
    finally:
        # Clean up uploads
        if excel_path.exists(): os.remove(excel_path)
        if zip_path.exists(): os.remove(zip_path)


@app.route('/analytics/ideas_summary', methods=['GET'])
@login_required
def analytics_ideas_summary():
    """
    Advanced big-data style analytics endpoint.
    Runs aggregated queries over the DuckDB-based data lake (Parquet files).
    If no lake data exists yet, it will attempt a one-time backfill from the
    current PostgreSQL 'ideas' table and then return fresh stats.
    """
    try:
        stats = data_lake.analytics_aggregate_ideas()

        # If lake is empty but the main DB has data, try to backfill once
        if stats.get("total_ideas", 0) == 0:
            try:
                if 'generate_table_csv' in globals():
                    csv_filename = generate_table_csv([])
                    if csv_filename:
                        csv_path = UPLOAD_FOLDER / csv_filename
                        raw_path = data_lake.ingest_uploaded_csv(csv_path, logical_dataset="ideas")
                        data_lake.run_etl_on_raw_csv(raw_path, logical_dataset="ideas")
                        # Recompute stats after backfill
                        stats = data_lake.analytics_aggregate_ideas()
            except Exception as etl_err:
                logger.error(f"Analytics backfill to data lake failed (non-fatal): {etl_err}")

        return jsonify({"success": True, "data": stats})
    except Exception as e:
        logger.error(f"Analytics endpoint error: {e}")
        return jsonify({"success": False, "error": str(e)}), 500


@app.route('/analytics/lake_status', methods=['GET'])
@login_required
def analytics_lake_status():
    """
    Big-data style monitoring endpoint for the DuckDB lake.
    Returns load counts and rows per load_date for 'ideas'.
    """
    if current_user.role not in ['SuperAdmin', 'Admin']:
        return jsonify({"success": False, "error": "Permission denied"}), 403

    try:
        stats = data_lake.lake_status("ideas")
        return jsonify({"success": True, "data": stats})
    except Exception as e:
        logger.error(f"Lake status endpoint error: {e}")
        return jsonify({"success": False, "error": str(e)}), 500


def ensure_data_uploads_table():
    """
    Ensure the PostgreSQL table used to track knowledge-base uploads exists.
    Safe to call before inserts/selects.
    """
    try:
        conn = get_db_connection()
        cursor = conn.cursor()
        cursor.execute(
            """
            CREATE TABLE IF NOT EXISTS data_uploads (
                id SERIAL PRIMARY KEY,
                excel_filename TEXT,
                zip_filename TEXT,
                uploaded_by TEXT,
                status TEXT,
                message TEXT,
                created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
            );
            """
        )
        conn.commit()
        cursor.close()
        conn.close()
    except Exception as e:
        logger.error(f"Failed to ensure data_uploads table exists: {e}")


@app.route('/upload_database', methods=['POST'])
@login_required
def upload_database():
    """
    Database management upload used by SPA.
    """
    if current_user.role not in ['SuperAdmin', 'Admin']:
        return jsonify({"success": False, "error": "Permission denied"}), 403

    # Make sure the audit table exists before we try to insert into it
    ensure_data_uploads_table()

    # Expect both Excel and Zip (images) to fully refresh the knowledge base
    excel_file = request.files.get('excel_file')
    zip_file = request.files.get('zip_file')

    if not excel_file or excel_file.filename == '':
        return jsonify({"success": False, "error": "Please select an Excel file"}), 400
    if not zip_file or zip_file.filename == '':
        return jsonify({"success": False, "error": "Please select an Images ZIP file"}), 400

    excel_path = UPLOAD_FOLDER / secure_filename(excel_file.filename)
    zip_path = UPLOAD_FOLDER / secure_filename(zip_file.filename)

    try:
        # Save both files
        excel_file.save(str(excel_path))
        zip_file.save(str(zip_path))

        logger.info(f"Database Excel uploaded to {excel_path}")
        logger.info(f"Images ZIP uploaded to {zip_path}")

        # DB config for processor
        db_config = {
            "dbname": os.getenv("DB_NAME"),
            "user": os.getenv("DB_USER"),
            "password": os.getenv("DB_PASS"),
            "host": os.getenv("DB_HOST", "localhost"),
            "port": os.getenv("DB_PORT", "5432"),
        }

        # Run full preprocessing + image extraction + DB load
        success, message = data_processor.process_knowledge_base_upload(
            excel_path, zip_path, db_config
        )

        # Log upload in PostgreSQL
        try:
            conn = get_db_connection()
            cursor = conn.cursor()
            cursor.execute(
                """
                INSERT INTO data_uploads (
                    excel_filename, zip_filename, uploaded_by, status, message
                ) VALUES (%s, %s, %s, %s, %s)
                """,
                (
                    excel_file.filename,
                    zip_file.filename,
                    current_user.username,
                    "success" if success else "failed",
                    message,
                ),
            )
            conn.commit()
            cursor.close()
            conn.close()
        except Exception as log_err:
            logger.error(f"Failed to log data_uploads entry: {log_err}")

        if not success:
            return jsonify({"success": False, "error": message}), 500

        # Optionally regenerate vector DB and data lake
        try:
            build_vector_db()
        except Exception as e:
            logger.error(f"Vector DB rebuild after upload failed: {e}")

        # Best-effort: push a CSV snapshot into data lake for big-data analytics
        try:
            if 'generate_table_csv' in globals():
                csv_filename = generate_table_csv([])
                if csv_filename:
                    csv_path = UPLOAD_FOLDER / csv_filename
                    raw_path = data_lake.ingest_uploaded_csv(csv_path, logical_dataset="ideas")
                    data_lake.run_etl_on_raw_csv(raw_path, logical_dataset="ideas")
        except Exception as etl_err:
            logger.error(f"ETL to data lake after upload_database failed (non-fatal): {etl_err}")

        return jsonify({"success": True, "message": message})
    except Exception as e:
        logger.error(f"Upload_database error: {e}")
        return jsonify({"success": False, "error": str(e)}), 500
    finally:
        # Clean up temp files
        try:
            if excel_path.exists():
                excel_path.unlink()
            if zip_path.exists():
                zip_path.unlink()
        except Exception as cleanup_err:
            logger.warning(f"Failed to clean upload temp files: {cleanup_err}")


@app.route('/database_status', methods=['GET'])
@login_required
def database_status():
    """
    Returns basic status of the ideas database for the dashboard.
    """
    if current_user.role not in ['SuperAdmin', 'Admin']:
        return jsonify({"success": False, "error": "Permission denied"}), 403

    total_records = 0
    last_updated = None
    size = "-"

    try:
        conn = get_db_connection()
        cursor = conn.cursor()

        # Total records
        cursor.execute("SELECT COUNT(*) FROM ideas")
        total_records = cursor.fetchone()[0] or 0

        # Last updated
        try:
            cursor.execute("SELECT MAX(created_at) FROM ideas")
            last_updated = cursor.fetchone()[0]
        except Exception:
            last_updated = None

        # Approx size of ideas table
        try:
            cursor.execute("SELECT pg_size_pretty(pg_total_relation_size('public.ideas'));")
            size = cursor.fetchone()[0]
        except Exception:
            size = "-"

        cursor.close()
        conn.close()
    except Exception as e:
        logger.error(f"Database status error: {e}")

    return jsonify({
        "success": True,
        "total_records": total_records,
        "last_updated": str(last_updated) if last_updated else None,
        "size": size,
    })


@app.route('/database_uploads', methods=['GET'])
@login_required
def database_uploads():
    """
    Returns history of knowledge base uploads for admin UI.
    """
    if current_user.role not in ['SuperAdmin', 'Admin']:
        return jsonify({"success": False, "error": "Permission denied"}), 403

    try:
        # Ensure table exists before querying
        ensure_data_uploads_table()

        conn = get_db_connection()
        cursor = conn.cursor()
        cursor.execute(
            """
            SELECT id, excel_filename, zip_filename, uploaded_by, status, message, created_at
            FROM data_uploads
            ORDER BY created_at DESC
            LIMIT 50
            """
        )
        rows = cursor.fetchall()
        cursor.close()
        conn.close()

        history = [
            {
                "id": r[0],
                "excel_filename": r[1],
                "zip_filename": r[2],
                "uploaded_by": r[3],
                "status": r[4],
                "message": r[5],
                "created_at": r[6],
            }
            for r in rows
        ]

        return jsonify({"success": True, "uploads": history})
    except Exception as e:
        logger.error(f"database_uploads error: {e}")
        return jsonify({"success": False, "error": str(e)}), 500

@app.route('/admin/logs')
@login_required
def get_logs():
    if current_user.role != 'SuperAdmin':
        return jsonify({'error': 'Permission denied'}), 403
    
    log_file_path = "logger.log"
    try:
        if os.path.exists(log_file_path):
            # Read last 1000 lines
            with open(log_file_path, 'r') as f:
                lines = f.readlines()
                return jsonify({'logs': ''.join(lines[-1000:])}) 
        return jsonify({'logs': 'No log file found.'})
    except Exception as e:
        return jsonify({'logs': f'Error reading logs: {str(e)}'})

@app.route('/admin/clear_logs', methods=['POST'])
@login_required
def clear_logs():
    if current_user.role != 'SuperAdmin':
        return jsonify({'success': False, 'error': 'Permission denied'})
    
    try:
        with open("logger.log", 'w') as f:
            f.write(f"Logs cleared by {current_user.username} at {datetime.now()}\n")
        return jsonify({'success': True})
    except Exception as e:
        return jsonify({'success': False, 'error': str(e)})

@app.route('/upload_image', methods=['POST'])
@login_required
def upload_image():
    """Handle image upload."""
    try:
        if 'file' not in request.files:
            return jsonify({"error": "No file uploaded"}), 400

        file = request.files['file']
        if file.filename == '':
            return jsonify({"error": "No file selected"}), 400

        if file and allowed_file(file.filename):
            filename = secure_filename(f"{uuid.uuid4()}_{file.filename}")
            filepath = UPLOAD_FOLDER / filename
            file.save(str(filepath))
            
            return jsonify({
                "message": "File uploaded successfully",
                "filename": filename,
                "filepath": str(filepath)
            })
        
        return jsonify({"error": "Invalid file type"}), 400

    except Exception as e:
        logger.error(f"Error in upload_image: {str(e)}")
        return jsonify({"error": f"Upload failed: {str(e)}"}), 500


@app.errorhandler(413)
def request_entity_too_large(error):
    """
    Return a clean JSON error for payload-too-large instead of a generic HTML.
    This is used by the SPA to show a friendly message.
    """
    return (
        jsonify(
            {
                "success": False,
                "error": "Uploaded file is too large. Please upload a smaller file or contact the admin to increase the limit.",
            }
        ),
        413,
    )

@app.route('/download_excel/<filename>')
@login_required
def download_excel(filename):
    """Serve the generated Excel file."""
    try:
        file_path = Path(app.config['UPLOAD_FOLDER']) / filename
        if file_path.exists():
            logger.info(f"Serving Excel file: {file_path}")
            return send_file(file_path, as_attachment=True, download_name=filename)
        else:
            logger.error(f"Excel file not found: {file_path}")
            return jsonify({"error": "File not found"}), 404
    except Exception as e:
        logger.error(f"Error serving Excel file {filename}: {str(e)}")
        return jsonify({"error": f"Download failed: {str(e)}"}), 500

@app.route('/download_ppt/<filename>')
@login_required
def download_ppt(filename):
    """Serve the generated PowerPoint file."""
    try:
        # Check static/reports first (new logic)
        static_report_path = STATIC_DIR / "reports" / filename
        if static_report_path.exists():
            logger.info(f"Serving PPT from static/reports: {static_report_path}")
            return send_file(static_report_path, as_attachment=True)

        # Fallback to UPLOAD_FOLDER (legacy logic)
        file_path = app.config['UPLOAD_FOLDER'] / filename
        if file_path.exists():
            logger.info(f"Serving PPT from uploads: {file_path}")
            return send_file(file_path, as_attachment=True)
            
        logger.error(f"PPT file not found: {filename}")
        return jsonify({"error": "File not found"}), 404
    except Exception as e:
        logger.error(f"Error serving PPT file {filename}: {str(e)}")
        return jsonify({"error": f"Download failed: {str(e)}"}), 500

@app.route('/generate_excel', methods=['POST'])
@login_required
def generate_excel_route():
    """Generate Excel file on-demand from table data."""
    try:
        data = request.json
        table_data = data.get('table_data', [])
        
        if not table_data:
            return jsonify({"success": False, "error": "No data provided"}), 400

        # Generate Excel file directly in upload folder
        upload_folder = Path(app.config['UPLOAD_FOLDER'])
        upload_folder.mkdir(parents=True, exist_ok=True)
        
        filename, file_path = generate_excel_from_table_in_memory(table_data, upload_folder=str(upload_folder))
        
        if not filename:
            return jsonify({"success": False, "error": "Failed to generate Excel file"}), 500
        
        # Return download URL
        download_url = url_for('download_excel', filename=filename)
        
        return jsonify({
            "success": True,
            "download_url": download_url,
            "filename": filename,
            "message": "Excel file generated successfully."
        })

    except Exception as e:
        logger.error(f"Excel Generation Error: {e}")
        import traceback
        traceback.print_exc()
        return jsonify({"success": False, "error": str(e)}), 500


@app.route('/generate_ppt', methods=['POST'])
@login_required
def generate_ppt_route():
    """Generate comprehensive PowerPoint file on-demand using new VAVE Presentation Engine."""
    try:
        data = request.json
        table_data = data.get('table_data', [])
        response_text = data.get('response_text', '')
        
        if not table_data:
            return jsonify({"success": False, "error": "No data provided"}), 400

        # Generate a unique filename
        filename = f"VAVE_Detailed_Report_{uuid.uuid4().hex[:8]}.pptx"
        output_path = os.path.join(STATIC_DIR, "reports", filename)
        
        # Ensure directory exists
        os.makedirs(os.path.dirname(output_path), exist_ok=True)

        # Use the new comprehensive VAVE Presentation Engine
        # This uses Gemini API for LLM enrichment and generates boardroom-ready presentation
        logger.info(f"Starting PPT generation for {len(table_data)} ideas...")
        
        try:
            generate_deep_dive_ppt(table_data, output_path)
            logger.info(f"PPT generated successfully: {output_path}")
        except Exception as gen_error:
            logger.error(f"Error during PPT generation: {gen_error}")
            import traceback
            traceback.print_exc()
            return jsonify({
                "success": False, 
                "error": f"Failed to generate PPT: {str(gen_error)}"
            }), 500
        
        # Verify file was created
        if not os.path.exists(output_path):
            logger.error(f"PPT file was not created at {output_path}")
            return jsonify({
                "success": False, 
                "error": "PPT file generation completed but file not found"
            }), 500
        
        # Return the URL for download
        download_url = url_for('static', filename=f'reports/{filename}')
        
        return jsonify({
            "success": True, 
            "download_url": download_url,
            "filename": filename,
            "message": "PPT generated with deep-dive analysis using LLM enrichment."
        })

    except Exception as e:
        logger.error(f"PPT Generation Error: {e}")
        import traceback
        traceback.print_exc()
        return jsonify({"success": False, "error": f"PPT generation failed: {str(e)}"}), 500
    
# SQL without Endpoints
'''@app.route('/stats')
def stats():
    """Get application statistics."""
    try:
        with db_lock:
            with sqlite3.connect(DB_PATH) as conn:
                cursor = conn.cursor()
                
                # Count total ideas
                cursor.execute("SELECT COUNT(*) FROM ideas")
                total_ideas = cursor.fetchone()[0]
                
                # Count by status
                cursor.execute("SELECT status, COUNT(*) FROM ideas GROUP BY status")
                status_counts = dict(cursor.fetchall())
                
                # Count by department
                cursor.execute("SELECT dept, COUNT(*) FROM ideas GROUP BY dept")
                dept_counts = dict(cursor.fetchall())
                
                return jsonify({
                    "total_ideas": total_ideas,
                    "status_distribution": status_counts,
                    "department_distribution": dept_counts,
                    "vlm_initialized": vlm_initialized,
                    "vector_db_size": len(idea_texts)
                })
                
    except Exception as e:
        logger.error(f"Error getting stats: {str(e)}")
        return jsonify({"error": f"Failed to get stats: {str(e)}"}), 500'''
    

@app.route('/stats')
@login_required
def stats():
    """Get application statistics."""
    conn = None
    cursor = None
    try:
        conn = get_db_connection()
        cursor = conn.cursor(cursor_factory=DictCursor)
        
        cursor.execute("SELECT COUNT(*) FROM ideas")
        total_ideas = cursor.fetchone()[0]
        
        cursor.execute("SELECT status, COUNT(*) as count FROM ideas WHERE status IS NOT NULL GROUP BY status")
        status_counts = {row['status']: row['count'] for row in cursor.fetchall()}
        
        cursor.execute("SELECT dept, COUNT(*) as count FROM ideas WHERE dept IS NOT NULL AND dept != 'N/A' GROUP BY dept")
        dept_counts = {row['dept']: row['count'] for row in cursor.fetchall()}
        
        return jsonify({
            "total_ideas": total_ideas,
            "status_distribution": status_counts,
            "department_distribution": dept_counts,
            "vlm_initialized": vlm_initialized,
            "vector_db_size": len(idea_texts)
        })
            
    except Exception as e:
        logger.error(f"Error getting stats: {str(e)}")
        return jsonify({"error": f"Failed to get stats: {str(e)}"}), 500
    finally:
        if cursor:
            cursor.close()
        if conn:
            conn.close()


def analyze_query_for_agentic(query: str) -> dict:
    """
    Fine-detail query analysis for Agentic RAG.
    Extracts intent, keywords, filters, sub-phrases, and preserves full query for complex prompts.
    """
    q = (query or "").strip()
    q_lower = q.lower()
    tokens = q_lower.split()
    stop_words = {"the", "a", "an", "and", "or", "but", "in", "on", "at", "to", "for", "of", "with", "by", "is", "it", "that", "this"}

    # Question words
    question_words = ["what", "how", "why", "where", "when", "who", "which", "can", "should", "could", "give", "show", "list"]
    has_question_word = any(w in tokens for w in question_words)

    # Query intent (multi-intent aware)
    intents = []
    if any(w in q_lower for w in ["save", "saving", "cost", "reduce", "cut", "optimize"]):
        intents.append("cost_reduction")
    if any(w in q_lower for w in ["idea", "ideas", "suggestion", "proposal", "recommend"]):
        intents.append("idea_search")
    if any(w in q_lower for w in ["compare", "difference", "versus", "vs", "benchmark"]):
        intents.append("comparison")
    if any(w in q_lower for w in ["filter", "show", "list", "find", "get", "all"]):
        intents.append("filter_search")
    if any(w in q_lower for w in ["weight", "lightweight", "mass", "kg"]):
        intents.append("weight")
    intent = intents[0] if intents else "general"

    # Domain terms (expanded for fine detail)
    domain_terms = [
        "paint", "tire", "wheel", "brake", "caliper", "rotor", "suspension", "damper", "spring", "strut",
        "engine", "door", "hood", "fender", "panel", "frame", "bracket", "coating", "thickness", "weight",
        "material", "component", "assembly", "part", "chassis", "body", "interior", "trim", "seat",
        "hvac", "blower", "battery", "motor", "transmission", "exhaust", "catalytic"
    ]
    domain_keywords = [t for t in domain_terms if t in q_lower]

    # Sub-phrases: split on common conjunctions for complex prompts
    sub_phrases = re.split(r"\s+and\s+|\s+with\s+|\s+;\s+|\s+,\s+(?=[a-z])", q_lower)
    sub_phrases = [p.strip() for p in sub_phrases if len(p.strip()) > 3][:5]
    meaningful_tokens = [w for w in tokens if w not in stop_words and len(w) > 2]
    key_phrases = list(dict.fromkeys(meaningful_tokens[-6:] + meaningful_tokens[:4]))

    has_cost_filter = any(w in q_lower for w in [
        "more than", "greater than", "above", "at least", "less than", "below", "under", "minimum", "maximum", "inr", "lakh", "crore"
    ])

    return {
        "length": len(tokens),
        "is_short": len(tokens) <= 5,
        "is_complex": len(tokens) > 15,
        "full_query": q,
        "has_question_word": has_question_word,
        "intent": intent,
        "intents": intents,
        "domain_keywords": domain_keywords,
        "sub_phrases": sub_phrases,
        "key_phrases": key_phrases,
        "has_cost_filter": has_cost_filter,
    }


def rewrite_query_for_agentic(query: str, analysis: dict) -> str:
    """
    Enhanced query rewrite for better retrieval.
    Uses fine-detail analysis; for complex prompts preserves key sub-phrases and domain terms.
    """
    q = (query or "").strip()
    if not q:
        return q

    domain_keywords = analysis.get("domain_keywords", [])
    key_phrases = analysis.get("key_phrases", [])
    sub_phrases = analysis.get("sub_phrases", [])
    intent = analysis.get("intent", "general")
    is_complex = analysis.get("is_complex", False)

    # For complex prompts: build retrieval query from sub-phrases + domain + key phrases (no truncation)
    if is_complex and (sub_phrases or key_phrases):
        parts = []
        if sub_phrases:
            parts.append(" ".join(sub_phrases[:3]))
        if domain_keywords:
            parts.append(" ".join(domain_keywords[:3]))
        if key_phrases and not parts:
            parts.append(" ".join(key_phrases[:5]))
        if parts:
            base = " ".join(parts)
            if "cost" not in base and "saving" not in base:
                base = f"cost reduction {base}"
            return base[:MAX_QUERY_LENGTH_FOR_EMBEDDING]

    # For very short queries, expand with domain context
    if analysis.get("is_short"):
        if domain_keywords:
            return f"cost reduction ideas automotive {q} {' '.join(domain_keywords[:2])}"
        return f"cost reduction ideas automotive MG JSW MGI {q}"

    # For queries missing domain context, add it
    if not domain_keywords and intent == "idea_search":
        stop_words = {"the", "a", "an", "and", "or", "but", "in", "on", "at", "to", "for", "of", "with", "by"}
        main_words = [w for w in q.lower().split() if w not in stop_words][:5]
        return f"{q} automotive cost reduction ideas {' '.join(main_words)}"

    # For cost filter queries, ensure we include saving/cost terms
    if analysis.get("has_cost_filter") and "save" not in q.lower() and "saving" not in q.lower():
        return f"{q} cost savings saving value"

    # For domain-specific queries, ensure automotive context
    if domain_keywords and "automotive" not in q.lower() and "mg" not in q.lower():
        return f"{q} automotive MG Motor India"

    return q


def grade_retrieval(table_data) -> int:
    """
    Grade retrieval quality based on number of results and relevance.
    Returns a score indicating retrieval quality.
    """
    if not table_data:
        return 0
    
    count = len(table_data)
    
    # Bonus points for having multiple results (indicates good retrieval)
    if count >= 5:
        return count + 2
    elif count >= 3:
        return count + 1
    
    return count


'''def agentic_rag_chat(username: str, user_query: str, image_path=None) -> dict:
    """
    Enhanced Agentic RAG using VAVEAgent's Hybrid Logic Engine with Router Architecture.
    
    The router automatically selects the correct strategy:
    - SQL_FILTER: For filtering, math, lookup, duplicate detection (Category 2 & 4)
    - ENGINEERING_LOGIC: For generation, creative engineering (Category 1)
    - VECTOR_SEARCH: Default for semantic search (Category 3)
    """
    # Image path -> keep existing VLM flow as-is
    if image_path:
        logger.info("Agentic RAG: image provided, delegating directly to generate_response (VLM path).")
        return generate_response(username, user_query, image_path=image_path)

    # Use VAVEAgent's router architecture if available
    global vave_agent
    if vave_agent is not None:
        try:
            logger.info(f"VAVE Agent: Processing query with router architecture: {user_query[:100]}")
            agent_response = vave_agent.run(user_query)
            
            # Get raw data from agent
            table_data = vave_agent.get_last_result_data()
            response_text = agent_response
            
            # Enforce images for all origins (re-run VLM if missing)
            if table_data and vave_agent.vlm:
                for idea in table_data:
                    if idea.get("Current Scenario Image") in ["N/A", "NaN", None, ""]:
                        # Force re-retrieval
                        idea_text = idea.get("Cost Reduction Idea", "")
                        origin = idea.get("Origin", "Existing Database")
                        images = vave_agent.vlm.get_images_for_idea(
                            idea_text, origin, {"idea_text": idea_text}
                        )
                        idea.update({
                            "Current Scenario Image": images.get('current_scenario_image', 'static/images/mg/hector_rear.jpg'),
                            "Proposal Scenario Image": images.get('proposal_scenario_image', 'static/defaults/proposal_placeholder.jpg'),
                            "Competitor Image": images.get('competitor_image', 'static/images/competitor/1.jpg'),
                            "current_scenario_image": images.get('current_scenario_image', 'static/images/mg/hector_rear.jpg'),
                            "proposal_scenario_image": images.get('proposal_scenario_image', 'static/defaults/proposal_placeholder.jpg'),
                            "competitor_image": images.get('competitor_image', 'static/images/competitor/1.jpg')
                        })
            
            # Convert agent data format to expected table format
            if table_data and isinstance(table_data, list):
                # Convert agent's dict format to table format expected by frontend
                formatted_table_data = []
                for row in table_data:
                    formatted_row = {}
                    # Map agent's column names to frontend expected names with safe conversions
                    idea_id_val = row.get('idea_id', row.get('Idea Id', 'N/A'))
                    saving_val = safe_float_convert(row.get('saving_value_inr', row.get('Saving Value (INR)', 0)))
                    weight_val = safe_float_convert(row.get('weight_saving', row.get('Weight Saving (Kg)', 0)))

                    formatted_row['Idea Id'] = str(idea_id_val)
                    formatted_row['Cost Reduction Idea'] = str(row.get('cost_reduction_idea', row.get('Cost Reduction Idea', 'N/A')))
                    formatted_row['Saving Value (INR)'] = f"INR {saving_val:,.2f}"
                    formatted_row['Weight Saving (Kg)'] = f"{weight_val:.2f}"
                    capex_val = row.get('capex', row.get('CAPEX', row.get('investment_cr', row.get('Investment (Cr)'))))
                    formatted_row['CAPEX'] = f"{safe_float_convert(capex_val):,.2f}" if capex_val is not None and str(capex_val).strip() else "N/A"
                    formatted_row['Status'] = str(row.get('status', row.get('Status', 'N/A')))
                    formatted_row['Dept'] = str(row.get('dept', row.get('Dept', 'N/A')))
                    formatted_row['Way Forward'] = str(row.get('way_forward', row.get('Way Forward', 'N/A')))

                    # Fields needed for PPT/chart generation
                    formatted_row['Idea ID'] = str(idea_id_val)
                    formatted_row['Proposal'] = formatted_row['Cost Reduction Idea']
                    formatted_row['Imp'] = formatted_row['Way Forward']
                    formatted_row['Saving Value'] = f"INR {saving_val:,.2f}"
                    formatted_row['Responsibility'] = str(row.get('resp', row.get('Responsibility', 'N/A')))
                    formatted_row['Date'] = str(row.get('target_date', row.get('Date', 'N/A')))

                    # Add any other fields
                    for key, value in row.items():
                        if key not in [
                            'idea_id', 'cost_reduction_idea', 'saving_value_inr', 'weight_saving',
                            'status', 'dept', 'way_forward', 'resp', 'target_date'
                        ]:
                            formatted_row[key.replace('_', ' ').title()] = str(value) if value is not None else 'N/A'
                    formatted_table_data.append(formatted_row)
                table_data = formatted_table_data
            
            # If no table data from agent, try to get it from vector search as fallback
            if not table_data:
                table_data, context_str = retrieve_context(user_query)
            
            # Files are now generated on-demand when download buttons are clicked
            
            result = {
                "success": True,
                "response_text": response_text,
                "table_data": table_data,
                "image_urls": [],
                "excel_url": "",
                "ppt_url": "",
                "agentic_metadata": {
                    "strategy": "VAVE_AGENT_ROUTER",
                    "original_query": user_query
                }
            }
            
            logger.info(f"VAVE Agent: Response generated successfully with {len(table_data)} rows")
            return result
            
        except Exception as agent_error:
            logger.error(f"VAVE Agent error: {agent_error}. Falling back to legacy agentic_rag_chat.")
            # Fall through to legacy implementation
    
    # Legacy Agentic RAG implementation (fallback)
    analysis = analyze_query_for_agentic(user_query)
    logger.info(f"Agentic RAG: query analysis = {analysis}")

    # First retrieval with original query
    initial_table, _ = retrieve_context(user_query)
    initial_score = grade_retrieval(initial_table)
    logger.info(f"Agentic RAG: initial retrieval rows = {len(initial_table)}, score = {initial_score}")

    chosen_query = user_query
    chosen_table = initial_table

    # If poor retrieval (score < 3), try query rewrite
    if initial_score < 3:
        rewritten = rewrite_query_for_agentic(user_query, analysis)
        if rewritten != user_query:
            logger.info(f"Agentic RAG: attempting rewritten query: {rewritten!r}")
            rewritten_table, _ = retrieve_context(rewritten)
            rewritten_score = grade_retrieval(rewritten_table)
            logger.info(f"Agentic RAG: rewritten retrieval rows = {len(rewritten_table)}, score = {rewritten_score}")

            if rewritten_score > initial_score:
                chosen_query = rewritten
                chosen_table = rewritten_table
                logger.info("Agentic RAG: using rewritten query for generation.")
            else:
                logger.info("Agentic RAG: keeping original query for generation.")
    
    # If still poor results and we have domain keywords, try domain-focused query
    if len(chosen_table) < 3 and analysis.get("domain_keywords"):
        domain_query = f"{user_query} {' '.join(analysis['domain_keywords'][:2])}"
        logger.info(f"Agentic RAG: attempting domain-focused query: {domain_query!r}")
        domain_table, _ = retrieve_context(domain_query)
        domain_score = grade_retrieval(domain_table)
        logger.info(f"Agentic RAG: domain-focused retrieval rows = {len(domain_table)}, score = {domain_score}")
        
        if domain_score > grade_retrieval(chosen_table):
            chosen_query = domain_query
            chosen_table = domain_table
            logger.info("Agentic RAG: using domain-focused query for generation.")

    # Generation step – use pre-retrieved data for better results
    # Retrieve final context with chosen query for optimal results
    if chosen_table and len(chosen_table) > 0:
        final_table, final_context = retrieve_context(chosen_query)
        if final_table and len(final_table) > 0:
            # Use the retrieved data directly in generate_response
            result = generate_response(username, user_query, image_path=None, pre_retrieved_data=(final_table, final_context))
        else:
            result = generate_response(username, user_query, image_path=None)
    else:
        result = generate_response(username, user_query, image_path=None)
    
    # Annotate metadata for analytics
    result.setdefault("agentic_metadata", {})
    result["agentic_metadata"]["original_query"] = user_query
    result["agentic_metadata"]["final_query"] = chosen_query
    result["agentic_metadata"]["retrieval_score"] = grade_retrieval(result.get("table_data", []))
    result["agentic_metadata"]["query_analysis"] = analysis

    return result'''
def convert_image_paths_to_urls(table_data: list) -> list:
    """
    Convert image paths in table data to proper URLs for frontend display.
    
    Args:
        table_data: List of idea dictionaries with image paths
        
    Returns:
        List of idea dictionaries with image URLs
    """
    if not table_data:
        return table_data
    
    processed_data = []
    for row in table_data:
        processed_row = row.copy()
        
        # Process Current Scenario Image
        current_img = processed_row.get('current_scenario_image') or processed_row.get('Current Scenario Image')
        if current_img and current_img != "N/A" and str(current_img).strip():
            current_img = str(current_img).strip()
            if os.path.exists(current_img):
                # Absolute path - copy to static and convert to URL
                filename = os.path.basename(current_img)
                static_path = STATIC_DIR / "generated" / filename
                if not static_path.exists():
                    try:
                        shutil.copy(current_img, static_path)
                    except Exception as e:
                        logger.error(f"Failed to copy current image: {e}")
                processed_row['current_scenario_image'] = url_for('static', filename=f'generated/{filename}')
                processed_row['Current Scenario Image'] = url_for('static', filename=f'generated/{filename}')
            elif current_img.startswith("static/") or current_img.startswith("/static/"):
                # Already a static path - convert to URL
                url_path = current_img.replace("static/", "").replace("/static/", "").lstrip("/")
                processed_row['current_scenario_image'] = url_for('static', filename=url_path)
                processed_row['Current Scenario Image'] = url_for('static', filename=url_path)
            elif current_img.startswith("http") or current_img.startswith("/"):
                # Already a URL - keep as is
                processed_row['current_scenario_image'] = current_img
                processed_row['Current Scenario Image'] = current_img
            else:
                # Invalid path - try to check if it's a relative path in generated folder
                if not os.path.isabs(current_img):
                    # Try as relative path in generated folder
                    potential_path = STATIC_DIR / "generated" / current_img
                    if potential_path.exists():
                        processed_row['current_scenario_image'] = url_for('static', filename=f'generated/{current_img}')
                        processed_row['Current Scenario Image'] = url_for('static', filename=f'generated/{current_img}')
                    else:
                        processed_row['current_scenario_image'] = "N/A"
                        processed_row['Current Scenario Image'] = "N/A"
                else:
                    processed_row['current_scenario_image'] = "N/A"
                    processed_row['Current Scenario Image'] = "N/A"
        else:
            processed_row['current_scenario_image'] = "N/A"
            processed_row['Current Scenario Image'] = "N/A"
        
        # Process Proposal Scenario Image
        proposal_img = processed_row.get('proposal_scenario_image') or processed_row.get('Proposal Scenario Image')
        if proposal_img and proposal_img != "N/A" and str(proposal_img).strip():
            proposal_img = str(proposal_img).strip()
            if os.path.exists(proposal_img):
                # Absolute path - copy to static and convert to URL
                filename = os.path.basename(proposal_img)
                static_path = STATIC_DIR / "generated" / filename
                if not static_path.exists():
                    try:
                        shutil.copy(proposal_img, static_path)
                    except Exception as e:
                        logger.error(f"Failed to copy proposal image: {e}")
                processed_row['proposal_scenario_image'] = url_for('static', filename=f'generated/{filename}')
                processed_row['Proposal Scenario Image'] = url_for('static', filename=f'generated/{filename}')
            elif proposal_img.startswith("static/") or proposal_img.startswith("/static/"):
                # Already a static path - convert to URL
                url_path = proposal_img.replace("static/", "").replace("/static/", "").lstrip("/")
                processed_row['proposal_scenario_image'] = url_for('static', filename=url_path)
                processed_row['Proposal Scenario Image'] = url_for('static', filename=url_path)
            elif proposal_img.startswith("http") or proposal_img.startswith("/"):
                # Already a URL - keep as is
                processed_row['proposal_scenario_image'] = proposal_img
                processed_row['Proposal Scenario Image'] = proposal_img
            else:
                # Invalid path - try to check if it's a relative path in generated folder
                if not os.path.isabs(proposal_img):
                    # Try as relative path in generated folder
                    potential_path = STATIC_DIR / "generated" / proposal_img
                    if potential_path.exists():
                        processed_row['proposal_scenario_image'] = url_for('static', filename=f'generated/{proposal_img}')
                        processed_row['Proposal Scenario Image'] = url_for('static', filename=f'generated/{proposal_img}')
                    else:
                        processed_row['proposal_scenario_image'] = "N/A"
                        processed_row['Proposal Scenario Image'] = "N/A"
                else:
                    processed_row['proposal_scenario_image'] = "N/A"
                    processed_row['Proposal Scenario Image'] = "N/A"
        else:
            processed_row['proposal_scenario_image'] = "N/A"
            processed_row['Proposal Scenario Image'] = "N/A"
        
        # Process Competitor Image
        competitor_img = processed_row.get('competitor_image') or processed_row.get('Competitor Image')
        if competitor_img and competitor_img != "N/A" and str(competitor_img).strip():
            competitor_img = str(competitor_img).strip()
            if os.path.exists(competitor_img):
                # Absolute path - copy to static and convert to URL
                filename = os.path.basename(competitor_img)
                static_path = STATIC_DIR / "generated" / filename
                if not static_path.exists():
                    try:
                        shutil.copy(competitor_img, static_path)
                    except Exception as e:
                        logger.error(f"Failed to copy competitor image: {e}")
                processed_row['competitor_image'] = url_for('static', filename=f'generated/{filename}')
                processed_row['Competitor Image'] = url_for('static', filename=f'generated/{filename}')
            elif competitor_img.startswith("static/") or competitor_img.startswith("/static/"):
                # Already a static path - convert to URL
                url_path = competitor_img.replace("static/", "").replace("/static/", "").lstrip("/")
                processed_row['competitor_image'] = url_for('static', filename=url_path)
                processed_row['Competitor Image'] = url_for('static', filename=url_path)
            elif competitor_img.startswith("http") or competitor_img.startswith("/"):
                # Already a URL - keep as is
                processed_row['competitor_image'] = competitor_img
                processed_row['Competitor Image'] = competitor_img
            else:
                # Invalid path - try to check if it's a relative path in generated folder
                if not os.path.isabs(competitor_img):
                    # Try as relative path in generated folder
                    potential_path = STATIC_DIR / "generated" / competitor_img
                    if potential_path.exists():
                        processed_row['competitor_image'] = url_for('static', filename=f'generated/{competitor_img}')
                        processed_row['Competitor Image'] = url_for('static', filename=f'generated/{competitor_img}')
                    else:
                        processed_row['competitor_image'] = "N/A"
                        processed_row['Competitor Image'] = "N/A"
                else:
                    processed_row['competitor_image'] = "N/A"
                    processed_row['Competitor Image'] = "N/A"
        else:
            processed_row['competitor_image'] = "N/A"
            processed_row['Competitor Image'] = "N/A"
        
        processed_data.append(processed_row)
    
    return processed_data

def agentic_rag_chat(username: str, user_query: str, image_path=None) -> dict:
    """
    Wrapper to call the VAVE Agent (Logic Engine).
    """
    global vave_agent

    # 1. Initialize Agent (Lazy Loading)
    if vave_agent is None:
        try:
            # We pass the necessary functions and paths to the agent
            vave_agent = VAVEAgent(
                db_path=str(DB_PATH), 
                vector_db_func=retrieve_context, 
                pg_conn_func=get_db_connection,
                db_conn=get_db_connection,  # Pass DB connection function for VLM Engine
                faiss_index=faiss_index,  # Pass FAISS index for VLM Engine
                sentence_model=embedding_model  # Pass sentence model for VLM Engine
            )
            logger.info("VAVE Agent initialized successfully.")
        except Exception as e:
            logger.error(f"Failed to init Agent: {e}")
            # Fallback to standard Legacy Chat if Agent dies
            return generate_response(username, user_query)

    # 2. Run the Reasoning Engine (full user_query preserved for complex prompts)
    try:
        response_text = vave_agent.run(user_query)
        table_data = vave_agent.get_last_result_data()
        
        # Enforce images for all origins (re-run VLM if missing)
        if table_data and vave_agent.vlm:
            for idea in table_data:
                if idea.get("Current Scenario Image") in ["N/A", "NaN", None, ""]:
                    # Force re-retrieval
                    idea_text = idea.get("Cost Reduction Idea", "")
                    origin = idea.get("Origin", "Existing Database")
                    images = vave_agent.vlm.get_images_for_idea(
                        idea_text, origin, {"idea_text": idea_text}
                    )
                    idea.update({
                        "Current Scenario Image": images.get('current_scenario_image', 'static/images/mg/hector_rear.jpg'),
                        "Proposal Scenario Image": images.get('proposal_scenario_image', 'static/defaults/proposal_placeholder.jpg'),
                        "Competitor Image": images.get('competitor_image', 'static/images/competitor/1.jpg'),
                        "current_scenario_image": images.get('current_scenario_image', 'static/images/mg/hector_rear.jpg'),
                        "proposal_scenario_image": images.get('proposal_scenario_image', 'static/defaults/proposal_placeholder.jpg'),
                        "competitor_image": images.get('competitor_image', 'static/images/competitor/1.jpg')
                    })
        
        # 3. Convert image paths to URLs for frontend
        table_data = convert_image_paths_to_urls(table_data)
        
        # 4. Files are now generated on-demand when download buttons are clicked
        # This saves API service costs by not generating files automatically

        # 5. Return to Frontend
        return {
            "success": True,
            "response_text": response_text,
            "table_data": table_data,
            "image_urls": [], 
            "excel_url": "",  # Will be generated on-demand
            "ppt_url": "",  # Will be generated on-demand
        }

    except Exception as e:
        logger.error(f"Agent Logic Failure (complex prompt or runtime error): {e}", exc_info=True)
        # Robust Fallback: use full user_query so analysis/retrieval still runs in fine detail
        return generate_response(username, user_query, image_path=image_path) if user_query else generate_response(username, "", image_path=image_path)

'''@app.route('/history')
def history():
    """Get recent query history."""
    try:
        with db_lock:
            with sqlite3.connect(DB_PATH) as conn:
                cursor = conn.cursor()
                cursor.execute("""
                    SELECT user_input, response_text, created_at 
                    FROM ideas 
                    WHERE user_input IS NOT NULL 
                    ORDER BY created_at DESC 
                    LIMIT 10
                """)
                history_data = cursor.fetchall()
                
                return jsonify({
                    "history": [
                        {
                            "query": row[0],
                            "response": row[1][:200] + "..." if len(row[1]) > 200 else row[1],
                            "timestamp": row[2]
                        }
                        for row in history_data
                    ]
                })
                
    except Exception as e:
        logger.error(f"Error getting history: {str(e)}")
        return jsonify({"error": f"Failed to get history: {str(e)}"}), 500'''

@app.route('/history')
@login_required # <-- MODIFIED & PROTECTED
def history():
    """Get RECENT chat history FOR THE CURRENT USER."""
    conn = None
    cursor = None
    try:
        conn = get_db_connection()
        cursor = conn.cursor(cursor_factory=DictCursor)
        
        # --- MODIFIED: Query chat_history table for the current user ---
        cursor.execute("""
            SELECT user_query, response_text, created_at 
            FROM chat_history 
            WHERE username = %s
            ORDER BY created_at DESC 
            LIMIT 10
        """, (current_user.username,)) # <-- Filter by user
        
        history_data = cursor.fetchall()
            
        return jsonify({
            "history": [
                {
                    "query": row['user_query'],
                    "response": row['response_text'][:200] + "..." if row['response_text'] and len(row['response_text']) > 200 else row['response_text'],
                    "timestamp": row['created_at']
                }
                for row in history_data
            ]
        })
            
    except Exception as e:
        logger.error(f"Error getting history: {str(e)}")
        return jsonify({"error": f"Failed to get history: {str(e)}"}), 500
    finally:
        if cursor:
            cursor.close()
        if conn:
            conn.close()

def cleanup_temp_files():
    """Clean up old temporary files."""
    try:
        import time
        current_time = time.time()
        
        # Only clean up files in TEMP_DIR, not UPLOAD_FOLDER
        for file_path in TEMP_DIR.glob("*"):
            if file_path.is_file() and file_path.parent == TEMP_DIR:
                # Delete files older than 1 hour
                if current_time - file_path.stat().st_mtime > 3600:
                    file_path.unlink()
                    logger.info(f"Cleaned up old temp file: {file_path}")
                    
    except Exception as e:
        logger.error(f"Error cleaning up temp files: {str(e)}")

@app.route('/reset_password/<int:user_id>')
@login_required
def reset_password(user_id):
    """
    Legacy route: resets password via redirect/flash, used by older UI.
    """
    if current_user.role != 'SuperAdmin':
        flash("Permission denied", "error")
        return redirect(url_for('chat_app', section='dashboard'))

    import string
    import random
    chars = string.ascii_letters + string.digits
    new_pass = ''.join(random.choice(chars) for _ in range(8))
    hashed_pw = generate_password_hash(new_pass)
    
    conn = get_user_db_conn()
    conn.execute("UPDATE users SET password_hash = ? WHERE id = ?", (hashed_pw, user_id))
    conn.commit()
    conn.close()
    
    flash(f'Password reset! New password for user ID {user_id} is: {new_pass}', 'success')
    return redirect(url_for('chat_app', section='dashboard'))


@app.route('/reset_password', methods=['POST'])
@login_required
def reset_password_api():
    """
    JSON API: reset password by username, used by SPA.
    """
    if current_user.role != 'SuperAdmin':
        return jsonify({"success": False, "error": "Permission denied"}), 403

    data = request.get_json(silent=True) or {}
    username = data.get('username')
    if not username:
        return jsonify({"success": False, "error": "Username is required"}), 400

    import string
    import random
    chars = string.ascii_letters + string.digits
    new_pass = ''.join(random.choice(chars) for _ in range(8))
    hashed_pw = generate_password_hash(new_pass)

    conn = get_user_db_conn()
    cursor = conn.cursor()
    try:
        cursor.execute("UPDATE users SET password_hash = ? WHERE username = ?", (hashed_pw, username))
        if cursor.rowcount == 0:
            return jsonify({"success": False, "error": "User not found"}), 404
        conn.commit()
        return jsonify({"success": True, "new_password": new_pass})
    finally:
        conn.close()

# Initialize the application
if __name__ == '__main__':
    try:
        # 0. Ensure users.db schema (roles, SuperAdmin) is correct
        init_user_db_schema()

        # 1. Ensure database connection and tables exist
        init_db()
        # 2. Load GPT-2 
        setup_model()
        # 3. Load text data from DB for FAISS index
        build_vector_db()
        # 4. Load VLM model and image data from DB
        setup_vlm()
        
        # Clean up old temp files on startup
        cleanup_temp_files()
        
        logger.info("Application initialization complete. Starting server...")
        app.run(host='0.0.0.0', port=5000, debug=False)
    except Exception as e:
        logger.error(f"Failed to start application: {str(e)}")
        raise

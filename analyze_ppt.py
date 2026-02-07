import os
import sys

try:
    from pptx import Presentation
except ImportError:
    print("python-pptx not installed. Skipping PPT analysis.")
    sys.exit(0)

ppt_path = "VAVE_Detailed_Report_a23ed2c8.pptx"

try:
    prs = Presentation(ppt_path)
    print(f"Analyzing PPT: {ppt_path}")
    print(f"Total Slides: {len(prs.slides)}")
    
    for i, slide in enumerate(prs.slides[:10]): # Check first 10 slides
        print(f"\nSlide {i+1}:")
        # Extract text
        text = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text.replace("\n", " "))
        print(f"  Text: {' | '.join(text)[:100]}...")
        
        # Check for images
        images = []
        for shape in slide.shapes:
            if shape.shape_type == 13: # PICTURE
                images.append("Image Found")
        if images:
            print(f"  Images: {len(images)} found")
        else:
            print("  Images: None")
            
except Exception as e:
    print(f"Error reading PPT: {e}")

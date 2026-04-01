import os
import json
import logging
import re
from datetime import date
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.chart.data import CategoryChartData
import google.generativeai as genai

# --- CONFIGURATION ---
logger = logging.getLogger(__name__)

# Slide bounds (16:9 widescreen) - all content must stay within these
SLIDE_WIDTH_IN = 13.333
SLIDE_HEIGHT_IN = 7.5
SLIDE_MARGIN_IN = 0.2
CONTENT_LEFT_IN = 0.3
CONTENT_RIGHT_IN = SLIDE_WIDTH_IN - SLIDE_MARGIN_IN
CONTENT_TOP_IN = 1.0
CONTENT_BOTTOM_IN = SLIDE_HEIGHT_IN - SLIDE_MARGIN_IN

# Ensure API Key is loaded
if "GOOGLE_API_KEY" in os.environ:
    genai.configure(api_key=os.environ["GOOGLE_API_KEY"])

def clean_text(text):
    """Removes special characters that break PPT generation."""
    if not text: return ""
    return str(text).encode('ascii', 'ignore').decode('ascii').strip()

def truncate_text(text, max_chars=500, suffix="..."):
    """Truncate text so it fits within slide elements and does not overflow."""
    if not text or len(str(text).strip()) <= max_chars:
        return (text or "").strip()
    s = str(text).strip()
    return s[: max_chars - len(suffix)].rstrip() + suffix

def expand_idea_with_llm(simple_idea_dict):
    """
    Calls Gemini 1.5 Pro to act as a Senior Engineer and generate deep content.
    """
    idea_title = simple_idea_dict.get('Cost Reduction Idea Proposal', 'Cost Reduction Idea')
    current_mg = simple_idea_dict.get('MG Product Scenario', 'Standard MG Spec')
    competitor = simple_idea_dict.get('Competitor Product Scenario', 'Standard Competitor Spec')
    
    # STRICT prompt for dense engineering content
    prompt = f"""
    Act as a Chief Value Engineer at MG Motors. Analyze this Cost Reduction Idea:
    
    IDEA: "{idea_title}"
    CONTEXT: MG Current: "{current_mg}" | Competitor: "{competitor}"
    
    Generate a JSON object with 4 detailed sections. 
    Write professionally. Use numbers, material grades (e.g. PA66, HSS), and specific processes.
    
    JSON FORMAT:
    {{
        "proposal_text": "Technical description (40-50 words). Focus on HOW. Mention materials/process.",
        "benchmarking_text": "Comparison analysis (40-50 words). Contrast MG vs Competitor explicitly.",
        "financial_text": "Cost breakdown (Material, Process, Tooling). Justify the saving.",
        "feasibility_text": "Validation plan. Mention specific tests (e.g., vibration, heat aging, crash)."
    }}
    """
    
    try:
        model = genai.GenerativeModel('gemini-3.1-flash-lite')
        response = model.generate_content(prompt)
        content = response.text.replace("```json", "").replace("```", "").strip()
        return json.loads(content)
    except Exception as e:
        logger.error(f"LLM Generation Failed: {e}")
        # Fallback content if AI fails
        return {
            "proposal_text": f"Engineering Proposal: {idea_title}. Change implemented to optimize BOM cost.",
            "benchmarking_text": f"Competitor Reference: {competitor}. Gap analysis indicates potential for VAVE.",
            "financial_text": "Estimated savings based on raw material delta. Tooling amortization TBD.",
            "feasibility_text": "Standard validation (DV/PV) required. Check homologation impacts."
        }

def set_widescreen(prs):
    """Forces 16:9 Layout for detailed content."""
    prs.slide_width = Inches(SLIDE_WIDTH_IN)
    prs.slide_height = Inches(SLIDE_HEIGHT_IN)

def format_textbox(shape, text, font_size=11, bold_title=False, max_text_chars=500):
    """Format text inside a box with word wrap and fit-to-shape so nothing overflows."""
    tf = shape.text_frame
    tf.clear()
    tf.margin_top = Pt(5)
    tf.margin_bottom = Pt(5)
    tf.margin_left = Pt(10)
    tf.margin_right = Pt(10)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.word_wrap = True

    raw = clean_text(truncate_text(text, max_chars=max_text_chars))
    if not raw:
        raw = "-"
    p = tf.paragraphs[0]
    p.text = raw
    p.font.size = Pt(font_size)
    p.font.name = "Calibri"
    p.font.color.rgb = RGBColor(0, 0, 0)
    if bold_title:
        p.font.bold = True
        p.font.size = Pt(min(font_size + 2, 14))
    try:
        tf.fit_text(font_family="Calibri", max_size=font_size, bold=bold_title)
    except Exception:
        pass

def generate_deep_dive_ppt(ideas_list, output_path):
    """Main Generator Function."""
    prs = Presentation()
    set_widescreen(prs) # CRITICAL: Switch to 16:9
    
    # --- SLIDE 1: Title (within slide) ---
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    title.text = truncate_text("VAVE Cost Reduction: Deep Dive Report", max_chars=80)
    tf_title = title.text_frame
    tf_title.word_wrap = True
    tf_title.paragraphs[0].font.size = Pt(44)
    tf_title.paragraphs[0].font.bold = True

    subtitle = slide.placeholders[1]
    subtitle.text = truncate_text(
        f"Automated Engineering Analysis | {date.today()}\nDetailed Feasibility & Benchmarking",
        max_chars=200,
    )
    if hasattr(subtitle, "text_frame") and subtitle.text_frame:
        subtitle.text_frame.word_wrap = True
    
    # --- SLIDE 2: Dashboard ---
    create_dashboard_slide(prs, ideas_list)
    # --- SLIDE 3: Summary of ideas (table) ---
    create_summary_table_slide(prs, ideas_list)
    # --- SLIDE 4+: Deep Dives ---
    for i, idea in enumerate(ideas_list):
        # 1. AI Expansion
        print(f"Generating AI Analysis for Idea {i+1}/{len(ideas_list)}...")
        ai_data = expand_idea_with_llm(idea)
        full_data = {**idea, **ai_data}
        
        # 2. Create Slide
        create_technical_slide(prs, full_data)
        
    prs.save(output_path)
    return output_path

def _parse_saving_float(idea, default=0.0):
    """Parse saving value from idea dict to float for calculations."""
    val = idea.get("Saving Value(INR)") or idea.get("Saving Value (INR)") or idea.get("saving_value_inr") or idea.get("Saving Value") or "0"
    s = str(val).replace(",", "").replace("INR", "").replace(" ", "").strip()
    try:
        return float(re.sub(r"[^\d.-]", "", s) or 0)
    except Exception:
        return default


def create_summary_table_slide(prs, ideas_list):
    """
    One slide: Summary of ideas generated.
    Table columns: Idea ID, Idea, CAPEX (detailed/fine estimation: tooling, inventory, setup cost), Saving, Yearly volume (1 lakh default, editable),
    Saving overall (saving * volume), Status, Turnaround time, Required volume yearly.
    All content fits on slide with word wrap; no truncation or ellipsis.
    """
    if not ideas_list:
        return
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tbl_left = Inches(CONTENT_LEFT_IN)
    tbl_top = Inches(0.45)
    tbl_width = Inches(SLIDE_WIDTH_IN - 2 * CONTENT_LEFT_IN)
    tbl_height = Inches(SLIDE_HEIGHT_IN - 0.45 - SLIDE_MARGIN_IN)
    cols = 9
    rows = len(ideas_list) + 1
    table_shape = slide.shapes.add_table(rows, cols, tbl_left, tbl_top, tbl_width, tbl_height)
    table = table_shape.table
    header_row = [
        "Idea ID",
        "Idea",
        "CAPEX",
        "Saving",
        "Yearly volume (per 1 lakh)",
        "Saving overall (Saving × Volume)",
        "Status (Accept/Reject)",
        "Turnaround time",
        "Required volume yearly",
    ]
    font_header = 9
    font_cell = 8
    for c, label in enumerate(header_row):
        cell = table.cell(0, c)
        tf = cell.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].text = label
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.size = Pt(font_header)
        tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0, 51, 102)
    default_volume = 100000
    for r, idea in enumerate(ideas_list):
        row_idx = r + 1
        idea_id = str(idea.get("Idea Id") or idea.get("idea_id") or "")
        idea_text = str(idea.get("Cost Reduction Idea Proposal") or idea.get("Cost Reduction Idea") or idea.get("cost_reduction_idea") or "")
        capex = str(idea.get("CAPEX") or idea.get("Capex Investment") or idea.get("Investment (Cr)") or idea.get("capex") or idea.get("investment_cr") or "")
        saving_raw = idea.get("Saving Value(INR)") or idea.get("Saving Value (INR)") or idea.get("saving_value_inr") or idea.get("Saving Value") or "0"
        saving_str = str(saving_raw).strip()
        saving_float = _parse_saving_float(idea)
        yearly_vol = default_volume
        vol_raw = idea.get("volume") or idea.get("Yearly volume") or idea.get("yearly_volume")
        if vol_raw is not None and str(vol_raw).strip():
            try:
                yearly_vol = int(float(re.sub(r"[^\d.]", "", str(vol_raw))) or default_volume)
            except Exception:
                yearly_vol = default_volume
        saving_overall = saving_float * yearly_vol
        status = str(idea.get("Status") or idea.get("status") or idea.get("current_status") or "").strip()
        turnaround = str(idea.get("Turnaround time") or idea.get("turnaround_time") or "").strip()
        req_vol_yearly = str(idea.get("Required volume yearly") or "").strip() or str(yearly_vol)
        row_data = [
            idea_id,
            idea_text,
            capex,
            saving_str,
            str(yearly_vol),
            f"{saving_overall:,.0f}",
            status if status else "",
            turnaround,
            req_vol_yearly,
        ]
        for c, val in enumerate(row_data):
            cell = table.cell(row_idx, c)
            tf = cell.text_frame
            tf.word_wrap = True
            tf.margin_left = tf.margin_right = Pt(4)
            tf.margin_top = tf.margin_bottom = Pt(2)
            p = tf.paragraphs[0]
            p.text = (val or "").strip()
            p.font.size = Pt(font_cell)
            p.font.name = "Calibri"
            p.font.color.rgb = RGBColor(0, 0, 0)
    title_box = slide.shapes.add_textbox(tbl_left, Inches(0.05), tbl_width, Inches(0.3))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Summary of ideas generated"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)


def _abbrev_number(n):
    """Abbreviate large numbers so they fit in KPI boxes (e.g. 120000 -> 1.2L)."""
    if n >= 1_00_00_000:
        return f"{n / 1_00_00_000:.1f}Cr"
    if n >= 1_00_000:
        return f"{n / 1_00_000:.1f}L"
    if n >= 1_000:
        return f"{n / 1_000:.1f}K"
    return f"{n:,.0f}"

def create_dashboard_slide(prs, ideas):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    w, h = Inches(SLIDE_WIDTH_IN), Inches(SLIDE_HEIGHT_IN)

    # Header Bar (within slide)
    bar_h = 1.0
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(SLIDE_WIDTH_IN), Inches(bar_h))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(192, 0, 0)
    bar.line.fill.background()
    tf_bar = bar.text_frame
    tf_bar.word_wrap = True
    p = tf_bar.paragraphs[0]
    p.text = "EXECUTIVE SUMMARY"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    total = len(ideas)
    savings = sum([float(str(x.get('Saving Value(INR)', '0')).replace(',', '').replace('INR', '').strip() or 0) for x in ideas])
    avg = (savings / total) if total else 0

    # KPI Cards - stay within slide width
    kpi_w, kpi_h = 3.5, 1.8
    kpi_y = 1.5
    gap = (SLIDE_WIDTH_IN - 3 * kpi_w - 2 * 0.5) / 2  # margins and gaps
    kpi_x1 = SLIDE_MARGIN_IN
    kpi_x2 = kpi_x1 + kpi_w + gap
    kpi_x3 = kpi_x2 + kpi_w + gap

    def add_kpi(label, val, x_in):
        box = slide.shapes.add_shape(1, Inches(x_in), Inches(kpi_y), Inches(kpi_w), Inches(kpi_h))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(245, 245, 245)
        box.line.color.rgb = RGBColor(200, 200, 200)
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = Pt(8)
        p1 = tf.paragraphs[0]
        p1.text = truncate_text(label.upper(), max_chars=25)
        p1.font.size = Pt(11)
        p1.font.bold = True
        p1.font.color.rgb = RGBColor(100, 100, 100)
        p1.alignment = PP_ALIGN.CENTER
        p2 = tf.add_paragraph()
        p2.text = truncate_text(val, max_chars=20)
        p2.font.size = Pt(24)
        p2.font.bold = True
        p2.font.color.rgb = RGBColor(0, 51, 102)
        p2.alignment = PP_ALIGN.CENTER

    add_kpi("Total Ideas", str(total), kpi_x1)
    add_kpi("Total Savings (INR)", _abbrev_number(savings), kpi_x2)
    add_kpi("Avg per Idea", _abbrev_number(avg), kpi_x3)

def create_technical_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    w_in, h_in = SLIDE_WIDTH_IN, SLIDE_HEIGHT_IN

    # --- 1. Header (full width, full title - no truncation; wraps to show full topic) ---
    header_h = 1.15
    header = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(w_in), Inches(header_h))
    header.fill.solid()
    header.fill.fore_color.rgb = RGBColor(0, 51, 102)
    header.line.fill.background()
    id_txt = data.get('Idea Id', 'ID-XX')
    title_txt = data.get('Cost Reduction Idea Proposal', 'Idea Title') or data.get('Cost Reduction Idea', 'Idea Title')
    header_line = f" {id_txt} | {clean_text(title_txt)}"
    tf = header.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(12)
    tf.margin_right = Pt(12)
    p = tf.paragraphs[0]
    p.text = header_line
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.LEFT

    # --- 2. Grid (within content area, below taller header) ---
    col1_x = Inches(CONTENT_LEFT_IN)
    col1_w = Inches(6.0)
    col2_x = Inches(6.6)
    col2_w = Inches(min(6.4, CONTENT_RIGHT_IN - 6.6))
    content_start_y = 1.25
    # --- LEFT COLUMN: Engineering ---
    curr_y = Inches(content_start_y)
    
    # Block A: Proposal
    lbl = slide.shapes.add_textbox(col1_x, curr_y, col1_w, Inches(0.3))
    lbl.text_frame.text = "DETAILED ENGINEERING PROPOSAL"
    lbl.text_frame.paragraphs[0].font.bold = True
    lbl.text_frame.paragraphs[0].font.size = Pt(10)
    lbl.text_frame.paragraphs[0].font.color.rgb = RGBColor(128, 128, 128)
    
    box_prop = slide.shapes.add_shape(1, col1_x, curr_y + Inches(0.25), col1_w, Inches(1.6))
    box_prop.fill.solid()
    box_prop.fill.fore_color.rgb = RGBColor(250, 250, 250)
    box_prop.line.color.rgb = RGBColor(200, 200, 200)
    format_textbox(box_prop, data.get('proposal_text', '-'))
    
    curr_y += Inches(2.1)
    
    # Block B: Benchmarking
    lbl = slide.shapes.add_textbox(col1_x, curr_y, col1_w, Inches(0.3))
    lbl.text_frame.text = "BENCHMARKING & COMPARISON"
    lbl.text_frame.paragraphs[0].font.bold = True
    lbl.text_frame.paragraphs[0].font.size = Pt(10)
    lbl.text_frame.paragraphs[0].font.color.rgb = RGBColor(128, 128, 128)
    
    box_bench = slide.shapes.add_shape(1, col1_x, curr_y + Inches(0.25), col1_w, Inches(1.6))
    box_bench.fill.solid()
    box_bench.fill.fore_color.rgb = RGBColor(250, 250, 250)
    box_bench.line.color.rgb = RGBColor(200, 200, 200)
    format_textbox(box_bench, data.get('benchmarking_text', '-'))
    
    # --- RIGHT COLUMN: Business & Status ---
    curr_y = Inches(content_start_y)
    
    # Block C: Financials
    lbl = slide.shapes.add_textbox(col2_x, curr_y, col2_w, Inches(0.3))
    lbl.text_frame.text = "FINANCIAL & WEIGHT IMPACT"
    lbl.text_frame.paragraphs[0].font.bold = True
    lbl.text_frame.paragraphs[0].font.size = Pt(10)
    lbl.text_frame.paragraphs[0].font.color.rgb = RGBColor(128, 128, 128)
    
    box_fin = slide.shapes.add_shape(1, col2_x, curr_y + Inches(0.25), col2_w, Inches(1.6))
    box_fin.fill.solid()
    box_fin.fill.fore_color.rgb = RGBColor(250, 250, 250)
    box_fin.line.color.rgb = RGBColor(200, 200, 200)
    
    # Construct rich financial text
    fin_summary = (
        f"SAVING: INR {data.get('Saving Value(INR)', '0')} / vehicle\n"
        f"WEIGHT: {data.get('Weight Saving(Kg)', '0')} Kg\n\n"
        f"BREAKDOWN: {data.get('financial_text', '-')}"
    )
    format_textbox(box_fin, fin_summary)
    
    curr_y += Inches(2.1)
    
    # Block D: Feasibility
    lbl = slide.shapes.add_textbox(col2_x, curr_y, col2_w, Inches(0.3))
    lbl.text_frame.text = "FEASIBILITY & VALIDATION PLAN"
    lbl.text_frame.paragraphs[0].font.bold = True
    lbl.text_frame.paragraphs[0].font.size = Pt(10)
    lbl.text_frame.paragraphs[0].font.color.rgb = RGBColor(128, 128, 128)
    
    box_feas = slide.shapes.add_shape(1, col2_x, curr_y + Inches(0.25), col2_w, Inches(1.6))
    box_feas.fill.solid()
    box_feas.fill.fore_color.rgb = RGBColor(250, 250, 250)
    box_feas.line.color.rgb = RGBColor(200, 200, 200)
    format_textbox(box_feas, data.get('feasibility_text', '-'))
    
    # --- BOTTOM ROW: Image placeholders (within slide) ---
    img_row_h = 1.8
    img_y_in = min(5.55, SLIDE_HEIGHT_IN - img_row_h - SLIDE_MARGIN_IN)
    img_y = Inches(img_y_in)
    img_h = Inches(img_row_h)
    img_w_in = min(3.8, (CONTENT_RIGHT_IN - CONTENT_LEFT_IN - 0.6) / 3)
    img_w = Inches(img_w_in)
    gap_in = 0.3
    gap = Inches(gap_in)

    img_x1 = col1_x
    img_x2 = Inches(CONTENT_LEFT_IN + img_w_in + gap_in)
    img_x3 = Inches(CONTENT_LEFT_IN + (img_w_in + gap_in) * 2)
    placeholders = [
        (img_x1, "CURRENT MG PART", data.get('MG Vehicle Image')),
        (img_x2, "COMPETITOR REF", data.get('Competitor Vehicle Image')),
        (img_x3, "PROPOSAL SKETCH", data.get('Proposal Image')),
    ]

    for x_pos, title, path in placeholders:
        box = slide.shapes.add_shape(1, x_pos, img_y, img_w, img_h)
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(240, 240, 240)
        box.line.dash_style = 4
        box.line.color.rgb = RGBColor(160, 160, 160)
        tf_img = box.text_frame
        tf_img.word_wrap = True
        tf_img.margin_left = tf_img.margin_right = Pt(6)
        path_short = truncate_text(path or "No Image Available", max_chars=35)
        p = tf_img.paragraphs[0]
        p.text = f"[{title}]\n{path_short}"
        p.alignment = PP_ALIGN.CENTER
        p.font.color.rgb = RGBColor(120, 120, 120)
        p.font.size = Pt(9)
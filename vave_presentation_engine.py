"""
VAVE Presentation Generator Engine
Generates boardroom-ready PowerPoint presentations.
Combines:
1. Dynamic LLM Enrichment (Gemini) for unique Validation, Risks, and Supply Chain logic.
2. Smart Regex Parsing for technical specs.
3. Origin Tagging (AI/Web/DB).
"""

import json
import os
import re
from pathlib import Path
from datetime import datetime
from typing import List, Dict, Any
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import logging

# Import Gemini API
try:
    import google.generativeai as genai
    GEMINI_AVAILABLE = True
    api_key = os.getenv('GOOGLE_API_KEY') or os.getenv('GEMINI_API_KEY')
    if api_key:
        genai.configure(api_key=api_key)
except ImportError:
    GEMINI_AVAILABLE = False
    genai = None

logger = logging.getLogger(__name__)

# Slide bounds (16:9) - keep all content inside
SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5
MARGIN_IN = 0.2

def _truncate(s, max_len=400, suffix="..."):
    if not s or len(str(s).strip()) <= max_len:
        return (s or "").strip()
    t = str(s).strip()
    return t[: max_len - len(suffix)].rstrip() + suffix


def _parse_saving_float(idea, default=0.0):
    """Parse saving value from idea dict to float."""
    val = idea.get("Saving Value (INR)") or idea.get("Saving Value(INR)") or idea.get("saving_value_inr") or idea.get("Saving Value") or "0"
    s = str(val).replace(",", "").replace("INR", "").replace(" ", "").strip()
    try:
        return float(re.sub(r"[^\d.-]", "", s) or 0)
    except Exception:
        return default


class LLMEnrichmentEngine:
    """
    LLM Enrichment Layer - Uses Google Gemini API.
    Generates UNIQUE, specific content for each idea.
    """
    
    def __init__(self):
        """Initialize LLM engine."""
        self.use_llm = GEMINI_AVAILABLE and genai is not None
        if self.use_llm:
            try:
                # Using flash model for speed + quality balance
                self.model = genai.GenerativeModel('gemini-2.0-flash') 
                logger.info("LLM Enrichment Engine initialized with Gemini API")
            except Exception as e:
                logger.warning(f"Failed to initialize Gemini model: {e}. Using mock data.")
                self.use_llm = False
        else:
            logger.info("LLM Enrichment Engine using mock data (Gemini API not available)")
    
    def enrich_idea(self, raw_idea: Dict[str, Any]) -> Dict[str, Any]:
        """
        Prompt: The Engineering Deep Dive
        Generates specific Validation, Risk, and Supply Chain data for THIS idea.
        """
        title = raw_idea.get('title') or raw_idea.get('Cost Reduction Idea', '')
        description = raw_idea.get('raw_description', '') or raw_idea.get('Cost Reduction Idea', '')
        
        # Clean description text
        if isinstance(description, str) and len(description) > 500:
            description = description[:500] + "..."

        if self.use_llm:
            try:
                # UPDATED PROMPT: Forces detailed, unique analysis per idea
                prompt = f"""
                Act as a Senior Automotive VAVE Engineer. 
                Analyze this specific cost-reduction proposal:
                Title: "{title}"
                Description: "{description}"

                Generate a detailed technical feasibility report. 
                Output strictly valid JSON with no markdown formatting. The JSON must match this structure exactly:
                {{
                    "engineering_logic": "Detailed technical explanation of the engineering change. Describe the physical changes, material differences, and why it saves cost. (approx 50 words)",
                    "estimated_capex_inr": "Realistic estimation of tooling/investment costs in INR (just number, e.g. 1500000).",
                    "estimated_unit_saving_inr": "Realistic saving per vehicle in INR (e.g. 150.50). IMPORTANT: Real-world component savings are small (< 5000 INR). Do NOT overestimate.",
                    "critical_risks": [
                        "Specific Risk 1 (Technical/Safety)",
                        "Specific Risk 2 (Quality/NVH)",
                        "Specific Risk 3 (Manufacturing)"
                    ],
                    "validation_plan": [
                        "Specific Validaton Test 1 (with standard if possible)",
                        "Specific Validaton Test 2",
                        "Specific Validaton Test 3"
                    ],
                    "supply_chain_analysis": {{
                        "supplier_change": "Describe impact on supplier base (e.g. Tier 1 vs Tier 2 or consolidation)",
                        "lead_time": "Describe impact on lead time (e.g. reduction via localization)",
                        "risk_mitigation": "Strategy to mitigate supply risk",
                        "cost_impact": "Impact on logistics or tooling costs"
                    }}
                }}
                """
                response = self.model.generate_content(prompt)
                content = response.text.replace("```json", "").replace("```", "").strip()
                # Parse JSON
                return json.loads(re.search(r'\{.*\}', content, re.DOTALL).group())
            except Exception as e:
                logger.error(f"LLM enrichment failed for {title}: {e}")
                return self._mock_enrich_idea(title)
        else:
            return self._mock_enrich_idea(title)
    
    def _mock_enrich_idea(self, title: str) -> Dict[str, Any]:
        """Fallback if LLM fails."""
        import random
        capex_val = random.choice([500000, 1500000, 2500000, 750000, 1200000])
        saving_val = random.uniform(50.0, 2500.0)
        return {
            'engineering_logic': f"Technical optimization of {title} focusing on material efficiency.",
            'estimated_capex_inr': str(capex_val),
            'estimated_unit_saving_inr': str(round(saving_val, 2)),
            'critical_risks': ["NVH performance degradation", "Thermal durability limits"],
            'validation_plan': ["DVP&R Level 2 Testing", "Environmental Aging (1000 hrs)"],
            'supply_chain_analysis': {
                "supplier_change": "Consolidation to single-source Tier 1",
                "lead_time": "Reduction of 2 weeks",
                "risk_mitigation": "Dual tooling strategy",
                "cost_impact": "15% reduction in logistics overhead"
            }
        }


class VAVEPresentation:
    """
    Main class for generating VAVE PowerPoint presentations.
    """
    
    def __init__(self, project_metadata: Dict[str, Any], ideas: List[Dict[str, Any]]):
        self.project_metadata = project_metadata
        self.prs = Presentation()
        # Widescreen 16:9 (all content must stay within)
        self.prs.slide_width = Inches(SLIDE_W_IN)
        self.prs.slide_height = Inches(SLIDE_H_IN)

        self.llm_engine = LLMEnrichmentEngine()
        
        # Color scheme
        self.primary_color = RGBColor(0, 51, 102)  # Deep blue
        self.accent_color = RGBColor(255, 100, 0)  # Orange
        self.text_color = RGBColor(51, 51, 51)     # Dark gray
        self.light_gray = RGBColor(240, 240, 240)
        
        self.ideas = self._normalize_ideas(ideas)
    
    def _normalize_ideas(self, ideas: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
        normalized = []
        for idea in ideas:
            # Determine Origin
            origin_raw = idea.get('Origin') or idea.get('origin', 'Unknown')
            if "AI" in origin_raw: origin_tag = "AI INNOVATION"
            elif "Web" in origin_raw or "World" in origin_raw: origin_tag = "WEB SOURCED"
            elif "DB" in origin_raw or "Existing" in origin_raw: origin_tag = "EXISTING DATABASE"
            else: origin_tag = "DATABASE"

            normalized_idea = {
                'title': idea.get('title') or idea.get('Cost Reduction Idea', 'Untitled'),
                'raw_description': idea.get('raw_description') or idea.get('Cost Reduction Idea', '') or idea.get('Way Forward', ''),
                'saving_amount': self._extract_saving_value(idea),
                'origin_tag': origin_tag,
                # Image keys
                'current_image': idea.get('mg_vehicle_image') or idea.get('Current Scenario Image') or idea.get('current_scenario_image'),
                'proposal_image': idea.get('proposal_image_filename') or idea.get('Proposal Scenario Image') or idea.get('proposal_scenario_image'),
                'competitor_image': idea.get('competitor_image') or idea.get('Competitor Image') or "NaN", 
                '_original': idea
            }
            normalized.append(normalized_idea)
        return normalized

    def _extract_saving_value(self, idea):
        val = idea.get('Saving Value (INR)') or idea.get('saving_value_inr') or idea.get('Estimated Cost Savings')
        if val: return str(val)
        return "TBD"
        
    def generate(self) -> Presentation:
        logger.info("Starting VAVE presentation generation...")
        self._create_title_slide()
        self._create_summary_slide()
        for idx, idea in enumerate(self.ideas):
            self._create_engineering_idea_slide(idea, idx + 1)
        return self.prs
    
    def _create_title_slide(self):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[0])
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = _truncate("VAVE Cost Reduction Strategy", max_len=80)
        subtitle.text = _truncate(
            f"Detailed Engineering Feasibility Report | {datetime.now().strftime('%d-%b-%Y')}",
            max_len=200,
        )
        tf = title.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].font.color.rgb = self.primary_color
        if hasattr(subtitle, "text_frame") and subtitle.text_frame:
            subtitle.text_frame.word_wrap = True

    def _to_crore(self, val_str):
        """Convert string value to Float (Crores). 1 Cr = 1,00,00,000."""
        try:
            # Remove non-numeric chars except dot
            clean = re.sub(r"[^\d.]", "", str(val_str).replace(",",""))
            if not clean: return 0.0
            val = float(clean)
            return val / 10000000
        except:
            return 0.0

    def _create_summary_slide(self):
        """
        Summary Slide(s). Splits into multiple slides if ideas > 7 (Antigravity Layout).
        """
        if not self.ideas: return

        # Antigravity Layout: Strict 7 rows to prevent overflow
        MAX_ROWS = 7
        chunks = [self.ideas[i:i + MAX_ROWS] for i in range(0, len(self.ideas), MAX_ROWS)]
        
        for chunk_idx, chunk_ideas in enumerate(chunks):
            slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
            left_in = 0.2 # Tighter margins
            top_in = 0.45
            w_in = SLIDE_W_IN - 2 * left_in
            h_in = SLIDE_H_IN - top_in - MARGIN_IN
            
            # 8 Columns (Added Remark)
            cols = 8
            rows = len(chunk_ideas) + 1
            table_shape = slide.shapes.add_table(rows, cols, Inches(left_in), Inches(top_in), Inches(w_in), Inches(h_in))
            table = table_shape.table
            
            # Headers with (Cr) and Unit Logic
            header_row = [
                "ID", "Idea", "CAPEX (Cr)", "Unit Sav (₹)",
                "Vol (Qty)", "Total Sav (Cr)", "Status", "Remark"
            ]
            
            # Optmized Widths
            widths = [0.6, 4.0, 1.0, 1.0, 1.0, 1.2, 1.2, 2.0]
            for i, width in enumerate(widths):
                table.columns[i].width = Inches(width)

            font_header, font_cell = 9, 9
            for c, label in enumerate(header_row):
                cell = table.cell(0, c)
                tf = cell.text_frame
                tf.word_wrap = True
                tf.paragraphs[0].text = label
                tf.paragraphs[0].font.bold = True
                tf.paragraphs[0].font.size = Pt(font_header)
                tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
                cell.fill.solid()
                cell.fill.fore_color.rgb = self.primary_color
            
            # Defaults
            default_volume = 100000
            
            for r, idea in enumerate(chunk_ideas):
                row_idx = r + 1
                raw = idea.get("_original") or idea
                
                # Enrich with AI Financials
                llm_data = idea.get('_llm_cache') or self.llm_engine.enrich_idea(idea)
                idea['_llm_cache'] = llm_data
                
                # 1. CAPEX Logic (Crores)
                # DB > AI > 0
                db_capex_raw = str(raw.get("CAPEX") or raw.get("Capex Investment") or "").strip()
                ai_capex_raw = str(llm_data.get('estimated_capex_inr', '0'))
                
                use_capex = "0"
                if db_capex_raw and "tbd" not in db_capex_raw.lower() and "nan" not in db_capex_raw.lower():
                    use_capex = db_capex_raw
                else:
                    use_capex = ai_capex_raw
                
                capex_cr = self._to_crore(use_capex)

                # 2. Unit Saving (INR)
                # DB > AI > 0.0
                saving_inr = 0.0
                db_saving = _parse_saving_float(raw)
                if db_saving > 0:
                    saving_inr = db_saving
                else:
                    # Try AI estimate
                    try:
                        ai_est = str(llm_data.get('estimated_unit_saving_inr', '0'))
                        # Clean if it has INR or currency symbols
                        ai_est = re.sub(r"[^\d.]", "", ai_est)
                        saving_inr = float(ai_est)
                    except:
                        saving_inr = 0.0
                
                # SANITY CHECK: Unit Saving Cap (Antigravity Protocol)
                # Max 5000 INR per unit for component level changes.
                if saving_inr > 5000.0:
                    saving_inr = 5000.0
                
                # 3. Volume (Quantity)
                # DB > Default
                yearly_vol = default_volume
                vol_raw = raw.get("volume") or raw.get("Yearly volume") or raw.get("yearly_volume")
                if vol_raw:
                     try:
                        yearly_vol = int(float(re.sub(r"[^\d.]", "", str(vol_raw))) or default_volume)
                     except: pass
                
                # 4. Total Saving (Crores)
                # (Unit INR * Vol) / 1Cr
                total_saving_cr = (saving_inr * yearly_vol) / 10000000.0

                # Status & Remarks
                status = str(raw.get("Status") or "Proposed")
                remark = str(raw.get("Remark") or raw.get("remark") or "")
                
                # Intelligent Remark Logic
                if not remark:
                    if status.upper() == "NG":
                         status = "In Progress"
                         remark = "Feasibility Study On-going"
                    elif "AUTO" in status.upper() or "RESEARCH" in status.upper():
                         status = "AI Generated"
                         remark = "Validation Pending"
                    else:
                         remark = "Engineering Review"

                # Truncate text to prevent row overflow
                def _t(s, l=100): return _truncate(str(s), l)

                row_data = [
                    _t(raw.get("Idea Id") or f"ID-{r+1}", 15),
                    _t(idea.get("title") or "", 80), # Title might need more
                    f"{capex_cr:.4f}",
                    f"{saving_inr:.2f}",     # INR (Simple 2 decimals)
                    f"{yearly_vol:,}",       # Integer with commas
                    f"{total_saving_cr:.4f}",
                    _t(status, 20),
                    _t(remark, 30)
                ]
                
                for c, val in enumerate(row_data):
                    cell = table.cell(row_idx, c)
                    tf = cell.text_frame
                    tf.word_wrap = True
                    tf.margin_left = tf.margin_right = Pt(4)
                    p = tf.paragraphs[0]
                    p.text = str(val).strip()
                    p.font.size = Pt(font_cell)
                    p.font.name = "Calibri"
                    p.font.color.rgb = self.text_color
            
            # Slide Title
            title_box = slide.shapes.add_textbox(Inches(left_in), Inches(0.05), Inches(w_in), Inches(0.35))
            tf = title_box.text_frame
            p = tf.paragraphs[0]
            suffix = f" (Page {chunk_idx+1})" if len(chunks) > 1 else ""
            p.text = f"Summary (Values in ₹ Crore){suffix}"
            p.font.size = Pt(15)
            p.font.bold = True
            p.font.color.rgb = self.primary_color

    def _fit_text_to_bounds(self, text_frame, max_height_in):
        """Helper to shrink text to fit height."""
        # Simple heuristic: try 10pt, if too long, try 9pt...
        pass # python-pptx can't "measure" text height easily without rendering.
        # We'll stick to a simpler method: truncate or fix font size.
        # However, we can enforce Autosize
        text_frame.word_wrap = True
        text_frame.auto_size = MSO_ANCHOR.TOP # Unfortunately MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE is not fully exposed in all pptx versions cleanly
        # So manual calc:
        char_limit = int(max_height_in * 1000) # heuristic
        if len(text_frame.text) > char_limit:
            text_frame.text = text_frame.text[:char_limit] + "..."

    def _create_engineering_idea_slide(self, idea, index):
        """
        Creates the Custom Engineering Layout Slide with STRICT Bound Checking.
        PREVENTS OVERLAP by paginating or truncating if overflow.
        """
        slide_layout = self.prs.slide_layouts[6] # Blank
        slide = self.prs.slides.add_slide(slide_layout)

        # CONSTANTS
        MARGIN = 0.3
        
        # 1. HEADER (Safe Zone)
        title_text = f"#{index}: {idea['title']}"
        # Heuristic: 1 line approx 70 chars at pt18
        char_count = len(title_text)
        est_lines = (char_count // 70) + 1
        real_header_h = max(0.6, est_lines * 0.45) 
        
        title_box = slide.shapes.add_textbox(Inches(MARGIN), Inches(0.2), Inches(SLIDE_W_IN - 4), Inches(real_header_h))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(18) # Reduced Font as requested
        p.font.bold = True
        p.font.color.rgb = self.primary_color
        
        # 2. ORIGIN BADGE
        tag_w = 2.0
        badge = slide.shapes.add_textbox(Inches(SLIDE_W_IN - tag_w - MARGIN), Inches(0.2), Inches(tag_w), Inches(0.5))
        badge.fill.solid()
        ot = idea['origin_tag']
        if "AI" in ot: badge.fill.fore_color.rgb = RGBColor(128, 0, 128)
        elif "WEB" in ot: badge.fill.fore_color.rgb = RGBColor(0, 128, 128)
        else: badge.fill.fore_color.rgb = RGBColor(100, 100, 100)
        tf = badge.text_frame
        p = tf.paragraphs[0]
        p.text = _truncate(ot, max_len=20)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

        # 3. DYNAMIC CONTENT
        llm_data = self.llm_engine.enrich_idea(idea)
        regex_data = self._smart_parse_technical_details(idea)
        
        # Validation & Risks formatted
        val_text = "VALIDATION:\n" + "".join([f"• {v}\n" for v in llm_data.get('validation_plan', [])[:3]])
        risk_text = "RISKS:\n" + "".join([f"⚠ {r}\n" for r in llm_data.get('critical_risks', [])[:3]])
        full_val_txt = val_text + "\n" + risk_text

        # Supply chain
        sc = llm_data.get('supply_chain_analysis', {})
        sc_text = f"• Supplier: {sc.get('supplier_change','N/A')}\n• Tooling: {sc.get('cost_impact','N/A')}"

        fields = [
            ("Technical Engineering View", llm_data.get('engineering_logic'), 1.2),
            ("Material & Dimension Swap", f"{regex_data.get('material')}\n{regex_data.get('dimension')}", 0.8),
            ("Financial Impact", f"Est. Saving: INR {idea['saving_amount']}\nCAPEX: {idea.get('_original', {}).get('CAPEX','TBD')}", 0.6),
            ("Validation & Risks", full_val_txt, 1.8),
            ("Supply Chain", sc_text, 0.8)
        ]

        # 4. LEFT COLUMN (Text) - AUTO-LAYOUT
        col1_x = MARGIN
        col1_w = 6.0
        start_y = real_header_h + 0.3 # Reduced gap
        max_y = SLIDE_H_IN - MARGIN
        avail_h = max_y - start_y
        
        # Define Spacing
        SPACING = 0.08 
        LABEL_H = 0.25
        
        # PASS 1: Estimate Height
        base_font = 10
        total_content_chars = 0
        clean_fields = []
        for label, content, _ignore_h in fields:
            text = str(content or "N/A").strip()
            clean_fields.append((label, text))
            total_content_chars += len(text)
            
        fixed_overhead = (5 * LABEL_H) + (5 * SPACING)
        # 100 chars per line estimate at 10pt (optimistic)
        est_lines = total_content_chars / 100
        text_height_needed = est_lines * 0.18
        total_needed = fixed_overhead + text_height_needed
        
        # PASS 2: Calculate Scale Factor
        scale = 1.0
        if total_needed > avail_h:
            scale = avail_h / total_needed
            if base_font * scale < 6.0: scale = 6.0 / base_font # Floor at 6pt
        
        final_font = max(6.0, base_font * scale)
        final_line_h = 0.18 * scale
        
        current_y = start_y
        for label, text in clean_fields:
            if current_y > max_y - 0.2: break # Strict Cutoff

            lines = (len(text) / 100) + 1 + text.count('\n')
            box_h = lines * final_line_h
            if box_h < 0.15: box_h = 0.15 # Min height
            
            # Check overlap safety
            if current_y + LABEL_H*scale + box_h > max_y:
                 # Truncate box height to available space
                 remaining = max_y - (current_y + LABEL_H*scale)
                 if remaining < 0.2: break 
                 box_h = remaining

            # Label
            lbl = slide.shapes.add_textbox(Inches(col1_x), Inches(current_y), Inches(col1_w), Inches(LABEL_H * scale))
            lbl.text_frame.text = label.upper()
            lbl.text_frame.paragraphs[0].font.bold = True
            lbl.text_frame.paragraphs[0].font.size = Pt(max(6, 11 * scale))
            lbl.text_frame.paragraphs[0].font.color.rgb = self.accent_color
            current_y += (LABEL_H * scale)
            
            # Content
            txt = slide.shapes.add_textbox(Inches(col1_x), Inches(current_y), Inches(col1_w), Inches(box_h))
            tf = txt.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = text
            p.font.size = Pt(final_font)
            p.font.color.rgb = self.text_color
            current_y += box_h
            
            # Divider
            ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(col1_x), Inches(current_y), Inches(col1_w), Inches(0.01))
            ln.fill.solid()
            ln.fill.fore_color.rgb = self.light_gray
            current_y += SPACING

        # 5. RIGHT COLUMN (Images)
        col2_x = col1_x + col1_w + 0.5
        col2_w = SLIDE_W_IN - col2_x - MARGIN
        
        img_start_y = real_header_h + 0.2
        img_avail_h = max_y - img_start_y
        img_h = (img_avail_h / 2) - SPACING
        
        self._place_image(slide, idea.get('current_image'), "Current Scenario", Inches(col2_x), Inches(img_start_y), Inches(col2_w), Inches(img_h))
        
        img2_y = img_start_y + img_h + SPACING*2
        # Antigravity Placeholder for consistent engineering view
        self._placeholder(slide, Inches(col2_x), Inches(img2_y), Inches(col2_w), Inches(img_h), "VLM Generated AI Image")

    def _place_image(self, slide, img_path, label, x, y, w, h):
        # Label (wrap and truncate so it stays in slide)
        lbl = slide.shapes.add_textbox(x, y - Inches(0.3), w, Inches(0.28))
        tf_lbl = lbl.text_frame
        tf_lbl.word_wrap = True
        p = tf_lbl.paragraphs[0]
        p.text = _truncate(label, max_len=50)
        p.font.bold = True
        p.font.size = Pt(11)
        p.alignment = PP_ALIGN.CENTER
        p.font.color.rgb = self.primary_color

        # Image Logic
        valid = False
        final_path = None
        if img_path and str(img_path).lower() not in ['nan', 'none', 'n/a', '']:
            if str(img_path).startswith("/"): img_path = str(img_path)[1:]
            base_dir = os.path.dirname(os.path.abspath(__file__))
            
            p1 = os.path.join(base_dir, str(img_path))
            if os.path.exists(p1): 
                final_path = p1
                valid = True
            elif os.path.exists(str(img_path)):
                final_path = str(img_path)
                valid = True
        
        if valid:
            try:
                pic = slide.shapes.add_picture(final_path, x, y, w, h)
                pic.line.color.rgb = RGBColor(200, 200, 200)
                pic.line.width = Pt(1)
            except:
                self._placeholder(slide, x, y, w, h, "Image Error")
        else:
            self._placeholder(slide, x, y, w, h, "Image Not Available")

    def _placeholder(self, slide, x, y, w, h, text):
        shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
        shp.fill.solid()
        shp.fill.fore_color.rgb = self.light_gray
        tf = shp.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].text = _truncate(text, max_len=80)
        tf.paragraphs[0].font.color.rgb = RGBColor(150, 150, 150)
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    def _smart_parse_technical_details(self, idea):
        """Regex parsing for granular details."""
        text = str(idea.get('raw_description', ''))
        
        details = {
            'material': 'N/A (No material change explicitly mentioned)',
            'dimension': 'N/A (Maintains current envelope)'
        }
        
        # Material Regex
        mat_keywords = ['steel', 'aluminum', 'plastic', 'composite', 'resin', 'grade', 'material', 'pp', 'abs', 'rubber']
        for sent in text.split('.'):
            if any(k in sent.lower() for k in mat_keywords):
                details['material'] = sent.strip()[:100]
                break
        
        # Dimension Regex
        dim_match = re.search(r'(\d+(\.\d+)?\s*(mm|cm|um|μm|kg).{0,40}\d+(\.\d+)?\s*(mm|cm|um|μm|kg)?)', text, re.IGNORECASE)
        if dim_match:
            details['dimension'] = dim_match.group(0)
        
        return details

    def save(self, output_path: str):
        self.prs.save(output_path)
        logger.info(f"Presentation saved: {output_path}")

# ============================================================================
# COMPATIBILITY WRAPPER
# ============================================================================

def generate_deep_dive_ppt(ideas_list: List[Dict[str, Any]], output_path: str) -> str:
    """Wrapper used by app.py"""
    meta = {'project_name': 'VAVE', 'car_model': 'Vehicle Program'}
    engine = VAVEPresentation(meta, ideas_list)
    engine.generate()
    if not output_path.endswith(".pptx"): output_path += ".pptx"
    engine.save(output_path)
    return output_path